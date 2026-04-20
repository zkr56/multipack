"""
╔══════════════════════════════════════════════════════════════════╗
║  MULTIPACK SA — TABLEAU DE BORD DIRECTION GÉNÉRALE              ║
║  Version Finale v9  ·  Zone Industrielle de Yopougon, Abidjan   ║
║  Pilotage Commercial · Reporting · Contrôle de Gestion          ║
╚══════════════════════════════════════════════════════════════════╝
"""

import streamlit as st
import pandas as pd
import numpy as np
import plotly.graph_objects as go
from datetime import datetime, timedelta
from sklearn.linear_model import LinearRegression
from sklearn.preprocessing import PolynomialFeatures
import warnings
warnings.filterwarnings("ignore")
import io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

st.set_page_config(
    page_title="MULTIPACK SA — Direction",
    page_icon="📦",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ── Palette officielle MULTIPACK ──────────────────────────────────
C = {
    "bleu":       "#4F6EF7",   # indigo principal
    "bleu_clair": "#C7D2FE",   # indigo clair
    "bleu_fonce": "#3730A3",   # indigo foncé
    "rose":       "#F472B6",   # rose accent
    "violet":     "#7C3AED",   # violet profond
    "vert":       "#10B981",   # emeraude
    "rouge":      "#EF4444",   # rouge vif
    "orange":     "#F59E0B",   # ambre
    "or":         "#D97706",   # or foncé
    "cyan":       "#06B6D4",   # cyan
    "gris_bg":    "#F8FAFC",   # fond page
    "blanc":      "#FFFFFF",
    "texte":      "#1E293B",   # slate 800
    "muted":      "#64748B",   # slate 500
    "bordure":    "#E2E8F0",   # slate 200
    "sidebar":    "#0F172A",   # slate 950
    "card_bg":    "#FFFFFF",
}

# Palettes graphiques (14 couleurs distinctes)
PIE = [
    "#4F6EF7","#F472B6","#7C3AED","#F59E0B","#10B981",
    "#3B82F6","#EC4899","#8B5CF6","#FBBF24","#34D399",
    "#06B6D4","#EF4444","#A78BFA","#6EE7B7",
]

# ══════════════════════════════════════════════════════════════════
# CSS PREMIUM — MULTIPACK FINAL
# ══════════════════════════════════════════════════════════════════
st.markdown(f"""
<style>
@import url('https://fonts.googleapis.com/css2?family=Plus+Jakarta+Sans:wght@300;400;500;600;700;800;900&display=swap');

/* ── RESET & BASE ── */
*, *::before, *::after {{ box-sizing: border-box; margin: 0; padding: 0; }}
html, body, [class*="css"] {{
    font-family: 'Plus Jakarta Sans', sans-serif !important;
    background: {C['gris_bg']};
    color: {C['texte']};
    -webkit-font-smoothing: antialiased;
}}
.main .block-container {{
    background: {C['gris_bg']};
    padding: 0 2rem 3rem;
    max-width: 1700px;
}}

/* ── SIDEBAR ── */
[data-testid="stSidebar"] {{
    background: {C['sidebar']} !important;
    border-right: 1px solid rgba(255,255,255,0.04) !important;
}}
[data-testid="stSidebar"] * {{ color: #CBD5E1 !important; font-family: 'Plus Jakarta Sans', sans-serif !important; }}
[data-testid="stSidebar"] .stRadio div[role="radiogroup"] label {{
    background: rgba(255,255,255,0.03);
    border-radius: 10px;
    padding: 10px 14px !important;
    margin: 2px 0 !important;
    border: 1px solid rgba(255,255,255,0.04);
    font-size: 0.83rem !important;
    font-weight: 500 !important;
    color: #94A3B8 !important;
    transition: all 0.18s ease;
    cursor: pointer;
}}
[data-testid="stSidebar"] .stRadio div[role="radiogroup"] label:hover {{
    background: rgba(79,110,247,0.15) !important;
    border-color: rgba(79,110,247,0.4) !important;
    color: #E2E8F0 !important;
}}

/* ── TOP NAVBAR ── */
.top-navbar {{
    background: {C['sidebar']};
    border-radius: 16px;
    padding: 0 24px;
    margin-bottom: 20px;
    display: flex;
    align-items: center;
    justify-content: space-between;
    box-shadow: 0 4px 24px rgba(0,0,0,0.18);
    border: 1px solid rgba(255,255,255,0.05);
    overflow-x: auto;
    gap: 2px;
    min-height: 72px;
}}
.nav-brand {{
    display: flex;
    align-items: center;
    gap: 12px;
    padding: 16px 0;
    min-width: 170px;
    white-space: nowrap;
    border-right: 1px solid rgba(255,255,255,0.06);
    margin-right: 12px;
    padding-right: 20px;
}}
.nav-brand .logo {{ font-size: 1.8rem; }}
.nav-brand .name {{ font-size: 0.92rem; font-weight: 800; color: white; letter-spacing: -0.02em; line-height: 1.2; }}
.nav-brand .tagline {{ font-size: 0.6rem; color: #475569; font-weight: 500; text-transform: uppercase; letter-spacing: 0.08em; }}
.nav-links {{ display: flex; align-items: center; gap: 2px; flex: 1; justify-content: center; padding: 10px 0; flex-wrap: nowrap; }}
.nav-link {{
    display: flex;
    flex-direction: column;
    align-items: center;
    padding: 8px 11px;
    border-radius: 10px;
    cursor: pointer;
    transition: all 0.2s ease;
    border: 1px solid transparent;
    white-space: nowrap;
    min-width: 64px;
    text-decoration: none;
}}
.nav-link:hover {{ background: rgba(79,110,247,0.14); border-color: rgba(79,110,247,0.3); }}
.nav-link.active {{ background: rgba(79,110,247,0.22); border-color: rgba(79,110,247,0.5); }}
.nav-link .nl-ico {{ font-size: 1.15rem; line-height: 1; }}
.nav-link .nl-txt {{ font-size: 0.56rem; font-weight: 700; color: #475569; text-transform: uppercase; letter-spacing: 0.06em; margin-top: 4px; }}
.nav-link.active .nl-txt {{ color: {C['bleu']}; }}
.nav-right {{ display: flex; align-items: center; gap: 12px; padding: 10px 0; white-space: nowrap; border-left: 1px solid rgba(255,255,255,0.06); margin-left: 12px; padding-left: 20px; }}
.nav-dot {{ width: 8px; height: 8px; border-radius: 50%; background: {C['vert']}; box-shadow: 0 0 8px {C['vert']}; flex-shrink: 0; animation: pulse 2s infinite; }}
@keyframes pulse {{ 0%,100% {{ opacity:1; }} 50% {{ opacity:0.5; }} }}
.nav-date {{ font-size: 0.7rem; color: #475569; line-height: 1.5; }}

/* ── KPI CARDS ── */
.kpi {{
    background: {C['blanc']};
    border-radius: 16px;
    padding: 20px 22px;
    border: 1px solid {C['bordure']};
    box-shadow: 0 1px 3px rgba(0,0,0,0.04), 0 4px 16px rgba(0,0,0,0.04);
    position: relative;
    overflow: hidden;
    transition: transform 0.2s ease, box-shadow 0.2s ease;
    height: 100%;
}}
.kpi:hover {{ transform: translateY(-2px); box-shadow: 0 8px 32px rgba(0,0,0,0.10); }}
.kpi::after {{
    content: '';
    position: absolute;
    top: 0; left: 0;
    width: 100%; height: 3px;
    border-radius: 16px 16px 0 0;
}}
.kpi.bleu::after   {{ background: linear-gradient(90deg, {C['bleu']}, {C['bleu_clair']}); }}
.kpi.vert::after   {{ background: linear-gradient(90deg, {C['vert']}, #6EE7B7); }}
.kpi.rouge::after  {{ background: linear-gradient(90deg, {C['rouge']}, #FCA5A5); }}
.kpi.orange::after {{ background: linear-gradient(90deg, {C['orange']}, #FCD34D); }}
.kpi.violet::after {{ background: linear-gradient(90deg, {C['violet']}, #C4B5FD); }}
.kpi.rose::after   {{ background: linear-gradient(90deg, {C['rose']}, #FBCFE8); }}
.kpi.or::after     {{ background: linear-gradient(90deg, {C['or']}, #FDE68A); }}
.kpi.cyan::after   {{ background: linear-gradient(90deg, {C['cyan']}, #A5F3FC); }}
.kpi .ico {{ font-size: 1.5rem; margin-bottom: 10px; display: block; }}
.kpi .lbl {{ font-size: 0.67rem; font-weight: 700; color: {C['muted']}; text-transform: uppercase; letter-spacing: 0.08em; }}
.kpi .val {{ font-size: 1.6rem; font-weight: 900; color: {C['texte']}; margin: 6px 0 3px; line-height: 1.1; letter-spacing: -0.03em; }}
.kpi .sub {{ font-size: 0.72rem; color: {C['muted']}; font-weight: 500; }}

/* ── PREV KPI (dark) ── */
.prev-kpi {{
    background: linear-gradient(135deg, #0F172A, #1E293B);
    border-radius: 16px;
    padding: 20px 22px;
    border: 1px solid rgba(79,110,247,0.25);
    box-shadow: 0 4px 16px rgba(79,110,247,0.15);
    height: 100%;
}}
.prev-kpi .ico {{ font-size: 1.5rem; margin-bottom: 10px; display: block; }}
.prev-kpi .lbl {{ font-size: 0.67rem; font-weight: 700; color: #475569; text-transform: uppercase; letter-spacing: 0.08em; }}
.prev-kpi .val {{ font-size: 1.5rem; font-weight: 900; color: white; margin: 6px 0 3px; letter-spacing: -0.03em; }}
.prev-kpi .sub {{ font-size: 0.72rem; color: #64748B; font-weight: 500; }}

/* ── SECTION TITLES ── */
.sec {{
    display: flex;
    align-items: center;
    gap: 10px;
    font-size: 0.72rem;
    font-weight: 800;
    color: {C['muted']};
    text-transform: uppercase;
    letter-spacing: 0.1em;
    padding: 0 0 10px;
    margin: 28px 0 14px;
    border-bottom: 2px solid {C['bordure']};
}}
.sec::before {{
    content: '';
    display: block;
    width: 14px; height: 14px;
    background: linear-gradient(135deg, {C['bleu']}, {C['violet']});
    border-radius: 4px;
    flex-shrink: 0;
}}

/* ── CHART CARD WRAPPER ── */
.chart-wrap {{
    background: {C['blanc']};
    border-radius: 16px;
    border: 1px solid {C['bordure']};
    padding: 20px 22px;
    box-shadow: 0 1px 3px rgba(0,0,0,0.04), 0 4px 16px rgba(0,0,0,0.04);
    margin-bottom: 4px;
}}

/* ── COMMENT BOXES ── */
.comment-box {{
    background: linear-gradient(135deg, #EFF6FF, #EDE9FE);
    border-left: 4px solid {C['bleu']};
    border-radius: 0 12px 12px 0;
    padding: 14px 18px;
    margin: 10px 0 16px;
    font-size: 0.83rem;
    color: #1E3A5F;
    line-height: 1.6;
}}
.comment-box .ct {{
    font-weight: 800;
    color: {C['bleu']};
    font-size: 0.7rem;
    text-transform: uppercase;
    letter-spacing: 0.08em;
    margin-bottom: 6px;
    display: flex;
    align-items: center;
    gap: 6px;
}}
.prevision-box {{
    background: linear-gradient(135deg, #FFFBEB, #FEF3C7);
    border-left: 4px solid {C['orange']};
    border-radius: 0 12px 12px 0;
    padding: 14px 18px;
    margin: 10px 0 16px;
    font-size: 0.83rem;
    color: #451A03;
    line-height: 1.6;
}}
.prevision-box .ct {{
    font-weight: 800;
    color: {C['or']};
    font-size: 0.7rem;
    text-transform: uppercase;
    letter-spacing: 0.08em;
    margin-bottom: 6px;
}}

/* ── ALERTS ── */
.alert-r {{
    background: #FFF1F2;
    border-left: 4px solid {C['rouge']};
    border-radius: 0 10px 10px 0;
    padding: 12px 16px;
    margin: 6px 0;
    font-size: 0.82rem;
    color: #7F1D1D;
    display: flex;
    gap: 10px;
    align-items: flex-start;
}}
.alert-y {{
    background: #FFFBEB;
    border-left: 4px solid {C['orange']};
    border-radius: 0 10px 10px 0;
    padding: 12px 16px;
    margin: 6px 0;
    font-size: 0.82rem;
    color: #451A03;
    display: flex;
    gap: 10px;
    align-items: flex-start;
}}
.alert-g {{
    background: #ECFDF5;
    border-left: 4px solid {C['vert']};
    border-radius: 0 10px 10px 0;
    padding: 12px 16px;
    margin: 6px 0;
    font-size: 0.82rem;
    color: #052E16;
    display: flex;
    gap: 10px;
    align-items: flex-start;
}}

/* ── PROGRESS BARS ── */
.prog {{ margin: 8px 0; }}
.prog .pl {{ display: flex; justify-content: space-between; font-size: 0.73rem; color: {C['muted']}; margin-bottom: 5px; font-weight: 500; }}
.prog .pb {{ height: 8px; background: {C['bordure']}; border-radius: 20px; overflow: hidden; }}
.prog .pf {{ height: 100%; border-radius: 20px; transition: width 0.5s ease; }}

/* ── WELCOME HERO ── */
.welcome-hero {{
    background: linear-gradient(135deg, {C['sidebar']} 0%, #1E293B 45%, {C['violet']} 100%);
    border-radius: 20px;
    padding: 56px 52px;
    margin-bottom: 28px;
    position: relative;
    overflow: hidden;
    box-shadow: 0 20px 60px rgba(124,58,237,0.3);
}}
.welcome-hero::before {{
    content: "";
    position: absolute;
    top: -100px; right: -100px;
    width: 400px; height: 400px;
    border-radius: 50%;
    background: radial-gradient(circle, rgba(79,110,247,0.3) 0%, transparent 65%);
    pointer-events: none;
}}
.welcome-hero::after {{
    content: "📦";
    position: absolute;
    right: 56px; top: 50%;
    transform: translateY(-50%);
    font-size: 180px;
    opacity: 0.06;
    pointer-events: none;
}}
.hero-badge {{
    display: inline-flex;
    align-items: center;
    gap: 6px;
    background: rgba(79,110,247,0.25);
    border: 1px solid rgba(79,110,247,0.4);
    border-radius: 20px;
    padding: 5px 14px;
    font-size: 0.7rem;
    font-weight: 700;
    color: {C['bleu_clair']};
    text-transform: uppercase;
    letter-spacing: 0.08em;
    margin-bottom: 18px;
}}
.hero-title {{
    font-size: 2.8rem;
    font-weight: 900;
    color: white;
    margin: 0 0 12px;
    letter-spacing: -0.05em;
    line-height: 1.05;
}}
.hero-sub {{
    font-size: 1rem;
    color: rgba(255,255,255,0.60);
    font-weight: 400;
    margin-bottom: 32px;
    line-height: 1.6;
    max-width: 620px;
}}
.hero-stats {{ display: flex; gap: 16px; flex-wrap: wrap; }}
.hero-stat {{
    background: rgba(255,255,255,0.07);
    backdrop-filter: blur(10px);
    border: 1px solid rgba(255,255,255,0.12);
    border-radius: 14px;
    padding: 14px 22px;
    text-align: center;
    min-width: 100px;
    transition: transform 0.2s;
}}
.hero-stat:hover {{ transform: translateY(-2px); }}
.hero-stat .hs-val {{ font-size: 1.5rem; font-weight: 900; color: white; display: block; line-height: 1; letter-spacing: -0.03em; }}
.hero-stat .hs-lbl {{ font-size: 0.6rem; color: rgba(255,255,255,0.45); text-transform: uppercase; letter-spacing: 0.08em; margin-top: 5px; display: block; }}

/* ── MODULE CARDS ── */
.module-card {{
    background: {C['blanc']};
    border-radius: 16px;
    border: 1px solid {C['bordure']};
    padding: 24px;
    box-shadow: 0 1px 3px rgba(0,0,0,0.04), 0 4px 16px rgba(0,0,0,0.04);
    transition: transform 0.2s ease, box-shadow 0.2s ease;
    height: 100%;
    position: relative;
    overflow: hidden;
    cursor: default;
}}
.module-card:hover {{ transform: translateY(-4px); box-shadow: 0 12px 40px rgba(0,0,0,0.12); }}
.module-card::before {{
    content: '';
    position: absolute;
    top: 0; left: 0;
    width: 100%; height: 3px;
    border-radius: 16px 16px 0 0;
}}
.mc-bleu::before   {{ background: linear-gradient(90deg, {C['bleu']}, {C['violet']}); }}
.mc-vert::before   {{ background: linear-gradient(90deg, {C['vert']}, {C['cyan']}); }}
.mc-rouge::before  {{ background: linear-gradient(90deg, {C['rouge']}, {C['rose']}); }}
.mc-orange::before {{ background: linear-gradient(90deg, {C['orange']}, {C['rouge']}); }}
.mc-violet::before {{ background: linear-gradient(90deg, {C['violet']}, {C['rose']}); }}
.mc-or::before     {{ background: linear-gradient(90deg, {C['or']}, {C['orange']}); }}
.mc-cyan::before   {{ background: linear-gradient(90deg, {C['cyan']}, {C['vert']}); }}
.module-card .mc-ico {{ font-size: 2.2rem; margin-bottom: 14px; display: block; }}
.module-card .mc-title {{ font-size: 1rem; font-weight: 800; color: {C['texte']}; margin-bottom: 8px; letter-spacing: -0.02em; }}
.module-card .mc-desc {{ font-size: 0.78rem; color: {C['muted']}; line-height: 1.55; }}
.module-card .mc-tags {{ display: flex; flex-wrap: wrap; gap: 5px; margin-top: 14px; }}
.mc-tag {{
    background: {C['gris_bg']}; border: 1px solid {C['bordure']};
    border-radius: 20px; padding: 3px 9px;
    font-size: 0.62rem; font-weight: 700; color: {C['muted']};
    text-transform: uppercase; letter-spacing: 0.06em;
}}

/* ── GUIDE STEPS ── */
.guide-step {{
    display: flex;
    gap: 14px;
    align-items: flex-start;
    background: {C['blanc']};
    border: 1px solid {C['bordure']};
    border-radius: 14px;
    padding: 16px 18px;
    margin-bottom: 10px;
    box-shadow: 0 1px 4px rgba(0,0,0,0.04);
    transition: box-shadow 0.2s;
}}
.guide-step:hover {{ box-shadow: 0 4px 16px rgba(0,0,0,0.08); }}
.gs-num {{
    min-width: 32px; height: 32px;
    background: linear-gradient(135deg, {C['bleu']}, {C['violet']});
    border-radius: 50%;
    display: flex; align-items: center; justify-content: center;
    font-size: 0.78rem; font-weight: 900; color: white; flex-shrink: 0;
    box-shadow: 0 4px 12px rgba(79,110,247,0.4);
}}
.gs-content .gs-title {{ font-size: 0.9rem; font-weight: 700; color: {C['texte']}; margin-bottom: 4px; }}
.gs-content .gs-desc {{ font-size: 0.77rem; color: {C['muted']}; line-height: 1.5; }}

/* ── INFO ROWS ── */
.info-row {{
    display: flex; justify-content: space-between; align-items: center;
    padding: 9px 0; border-bottom: 1px solid {C['bordure']};
    font-size: 0.81rem;
}}
.info-row:last-child {{ border-bottom: none; }}
.info-row .ir-label {{ color: {C['muted']}; font-weight: 600; }}
.info-row .ir-value {{ color: {C['texte']}; font-weight: 700; }}

/* ── SCROLLBAR ── */
::-webkit-scrollbar {{ width: 4px; height: 4px; }}
::-webkit-scrollbar-track {{ background: transparent; }}
::-webkit-scrollbar-thumb {{ background: {C['bordure']}; border-radius: 2px; }}
::-webkit-scrollbar-thumb:hover {{ background: {C['muted']}; }}

/* ── TABS ── */
.stTabs [data-baseweb="tab-list"] {{
    gap: 4px;
    background: {C['gris_bg']};
    border-radius: 12px;
    padding: 4px;
    border: 1px solid {C['bordure']};
}}
.stTabs [data-baseweb="tab"] {{
    border-radius: 8px;
    font-weight: 600;
    font-size: 0.82rem;
    color: {C['muted']};
    padding: 8px 16px;
    border: none;
    background: transparent;
}}
.stTabs [aria-selected="true"] {{
    background: {C['blanc']} !important;
    color: {C['bleu']} !important;
    box-shadow: 0 1px 4px rgba(0,0,0,0.08);
}}

/* ── DATAFRAME ── */
[data-testid="stDataFrameResizable"] {{
    border: 1px solid {C['bordure']};
    border-radius: 12px;
    overflow: hidden;
}}

/* ── HIDE STREAMLIT CHROME ── */
#MainMenu, footer, header {{ visibility: hidden; }}
.stDeployButton {{ display: none; }}


/* ── SIDEBAR ÉPURÉE ── */
[data-testid="stSidebar"] {{
    background: {C['sidebar']} !important;
    border-right: none;
    min-width: 240px !important;
    max-width: 260px !important;
}}
[data-testid="stSidebar"] * {{ color: #E2E8F0 !important; }}
[data-testid="stSidebar"] .stRadio div[role="radiogroup"] label {{
    background: rgba(255,255,255,0.04); border-radius: 8px;
    padding: 9px 14px !important; margin: 3px 0 !important;
    border: 1px solid transparent; font-size: 0.85rem !important;
    font-weight: 500 !important; color: #CBD5E1 !important; transition: all 0.15s;
}}
[data-testid="stSidebar"] .stRadio div[role="radiogroup"] label:hover {{
    background: rgba(97,114,243,0.2) !important;
    border-color: rgba(97,114,243,0.5) !important;
}}

/* ── HEADER NAVIGATION HORIZONTAL ── */
.top-navbar {{
    background: {C['sidebar']};
    border-radius: 14px;
    padding: 0 20px;
    margin-bottom: 18px;
    display: flex;
    align-items: center;
    justify-content: space-between;
    box-shadow: 0 2px 12px rgba(0,0,0,0.15);
    border: 1px solid rgba(255,255,255,0.06);
    overflow-x: auto;
    gap: 4px;
}}
.nav-brand {{
    display: flex;
    align-items: center;
    gap: 10px;
    padding: 14px 0;
    min-width: 160px;
    white-space: nowrap;
}}
.nav-brand .logo {{ font-size: 1.6rem; }}
.nav-brand .name {{
    font-size: 0.9rem; font-weight: 800;
    color: white; letter-spacing: -0.01em; line-height: 1.2;
}}
.nav-brand .tagline {{ font-size: 0.62rem; color: #64748B; font-weight: 500; }}
.nav-links {{
    display: flex; align-items: center; gap: 2px; flex: 1;
    justify-content: center; padding: 8px 0;
}}
.nav-link {{
    display: flex; flex-direction: column; align-items: center;
    padding: 8px 12px; border-radius: 10px; cursor: pointer;
    text-decoration: none; transition: all 0.2s;
    border: 1px solid transparent; white-space: nowrap;
    min-width: 70px;
}}
.nav-link:hover {{
    background: rgba(97,114,243,0.15);
    border-color: rgba(97,114,243,0.3);
}}
.nav-link.active {{
    background: rgba(97,114,243,0.25);
    border-color: rgba(97,114,243,0.5);
}}
.nav-link .nl-ico {{ font-size: 1.1rem; line-height: 1; }}
.nav-link .nl-txt {{
    font-size: 0.58rem; font-weight: 600; color: #94A3B8;
    text-transform: uppercase; letter-spacing: 0.05em; margin-top: 3px;
}}
.nav-link.active .nl-txt {{ color: {C['bleu']}; }}
.nav-right {{
    display: flex; align-items: center; gap: 10px;
    padding: 8px 0; white-space: nowrap;
}}
.nav-date {{
    font-size: 0.72rem; color: #64748B; text-align: right; line-height: 1.4;
}}
.nav-dot {{
    width: 8px; height: 8px; border-radius: 50%;
    background: {C['vert']}; box-shadow: 0 0 6px {C['vert']};
}}

/* ── PAGE ACCUEIL ── */
.welcome-hero {{
    background: linear-gradient(135deg, {C['sidebar']} 0%, #1A3A6B 50%, {C['bleu']} 100%);
    border-radius: 20px;
    padding: 52px 48px;
    margin-bottom: 28px;
    position: relative;
    overflow: hidden;
    box-shadow: 0 8px 40px rgba(97,114,243,0.3);
}}
.welcome-hero::before {{
    content: "📦";
    position: absolute;
    right: 60px;
    top: 50%;
    transform: translateY(-50%);
    font-size: 160px;
    opacity: 0.07;
    pointer-events: none;
}}
.welcome-hero::after {{
    content: "";
    position: absolute;
    top: -60px; right: -60px;
    width: 300px; height: 300px;
    border-radius: 50%;
    background: radial-gradient(circle, rgba(97,114,243,0.3) 0%, transparent 70%);
    pointer-events: none;
}}
.hero-badge {{
    display: inline-block;
    background: rgba(97,114,243,0.3);
    border: 1px solid rgba(97,114,243,0.5);
    border-radius: 20px;
    padding: 4px 14px;
    font-size: 0.72rem;
    font-weight: 700;
    color: {C['bleu_clair']};
    text-transform: uppercase;
    letter-spacing: 0.08em;
    margin-bottom: 16px;
}}
.hero-title {{
    font-size: 2.6rem;
    font-weight: 900;
    color: white;
    margin: 0 0 10px;
    letter-spacing: -0.04em;
    line-height: 1.1;
}}
.hero-sub {{
    font-size: 1.05rem;
    color: rgba(255,255,255,0.70);
    font-weight: 400;
    margin-bottom: 28px;
    line-height: 1.5;
    max-width: 600px;
}}
.hero-stats {{
    display: flex;
    gap: 24px;
    flex-wrap: wrap;
}}
.hero-stat {{
    background: rgba(255,255,255,0.08);
    backdrop-filter: blur(8px);
    border: 1px solid rgba(255,255,255,0.15);
    border-radius: 12px;
    padding: 12px 20px;
    text-align: center;
    min-width: 100px;
}}
.hero-stat .hs-val {{
    font-size: 1.5rem; font-weight: 800; color: white; display: block; line-height: 1;
}}
.hero-stat .hs-lbl {{
    font-size: 0.65rem; color: rgba(255,255,255,0.55);
    text-transform: uppercase; letter-spacing: 0.07em; margin-top: 4px; display: block;
}}

/* ── CARDS DE MODULES ── */
.module-card {{
    background: {C['blanc']};
    border-radius: 14px;
    border: 1px solid {C['bordure']};
    padding: 22px 22px 18px;
    box-shadow: 0 2px 8px rgba(0,0,0,0.06);
    transition: transform 0.2s, box-shadow 0.2s;
    height: 100%;
    position: relative;
    overflow: hidden;
}}
.module-card:hover {{
    transform: translateY(-3px);
    box-shadow: 0 8px 24px rgba(0,0,0,0.12);
}}
.module-card::before {{
    content: '';
    position: absolute;
    top: 0; left: 0;
    width: 100%; height: 4px;
    border-radius: 14px 14px 0 0;
}}
.mc-bleu::before   {{ background: {C['bleu']}; }}
.mc-vert::before   {{ background: {C['vert']}; }}
.mc-rouge::before  {{ background: {C['rouge']}; }}
.mc-orange::before {{ background: {C['orange']}; }}
.mc-violet::before {{ background: {C['violet']}; }}
.mc-rose::before   {{ background: {C['rose']}; }}
.mc-or::before     {{ background: {C['or']}; }}
.module-card .mc-ico {{
    font-size: 2rem; margin-bottom: 12px; display: block;
}}
.module-card .mc-title {{
    font-size: 0.95rem; font-weight: 700; color: {C['texte']}; margin-bottom: 6px;
}}
.module-card .mc-desc {{
    font-size: 0.78rem; color: {C['muted']}; line-height: 1.5;
}}
.module-card .mc-tags {{
    display: flex; flex-wrap: wrap; gap: 5px; margin-top: 12px;
}}
.mc-tag {{
    background: {C['gris_bg']}; border: 1px solid {C['bordure']};
    border-radius: 20px; padding: 2px 8px;
    font-size: 0.65rem; font-weight: 600; color: {C['muted']};
    text-transform: uppercase; letter-spacing: 0.05em;
}}

/* ── GUIDE RAPIDE ── */
.guide-step {{
    display: flex; gap: 14px; align-items: flex-start;
    background: {C['blanc']}; border: 1px solid {C['bordure']};
    border-radius: 12px; padding: 16px 18px; margin-bottom: 10px;
    box-shadow: 0 1px 4px rgba(0,0,0,0.04);
}}
.gs-num {{
    min-width: 30px; height: 30px;
    background: {C['bleu']}; border-radius: 50%;
    display: flex; align-items: center; justify-content: center;
    font-size: 0.78rem; font-weight: 800; color: white; flex-shrink: 0;
}}
.gs-content .gs-title {{
    font-size: 0.88rem; font-weight: 700; color: {C['texte']}; margin-bottom: 3px;
}}
.gs-content .gs-desc {{
    font-size: 0.78rem; color: {C['muted']}; line-height: 1.45;
}}

/* ── CONTACTS & INFOS ── */
.info-row {{
    display: flex; justify-content: space-between; align-items: center;
    padding: 9px 0; border-bottom: 1px solid {C['bordure']};
    font-size: 0.82rem;
}}
.info-row:last-child {{ border-bottom: none; }}
.info-row .ir-label {{ color: {C['muted']}; font-weight: 600; }}
.info-row .ir-value {{ color: {C['texte']}; font-weight: 500; }}

/* ── SCROLL BAR DISCRÈTE ── */
::-webkit-scrollbar {{ width: 5px; height: 5px; }}
::-webkit-scrollbar-track {{ background: transparent; }}
::-webkit-scrollbar-thumb {{ background: {C['bordure']}; border-radius: 3px; }}

</style>
""", unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════
# HELPERS FONCTIONS
# ══════════════════════════════════════════════════════════════════
def excel_style(fig, height=320, show_legend=True):
    """Style unifié pour tous les graphiques — fond blanc, grille légère."""
    fig.update_layout(
        paper_bgcolor="#FFFFFF",
        plot_bgcolor="#FFFFFF",
        font=dict(family="'Plus Jakarta Sans', sans-serif", color=C["texte"], size=11),
        height=height,
        margin=dict(t=24, b=44, l=8, r=16),
        legend=dict(
            bgcolor="rgba(0,0,0,0)",
            font=dict(color=C["texte"], size=10, family="'Plus Jakarta Sans'"),
            orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1,
        ),
        hoverlabel=dict(
            bgcolor="white", font_size=12, bordercolor=C["bordure"],
            font_family="'Plus Jakarta Sans'",
        ),
        showlegend=show_legend,
    )
    fig.update_xaxes(
        gridcolor="#F1F5F9", showgrid=True, zeroline=False,
        linecolor=C["bordure"], tickfont=dict(color=C["muted"], size=10),
    )
    fig.update_yaxes(
        gridcolor="#F1F5F9", showgrid=True, zeroline=False,
        linecolor=C["bordure"], tickfont=dict(color=C["muted"], size=10),
    )
    return fig


def kpi(icon, label, value, sub="", color="bleu"):
    st.markdown(f"""
    <div class="kpi {color}">
        <span class="ico">{icon}</span>
        <div class="lbl">{label}</div>
        <div class="val">{value}</div>
        <div class="sub">{sub}</div>
    </div>""", unsafe_allow_html=True)


def prev_kpi(icon, label, value, sub=""):
    st.markdown(f"""
    <div class="prev-kpi">
        <span class="ico">{icon}</span>
        <div class="lbl">{label}</div>
        <div class="val">{value}</div>
        <div class="sub">{sub}</div>
    </div>""", unsafe_allow_html=True)


def section(title):
    st.markdown(f'<div class="sec">{title}</div>', unsafe_allow_html=True)


def comment(texte, titre="📌 Ce que ça signifie pour vous"):
    st.markdown(f"""
    <div class="comment-box">
        <div class="ct">{titre}</div>
        {texte}
    </div>""", unsafe_allow_html=True)


def prevision_comment(texte, titre="🔮 Prévision & Anticipation"):
    st.markdown(f"""
    <div class="prevision-box">
        <div class="ct">{titre}</div>
        {texte}
    </div>""", unsafe_allow_html=True)


def fmt(v, suffix="FCFA"):
    if v >= 1_000_000_000: return f"{v/1_000_000_000:.2f} Md {suffix}"
    if v >= 1_000_000:     return f"{v/1_000_000:.2f} M {suffix}"
    if v >= 1_000:         return f"{v/1_000:.1f} K {suffix}"
    return f"{v:,.0f} {suffix}"


def prog_bar(label, value, max_val, color=None):
    color = color or C["bleu"]
    pct = min(100, value / max_val * 100) if max_val else 0
    st.markdown(f"""
    <div class="prog">
        <div class="pl"><span>{label}</span><span><b>{value:,.0f} FCFA</b></span></div>
        <div class="pb"><div class="pf" style="width:{pct:.1f}%;background:{color};"></div></div>
    </div>""", unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════
# MODÈLE DE PRÉVISION (régression polynomiale deg=2)
# ══════════════════════════════════════════════════════════════════
def prevoir_ca(serie: pd.Series, n: int = 6):
    y  = serie.values.astype(float)
    X  = np.arange(len(y)).reshape(-1, 1)
    poly = PolynomialFeatures(degree=2)
    Xp = poly.fit_transform(X)
    mdl = LinearRegression().fit(Xp, y)
    std = (y - mdl.predict(Xp)).std()
    X_fut = np.arange(len(y), len(y) + n).reshape(-1, 1)
    y_fut = np.maximum(mdl.predict(poly.transform(X_fut)), 0)
    last  = serie.index[-1]
    peri  = [str(pd.Period(last, "M") + i + 1) for i in range(n)]
    return pd.DataFrame({"Période": peri, "CA_Prévu": y_fut,
                          "Bas": np.maximum(y_fut - 1.5*std, 0),
                          "Haut": y_fut + 1.5*std})

@st.cache_data(show_spinner=False)
def generer_donnees():
    import random as rnd
    rnd.seed(42)
    np.random.seed(42)

    PRODUITS = [
        # Sacs supermarché
        {"ref":"SAC-SM-001","designation":"Sac Supermarché 30L Standard",      "categorie":"Sacs Supermarché",        "prix_achat":1800,  "prix_vente":2500,  "seuil":200},
        {"ref":"SAC-SM-002","designation":"Sac Supermarché 50L Renforcé",       "categorie":"Sacs Supermarché",        "prix_achat":2600,  "prix_vente":3500,  "seuil":200},
        {"ref":"SAC-SM-003","designation":"Sac Supermarché 20L Mini",           "categorie":"Sacs Supermarché",        "prix_achat":1400,  "prix_vente":2000,  "seuil":300},
        {"ref":"SAC-SM-004","designation":"Sac Supermarché HD Noir 40L",        "categorie":"Sacs Supermarché",        "prix_achat":2200,  "prix_vente":3000,  "seuil":150},
        {"ref":"SAC-SM-005","designation":"Sac Caisse Transparent 25L",         "categorie":"Sacs Supermarché",        "prix_achat":4500,  "prix_vente":6200,  "seuil":100},
        {"ref":"SAC-SM-006","designation":"Sac Biosourcé Compostable 30L",      "categorie":"Sacs Supermarché",        "prix_achat":5200,  "prix_vente":7500,  "seuil":80},
        # Sacs poubelles
        {"ref":"SAC-PB-001","designation":"Sac Poubelle 30L Vert",              "categorie":"Sacs Poubelles",          "prix_achat":650,   "prix_vente":950,   "seuil":500},
        {"ref":"SAC-PB-002","designation":"Sac Poubelle 50L Noir Renforcé",     "categorie":"Sacs Poubelles",          "prix_achat":850,   "prix_vente":1200,  "seuil":500},
        {"ref":"SAC-PB-003","designation":"Sac Poubelle 110L Industrial",       "categorie":"Sacs Poubelles",          "prix_achat":1200,  "prix_vente":1700,  "seuil":300},
        {"ref":"SAC-PB-004","designation":"Sac Poubelle 20L Bleu Ménager",      "categorie":"Sacs Poubelles",          "prix_achat":550,   "prix_vente":800,   "seuil":400},
        {"ref":"SAC-PB-005","designation":"Sac Poubelle 240L Collectivité",     "categorie":"Sacs Poubelles",          "prix_achat":2200,  "prix_vente":3000,  "seuil":100},
        # Gobelets
        {"ref":"GOB-001",   "designation":"Gobelet Plastique 20cl Standard",    "categorie":"Gobelets",                "prix_achat":3500,  "prix_vente":5000,  "seuil":150},
        {"ref":"GOB-002",   "designation":"Gobelet Plastique 33cl Grand",       "categorie":"Gobelets",                "prix_achat":2800,  "prix_vente":4000,  "seuil":150},
        {"ref":"GOB-003",   "designation":"Gobelet Cristal 25cl",               "categorie":"Gobelets",                "prix_achat":3200,  "prix_vente":4500,  "seuil":100},
        {"ref":"GOB-004",   "designation":"Gobelet Solo Cup 50cl",              "categorie":"Gobelets",                "prix_achat":2400,  "prix_vente":3500,  "seuil":80},
        {"ref":"GOB-005",   "designation":"Gobelet Dégustation 10cl",           "categorie":"Gobelets",                "prix_achat":3800,  "prix_vente":5500,  "seuil":120},
        # Emballages alimentaires
        {"ref":"EMB-AL-001","designation":"Barquette Alimentaire PP 500ml",     "categorie":"Emballages Alimentaires", "prix_achat":4200,  "prix_vente":6000,  "seuil":100},
        {"ref":"EMB-AL-002","designation":"Barquette Alimentaire PP 1L",        "categorie":"Emballages Alimentaires", "prix_achat":3600,  "prix_vente":5200,  "seuil":80},
        {"ref":"EMB-AL-003","designation":"Film Étirable Alimentaire 300m",     "categorie":"Emballages Alimentaires", "prix_achat":5500,  "prix_vente":8000,  "seuil":50},
        {"ref":"EMB-AL-004","designation":"Sachet Zip PP 1L",                   "categorie":"Emballages Alimentaires", "prix_achat":1800,  "prix_vente":2600,  "seuil":200},
        # Emballages industriels
        {"ref":"EMB-IN-001","designation":"Big Bag Industriel 1T FIBC",         "categorie":"Emballages Industriels",  "prix_achat":8500,  "prix_vente":12000, "seuil":30},
        {"ref":"EMB-IN-002","designation":"Fût Plastique HDPE 220L",            "categorie":"Emballages Industriels",  "prix_achat":12000, "prix_vente":17000, "seuil":20},
        {"ref":"EMB-IN-003","designation":"Jerrycan HDPE 10L",                  "categorie":"Emballages Industriels",  "prix_achat":2800,  "prix_vente":4000,  "seuil":60},
        {"ref":"EMB-IN-004","designation":"Bidon HDPE 5L avec bouchon",         "categorie":"Emballages Industriels",  "prix_achat":1600,  "prix_vente":2400,  "seuil":80},
        {"ref":"EMB-IN-005","designation":"Sac Industriel PE 50Kg",             "categorie":"Emballages Industriels",  "prix_achat":6500,  "prix_vente":9000,  "seuil":40},
        # Matières premières
        {"ref":"MAT-PE-001","designation":"Granulés PE-HD",                     "categorie":"Matières Premières",      "prix_achat":14000, "prix_vente":None,  "seuil":100},
        {"ref":"MAT-PP-001","designation":"Granulés PP",                        "categorie":"Matières Premières",      "prix_achat":15500, "prix_vente":None,  "seuil":80},
        {"ref":"MAT-PVC-001","designation":"Granulés PVC Souple",               "categorie":"Matières Premières",      "prix_achat":18000, "prix_vente":None,  "seuil":60},
        {"ref":"MAT-REC-001","designation":"Plastique Recyclé Mix",             "categorie":"Matières Premières",      "prix_achat":7500,  "prix_vente":None,  "seuil":50},
        {"ref":"MAT-COL-001","designation":"Masterbatch Noir",                  "categorie":"Matières Premières",      "prix_achat":22000, "prix_vente":None,  "seuil":40},
        {"ref":"MAT-COL-002","designation":"Masterbatch Vert",                  "categorie":"Matières Premières",      "prix_achat":23000, "prix_vente":None,  "seuil":35},
    ]

    # ─── 13 SEGMENTS · 60+ CLIENTS · 8 ZONES ──────────────────────
    CLIENTS = [
        # GRANDE DISTRIBUTION (zone Sud/Nord)
        {"nom":"SOCOCÉ Abidjan Plateau",         "segment":"Grande Distribution",      "zone":"Zone Sud",    "ville":"Abidjan"},
        {"nom":"Carrefour Market Cocody",         "segment":"Grande Distribution",      "zone":"Zone Sud",    "ville":"Abidjan"},
        {"nom":"PROSUMA SA",                      "segment":"Grande Distribution",      "zone":"Zone Sud",    "ville":"Abidjan"},
        {"nom":"Leader Price Yopougon",           "segment":"Grande Distribution",      "zone":"Zone Nord",   "ville":"Abidjan"},
        {"nom":"CDCI Supermarché Plateau",        "segment":"Grande Distribution",      "zone":"Zone Nord",   "ville":"Abidjan"},
        {"nom":"Hypermarché Prima",               "segment":"Grande Distribution",      "zone":"Zone Sud",    "ville":"Abidjan"},
        # INDUSTRIEL AGROALIMENTAIRE
        {"nom":"NESTLE CI",                       "segment":"Industriel Agroalimentaire","zone":"Zone Sud",   "ville":"Abidjan"},
        {"nom":"SIC CACAOS",                      "segment":"Industriel Agroalimentaire","zone":"Zone Sud",   "ville":"Abidjan"},
        {"nom":"SIFCA Groupe",                    "segment":"Industriel Agroalimentaire","zone":"Zone Sud",   "ville":"Abidjan"},
        {"nom":"PALMCI San-Pédro",                "segment":"Industriel Agroalimentaire","zone":"Zone Ouest", "ville":"San-Pédro"},
        {"nom":"Ivoire Sucre",                    "segment":"Industriel Agroalimentaire","zone":"Zone Centre","ville":"Bouaké"},
        {"nom":"BLOHORN SA",                      "segment":"Industriel Agroalimentaire","zone":"Zone Sud",   "ville":"Abidjan"},
        # INDUSTRIEL MANUFACTURIER
        {"nom":"SOLIBRA Brasserie",               "segment":"Industriel Manufacturier", "zone":"Zone Nord",   "ville":"Abidjan"},
        {"nom":"GONFREVILLE Textile",             "segment":"Industriel Manufacturier", "zone":"Zone Nord",   "ville":"Abidjan"},
        {"nom":"CIE Côte d'Ivoire",               "segment":"Industriel Manufacturier", "zone":"Zone Sud",    "ville":"Abidjan"},
        {"nom":"SITARAIL CI",                     "segment":"Industriel Manufacturier", "zone":"Zone Centre", "ville":"Bouaké"},
        # DISTRIBUTEUR LOCAL ABIDJAN
        {"nom":"Commerce Général Yop",            "segment":"Distributeur Local Abidjan","zone":"Zone Nord",  "ville":"Abidjan"},
        {"nom":"Maquis & Resto Supplies",         "segment":"Distributeur Local Abidjan","zone":"Zone Sud",   "ville":"Abidjan"},
        {"nom":"NEGOCE MARCORY",                  "segment":"Distributeur Local Abidjan","zone":"Zone Sud",   "ville":"Abidjan"},
        {"nom":"DIST ABOBO EXPRESS",              "segment":"Distributeur Local Abidjan","zone":"Zone Nord",  "ville":"Abidjan"},
        {"nom":"COMMERCE ADJAMÉ",                 "segment":"Distributeur Local Abidjan","zone":"Zone Nord",  "ville":"Abidjan"},
        # DISTRIBUTEUR RÉGIONAL CI
        {"nom":"DIOULA NÉGOCE Bouaké",            "segment":"Distributeur Régional CI", "zone":"Zone Centre", "ville":"Bouaké"},
        {"nom":"AGRO DIST Korhogo",               "segment":"Distributeur Régional CI", "zone":"Zone Nord",   "ville":"Korhogo"},
        {"nom":"TRANSIT SAN-PEDRO",               "segment":"Distributeur Régional CI", "zone":"Zone Ouest",  "ville":"San-Pédro"},
        {"nom":"COMMERCE ABENGOUROU",             "segment":"Distributeur Régional CI", "zone":"Zone Est",    "ville":"Abengourou"},
        {"nom":"DIST DALOA CENTER",               "segment":"Distributeur Régional CI", "zone":"Zone Centre", "ville":"Daloa"},
        {"nom":"NÉGOCE MAN OUEST",                "segment":"Distributeur Régional CI", "zone":"Zone Ouest",  "ville":"Man"},
        # HÔTELLERIE & RESTAURATION
        {"nom":"SOFITEL Abidjan",                 "segment":"Hôtellerie & Restauration","zone":"Zone Sud",    "ville":"Abidjan"},
        {"nom":"Hôtel Ivoire Intercontinental",   "segment":"Hôtellerie & Restauration","zone":"Zone Sud",    "ville":"Abidjan"},
        {"nom":"Groupe LAICO",                    "segment":"Hôtellerie & Restauration","zone":"Zone Sud",    "ville":"Abidjan"},
        {"nom":"Palm Club Hôtel",                 "segment":"Hôtellerie & Restauration","zone":"Zone Sud",    "ville":"Abidjan"},
        {"nom":"Restaurant Le Phare",             "segment":"Hôtellerie & Restauration","zone":"Zone Sud",    "ville":"Abidjan"},
        # SANTÉ & PHARMACIE
        {"nom":"CHU de Treichville",              "segment":"Santé & Pharmacie",        "zone":"Zone Sud",    "ville":"Abidjan"},
        {"nom":"Clinique Biasa",                  "segment":"Santé & Pharmacie",        "zone":"Zone Nord",   "ville":"Abidjan"},
        {"nom":"LABOREX CI",                      "segment":"Santé & Pharmacie",        "zone":"Zone Sud",    "ville":"Abidjan"},
        {"nom":"Pharmivoire Groupe",              "segment":"Santé & Pharmacie",        "zone":"Zone Sud",    "ville":"Abidjan"},
        # ADMINISTRATION & COLLECTIVITÉS
        {"nom":"Mairie Cocody",                   "segment":"Administration & Collectivités","zone":"Zone Sud","ville":"Abidjan"},
        {"nom":"Mairie de Bouaké",                "segment":"Administration & Collectivités","zone":"Zone Centre","ville":"Bouaké"},
        {"nom":"ANADER CI",                       "segment":"Administration & Collectivités","zone":"Zone Sud","ville":"Abidjan"},
        {"nom":"ONAD (Assainissement)",           "segment":"Administration & Collectivités","zone":"Zone Sud","ville":"Abidjan"},
        # BTP & CONSTRUCTION
        {"nom":"SOGEBAC Construction",            "segment":"BTP & Construction",       "zone":"Zone Sud",    "ville":"Abidjan"},
        {"nom":"COLAS CI",                        "segment":"BTP & Construction",       "zone":"Zone Sud",    "ville":"Abidjan"},
        {"nom":"BATIPRO Group",                   "segment":"BTP & Construction",       "zone":"Zone Nord",   "ville":"Abidjan"},
        # LOGISTIQUE & TRANSPORT
        {"nom":"BOLLORÉ Logistics CI",            "segment":"Logistique & Transport",   "zone":"Zone Sud",    "ville":"Abidjan"},
        {"nom":"SDV Port Abidjan",                "segment":"Logistique & Transport",   "zone":"Zone Sud",    "ville":"Abidjan"},
        {"nom":"MAERSK CI",                       "segment":"Logistique & Transport",   "zone":"Zone Sud",    "ville":"Abidjan"},
        {"nom":"TransCi Freight",                 "segment":"Logistique & Transport",   "zone":"Zone Ouest",  "ville":"San-Pédro"},
        # AGRICULTURE & AGRO-INDUSTRIE
        {"nom":"SOGB Rubber",                     "segment":"Agriculture & Agro-Industrie","zone":"Zone Ouest","ville":"San-Pédro"},
        {"nom":"SAPH Hévéa CI",                   "segment":"Agriculture & Agro-Industrie","zone":"Zone Ouest","ville":"Aboisso"},
        {"nom":"COOPAGRI Daloa",                  "segment":"Agriculture & Agro-Industrie","zone":"Zone Centre","ville":"Daloa"},
        # EXPORT SOUS-RÉGIONAL
        {"nom":"WEST AFRICA TRADING CO.",         "segment":"Export Sous-Régional",     "zone":"Export",      "ville":"Dakar"},
        {"nom":"GHANA IMPORT GROUP",              "segment":"Export Sous-Régional",     "zone":"Export",      "ville":"Accra"},
        {"nom":"MALI PACKAGING SARL",             "segment":"Export Sous-Régional",     "zone":"Export",      "ville":"Bamako"},
        {"nom":"BURKINA EMBALLAGE",               "segment":"Export Sous-Régional",     "zone":"Export",      "ville":"Ouagadougou"},
        {"nom":"TOGO DISTRIBUTION",               "segment":"Export Sous-Régional",     "zone":"Export",      "ville":"Lomé"},
        # ARTISANAT & PME
        {"nom":"COOPEC Artisans Yop",             "segment":"Artisanat & PME",          "zone":"Zone Nord",   "ville":"Abidjan"},
        {"nom":"PME Couture Mode CI",             "segment":"Artisanat & PME",          "zone":"Zone Sud",    "ville":"Abidjan"},
        {"nom":"Traiteur Événements Prestige",    "segment":"Artisanat & PME",          "zone":"Zone Sud",    "ville":"Abidjan"},
        {"nom":"MARCHÉ ADJAMÉ GROSSISTE",         "segment":"Artisanat & PME",          "zone":"Zone Nord",   "ville":"Abidjan"},
    ]

    MODES = ["Virement bancaire","Chèque certifié","Espèce","Carte bancaire","Traite"]
    POIDS = [0.40, 0.25, 0.18, 0.10, 0.07]
    SAISON= [0.75,0.70,0.80,0.85,0.90,0.97,0.88,1.12,1.22,1.30,1.45,1.38]

    SEG_W = {
        "Grande Distribution":3.5, "Industriel Agroalimentaire":3.2,
        "Industriel Manufacturier":2.8, "Distributeur Local Abidjan":1.8,
        "Distributeur Régional CI":1.5, "Hôtellerie & Restauration":1.2,
        "Santé & Pharmacie":1.0, "Administration & Collectivités":1.3,
        "BTP & Construction":1.6, "Logistique & Transport":1.4,
        "Agriculture & Agro-Industrie":2.0, "Export Sous-Régional":2.5,
        "Artisanat & PME":0.8,
    }
    P_PAY = {
        "Grande Distribution":0.90, "Industriel Agroalimentaire":0.87,
        "Industriel Manufacturier":0.85, "Distributeur Local Abidjan":0.75,
        "Distributeur Régional CI":0.70, "Hôtellerie & Restauration":0.92,
        "Santé & Pharmacie":0.95, "Administration & Collectivités":0.60,
        "BTP & Construction":0.72, "Logistique & Transport":0.88,
        "Agriculture & Agro-Industrie":0.80, "Export Sous-Régional":0.78,
        "Artisanat & PME":0.65,
    }

    start = datetime(2022, 1, 1)
    pv = [p for p in PRODUITS if p["prix_vente"]]
    factures = []
    for i in range(1200):  # 1200 factures vs 850 avant
        cl  = rnd.choice(CLIENTS)
        d   = start + timedelta(days=rnd.randint(0, 1094))
        co  = SAISON[d.month-1] * SEG_W.get(cl["segment"], 1.5)
        nb  = rnd.choices([1,2,3,4], weights=[0.25,0.38,0.27,0.10])[0]
        sel = rnd.sample(pv, min(nb, len(pv)))
        mht = round(sum(p["prix_vente"] * rnd.randint(1, int(10*co)) for p in sel), 0)
        tva = round(mht * 0.18, 0)
        ttc = mht + tva
        ok  = rnd.random() < P_PAY.get(cl["segment"], 0.80)
        mode = rnd.choices(MODES, weights=POIDS)[0] if ok else ""
        factures.append({
            "N° FACTURE":       f"F{d.year}-{str(d.month).zfill(2)}-{str(i+1).zfill(4)}",
            "DATE":             d, "ANNEE": d.year, "MOIS": d.month,
            "CLIENT":           cl["nom"], "SEGMENT": cl["segment"],
            "ZONE":             cl["zone"], "VILLE": cl["ville"],
            "MONTANT HT":       mht, "TVA": tva, "MONTANT TTC": ttc,
            "ETAT DE PAIEMENT": "Payée" if ok else "Impayée",
            "MODE DE PAIEMENT": mode,
        })
    df_f = pd.DataFrame(factures).sort_values("DATE").reset_index(drop=True)

    # Stock
    inv = []
    for p in PRODUITS:
        si = rnd.randint(30, 400)
        inv.append({**p, "Stock initial": si, "Entrées": 0, "Sorties": 0, "Stock final": si})
    df_inv = pd.DataFrame(inv)

    ent_rows, sor_rows = [], []
    for p in PRODUITS:
        for _ in range(rnd.randint(12, 45)):
            d = start + timedelta(days=rnd.randint(0, 1094))
            qte = rnd.randint(10, 600)
            cout= round(p["prix_achat"] * rnd.uniform(0.91, 1.06), 0)
            ent_rows.append({"Date":d,"Référence":p["ref"],"Désignation":p["designation"],
                             "Catégorie":p["categorie"],"Coût d'achat":cout,"Quantité":qte,"Total":round(cout*qte,0)})
        if p["prix_vente"]:
            for _ in range(rnd.randint(8, 40)):
                d = start + timedelta(days=rnd.randint(0, 1094))
                qte = rnd.randint(5, 350)
                pv_ = round(p["prix_vente"] * rnd.uniform(0.97, 1.10), 0)
                sor_rows.append({"Date":d,"Référence":p["ref"],"Désignation":p["designation"],
                                 "Catégorie":p["categorie"],"Prix de vente":pv_,"Quantité":qte,"Total":round(pv_*qte,0)})
    df_ent = pd.DataFrame(ent_rows).sort_values("Date").reset_index(drop=True)
    df_sor = pd.DataFrame(sor_rows).sort_values("Date").reset_index(drop=True)

    for i, row in df_inv.iterrows():
        ref = row["ref"]
        e = df_ent[df_ent["Référence"]==ref]["Quantité"].sum()
        s = df_sor[df_sor["Référence"]==ref]["Quantité"].sum() if ref in df_sor["Référence"].values else 0
        sf= max(0, int(row["Stock initial"]) + int(e) - int(s))
        df_inv.at[i,"Entrées"] = int(e); df_inv.at[i,"Sorties"] = int(s)
        df_inv.at[i,"Stock final"] = sf
        df_inv.at[i,"Valeur"] = round(sf * row["prix_achat"], 0)
        df_inv.at[i,"Sorties_moy_mois"] = round(s / 36, 1)
        df_inv.at[i,"Statut"] = ("Non disponible" if sf==0 else "Stock faible" if sf<row["seuil"] else "Stock normal")

    machines = ["Presse 1 – Injection","Presse 2 – Soufflage","Extrudeuse A","Extrudeuse B","Thermoformeuse"]
    prod_rows = []
    for m in range(36):
        d = start + pd.DateOffset(months=m)
        co = SAISON[d.month-1]
        for mach in machines:
            r  = rnd.uniform(0.73, 0.97); pp = int(rnd.randint(50000,140000)*co)
            pr = int(pp*r); reb = int(pr*rnd.uniform(0.01,0.055))
            prod_rows.append({
                "Mois":d.strftime("%Y-%m"),"Mois_Label":d.strftime("%b %Y"),"Année":d.year,"Machine":mach,
                "Production planifiée":pp,"Production réelle":pr,"Rebuts":reb,
                "Taux rendement":round(r*100,1),"Taux rebut":round(reb/pr*100,2) if pr>0 else 0,
            })
    df_prod = pd.DataFrame(prod_rows)
    return df_f, df_inv, df_ent, df_sor, df_prod

# ══════════════════════════════════════════════════════════════════
# GÉNÉRATION RAPPORT POWERPOINT — MULTIPACK SA
# ══════════════════════════════════════════════════════════════════
def generer_rapport_pptx(df_fact, df_inv, df_sor, df_prod, df_f, now):
    """Génère un rapport PowerPoint de synthèse pour la Direction."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    import io as _io

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Couleurs MULTIPACK
    NAVY   = RGBColor(0x0F, 0x17, 0x2A)
    INDIGO = RGBColor(0x4F, 0x6E, 0xF7)
    WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
    GREEN  = RGBColor(0x10, 0xB9, 0x81)
    RED    = RGBColor(0xEF, 0x44, 0x44)
    AMBER  = RGBColor(0xF5, 0x9E, 0x0B)
    SLATE  = RGBColor(0x64, 0x74, 0x8B)
    LIGHT  = RGBColor(0xF8, 0xFA, 0xFC)

    def blank_slide(prs):
        blank_layout = prs.slide_layouts[6]
        return prs.slides.add_slide(blank_layout)

    def bg(slide, color=NAVY):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def txt(slide, text, x, y, w, h, size=18, bold=False, color=WHITE,
            align=PP_ALIGN.LEFT, wrap=True):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        return tb

    def rect(slide, x, y, w, h, color=INDIGO, radius=0.1):
        from pptx.util import Inches
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def kpi_box(slide, x, y, w, h, label, value, color=INDIGO, sub=""):
        rect(slide, x, y, w, h, LIGHT)
        # Top accent bar
        accent = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0.06))
        accent.fill.solid(); accent.fill.fore_color.rgb = color
        accent.line.fill.background()
        txt(slide, label.upper(), x+0.1, y+0.12, w-0.2, 0.25,
            size=7, bold=True, color=SLATE, align=PP_ALIGN.LEFT)
        txt(slide, value, x+0.1, y+0.35, w-0.2, 0.55,
            size=20, bold=True, color=RGBColor(0x1E,0x29,0x3B), align=PP_ALIGN.LEFT)
        if sub:
            txt(slide, sub, x+0.1, y+0.88, w-0.2, 0.22,
                size=7.5, color=SLATE, align=PP_ALIGN.LEFT)

    def hbar(slide, x, y, label, value_pct, value_txt, bar_color=INDIGO, h=0.28):
        """Petite barre horizontale pour top N."""
        txt(slide, label, x, y, 3.2, h, size=8.5, color=RGBColor(0x1E,0x29,0x3B))
        bg_bar = slide.shapes.add_shape(1, Inches(x+3.3), Inches(y+0.04), Inches(3.0), Inches(h-0.08))
        bg_bar.fill.solid(); bg_bar.fill.fore_color.rgb = RGBColor(0xE2,0xE8,0xF0)
        bg_bar.line.fill.background()
        fill_w = max(0.05, min(3.0, 3.0 * value_pct))
        fill_bar = slide.shapes.add_shape(1, Inches(x+3.3), Inches(y+0.04), Inches(fill_w), Inches(h-0.08))
        fill_bar.fill.solid(); fill_bar.fill.fore_color.rgb = bar_color
        fill_bar.line.fill.background()
        txt(slide, value_txt, x+6.4, y, 1.5, h, size=8, bold=True,
            color=RGBColor(0x1E,0x29,0x3B), align=PP_ALIGN.RIGHT)

    # ── Calculs synthèse ──────────────────────────────────────────
    ca_total  = df_f["MONTANT TTC"].sum()
    ca_paye   = df_f[df_f["ETAT DE PAIEMENT"]=="Payée"]["MONTANT TTC"].sum()
    ca_impaye = df_f[df_f["ETAT DE PAIEMENT"]=="Impayée"]["MONTANT TTC"].sum()
    taux_rec  = ca_paye/ca_total*100 if ca_total else 0
    nb_cl     = df_f["CLIENT"].nunique()
    val_stock = df_inv["Valeur"].sum()
    nb_rupt   = (df_inv["Statut"]=="Non disponible").sum()
    nb_faib   = (df_inv["Statut"]=="Stock faible").sum()
    rend_moy  = df_prod["Taux rendement"].mean()

    top_clients = df_f.groupby("CLIENT")["MONTANT TTC"].sum().sort_values(ascending=False).head(5)
    top_segs    = df_f.groupby("SEGMENT")["MONTANT TTC"].sum().sort_values(ascending=False)

    def fmtp(v):
        if v>=1e9: return f"{v/1e9:.2f} Md FCFA"
        if v>=1e6: return f"{v/1e6:.2f} M FCFA"
        if v>=1e3: return f"{v/1e3:.1f} K FCFA"
        return f"{v:,.0f} FCFA"

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 1 — COUVERTURE
    # ═══════════════════════════════════════════════════════════════
    s = blank_slide(prs); bg(s, NAVY)
    # Gradient overlay rectangle
    r = rect(s, 0, 0, 13.33, 7.5, RGBColor(0x1E,0x29,0x3B))
    # Accent line top
    accent_top = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.12))
    accent_top.fill.solid(); accent_top.fill.fore_color.rgb = INDIGO
    accent_top.line.fill.background()
    # Logo area
    rect(s, 0.5, 0.5, 1.2, 1.2, INDIGO)
    txt(s, "📦", 0.72, 0.6, 0.8, 0.9, size=36, align=PP_ALIGN.CENTER)
    txt(s, "MULTIPACK SA", 1.9, 0.55, 7, 0.7, size=32, bold=True, color=WHITE)
    txt(s, "Zone Industrielle de Yopougon · Abidjan, Côte d'Ivoire", 1.9, 1.22, 9, 0.4,
        size=11, color=RGBColor(0x64,0x74,0x8B))
    # Titre rapport
    rect(s, 0.5, 2.2, 8.5, 0.06, INDIGO)
    txt(s, "RAPPORT DE DIRECTION", 0.5, 2.4, 10, 0.8, size=28, bold=True, color=WHITE)
    txt(s, "Synthèse des performances commerciales · Contrôle de gestion · Prévisions",
        0.5, 3.1, 10, 0.5, size=13, color=RGBColor(0x94,0xA3,0xB8))
    # Date
    txt(s, now.strftime("Généré le %d %B %Y à %H:%M"), 0.5, 6.7, 8, 0.4,
        size=9, color=RGBColor(0x47,0x55,0x69))
    # KPIs couverture
    rect(s, 9.8, 1.8, 3.0, 4.5, RGBColor(0x1E,0x29,0x3B))
    accent2 = s.shapes.add_shape(1, Inches(9.8), Inches(1.8), Inches(0.06), Inches(4.5))
    accent2.fill.solid(); accent2.fill.fore_color.rgb = INDIGO; accent2.line.fill.background()
    txt(s, "CHIFFRES CLÉS", 10.0, 1.95, 2.6, 0.3, size=7, bold=True,
        color=RGBColor(0x64,0x74,0x8B), align=PP_ALIGN.CENTER)
    for i, (lbl, val, col) in enumerate([
        ("CA Total",      fmtp(ca_total), WHITE),
        ("Encaissé",      f"{taux_rec:.0f}%", GREEN),
        ("Impayés",       fmtp(ca_impaye), RGBColor(0xFB,0xBF,0x24)),
        ("Clients",       str(nb_cl), WHITE),
        ("Stock",         fmtp(val_stock), WHITE),
    ]):
        yy = 2.35 + i * 0.72
        txt(s, lbl, 10.0, yy, 2.6, 0.22, size=8, color=RGBColor(0x64,0x74,0x8B),
            align=PP_ALIGN.CENTER)
        txt(s, val, 10.0, yy+0.22, 2.6, 0.42, size=15, bold=True, color=col,
            align=PP_ALIGN.CENTER)

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 2 — SYNTHÈSE EXÉCUTIVE (KPIs)
    # ═══════════════════════════════════════════════════════════════
    s2 = blank_slide(prs); bg(s2, LIGHT)
    rect(s2, 0, 0, 13.33, 0.7, NAVY)
    txt(s2, "📊  SYNTHÈSE EXÉCUTIVE — INDICATEURS CLÉS", 0.3, 0.18, 10, 0.4,
        size=14, bold=True, color=WHITE)
    txt(s2, now.strftime("%d %b %Y"), 11.5, 0.18, 1.5, 0.4,
        size=10, color=RGBColor(0x94,0xA3,0xB8), align=PP_ALIGN.RIGHT)
    # 8 KPIs 4x2
    kpis = [
        ("Chiffre d'Affaires Total", fmtp(ca_total),   INDIGO, f"{len(df_f)} factures"),
        ("Montant Encaissé",         fmtp(ca_paye),     GREEN,  f"Taux : {taux_rec:.1f}%"),
        ("Créances Impayées",        fmtp(ca_impaye),   RED,    f"{(df_f['ETAT DE PAIEMENT']=='Impayée').sum()} factures"),
        ("Clients Actifs",           str(nb_cl),        INDIGO, f"{df_f['SEGMENT'].nunique()} segments"),
        ("Valeur du Stock",          fmtp(val_stock),   INDIGO, f"{len(df_inv)} références"),
        ("Rendement Usine",          f"{rend_moy:.1f}%",GREEN if rend_moy>=80 else RED, "Moy. machines"),
        ("Ruptures de Stock",        str(nb_rupt),      RED if nb_rupt>0 else GREEN, "Produits épuisés"),
        ("Stocks Faibles",           str(nb_faib),      AMBER,  "Sous seuil alerte"),
    ]
    for i, (lbl, val, col, sub) in enumerate(kpis):
        row, col_idx = i // 4, i % 4
        kpi_box(s2, 0.3 + col_idx*3.18, 0.9 + row*1.3, 3.0, 1.18, lbl, val, col, sub)

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 3 — TOP CLIENTS & SEGMENTS
    # ═══════════════════════════════════════════════════════════════
    s3 = blank_slide(prs); bg(s3, LIGHT)
    rect(s3, 0, 0, 13.33, 0.7, NAVY)
    txt(s3, "👥  TOP CLIENTS & SEGMENTS — Répartition du CA", 0.3, 0.18, 10, 0.4,
        size=14, bold=True, color=WHITE)
    # Left: top clients
    rect(s3, 0.3, 0.85, 6.0, 0.38, NAVY)
    txt(s3, "🏆 TOP 5 CLIENTS PAR CA", 0.4, 0.9, 5, 0.28, size=8.5, bold=True, color=WHITE)
    max_cl = top_clients.max() if len(top_clients) else 1
    for i, (cl, v) in enumerate(top_clients.items()):
        hbar(s3, 0.3, 1.35 + i*0.95, cl[:30], v/max_cl, fmtp(v),
             bar_color=INDIGO if i==0 else RGBColor(0x99,0xA7,0xF8))
    # Right: top segments
    rect(s3, 6.8, 0.85, 6.2, 0.38, NAVY)
    txt(s3, "🏷 CA PAR SEGMENT MARCHÉ", 6.9, 0.9, 5.5, 0.28, size=8.5, bold=True, color=WHITE)
    max_seg = top_segs.max() if len(top_segs) else 1
    colors_seg = [INDIGO, GREEN, AMBER, RED, RGBColor(0x7C,0x3A,0xED),
                  RGBColor(0x06,0xB6,0xD4), RGBColor(0xF4,0x72,0xB6),
                  RGBColor(0x3B,0x82,0xF6), RGBColor(0x10,0xB9,0x81), RGBColor(0xF5,0x9E,0x0B)]
    for i, (seg, v) in enumerate(top_segs.head(8).items()):
        hbar(s3, 6.8, 1.35 + i*0.72, seg[:28], v/max_seg, fmtp(v),
             bar_color=colors_seg[i % len(colors_seg)], h=0.26)

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 4 — PAIEMENTS & TRÉSORERIE
    # ═══════════════════════════════════════════════════════════════
    s4 = blank_slide(prs); bg(s4, LIGHT)
    rect(s4, 0, 0, 13.33, 0.7, NAVY)
    txt(s4, "💳  PAIEMENTS & TRÉSORERIE", 0.3, 0.18, 10, 0.4,
        size=14, bold=True, color=WHITE)
    # Taux recouvrement visuel
    rc = GREEN if taux_rec >= 80 else (AMBER if taux_rec >= 65 else RED)
    kpi_box(s4, 0.3,  0.9, 3.0, 1.3, "Taux de Recouvrement", f"{taux_rec:.1f}%", rc, "Objectif : 80%")
    kpi_box(s4, 3.5,  0.9, 3.0, 1.3, "CA Encaissé",    fmtp(ca_paye),   GREEN, "Payements reçus")
    kpi_box(s4, 6.7,  0.9, 3.0, 1.3, "CA Impayé",      fmtp(ca_impaye), RED,  "À recouvrer")
    kpi_box(s4, 9.9,  0.9, 3.1, 1.3, "Panier Moyen",   fmtp(df_f["MONTANT TTC"].mean()), INDIGO, "Par facture")
    # Modes de paiement
    rect(s4, 0.3, 2.5, 6.0, 0.38, NAVY)
    txt(s4, "💰 RÉPARTITION PAR MODE DE PAIEMENT", 0.4, 2.55, 5.5, 0.28, size=8.5, bold=True, color=WHITE)
    df_pay = df_f[df_f["ETAT DE PAIEMENT"]=="Payée"]
    modes = df_pay.groupby("MODE DE PAIEMENT")["MONTANT TTC"].sum().sort_values(ascending=False)
    max_m = modes.max() if len(modes) else 1
    for i, (m, v) in enumerate(modes.head(5).items()):
        hbar(s4, 0.3, 3.05+i*0.72, m, v/max_m, fmtp(v),
             bar_color=colors_seg[i%len(colors_seg)], h=0.26)
    # Top impayés
    rect(s4, 6.8, 2.5, 6.2, 0.38, NAVY)
    txt(s4, "⚠️ TOP CLIENTS IMPAYÉS", 6.9, 2.55, 5.5, 0.28, size=8.5, bold=True, color=WHITE)
    imp = df_f[df_f["ETAT DE PAIEMENT"]=="Impayée"].groupby("CLIENT")["MONTANT TTC"].sum().sort_values(ascending=False)
    max_imp = imp.max() if len(imp) else 1
    for i, (cl, v) in enumerate(imp.head(5).items()):
        hbar(s4, 6.8, 3.05+i*0.72, cl[:28], v/max_imp, fmtp(v), bar_color=RED, h=0.26)

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 5 — STOCKS & PRODUCTION
    # ═══════════════════════════════════════════════════════════════
    s5 = blank_slide(prs); bg(s5, LIGHT)
    rect(s5, 0, 0, 13.33, 0.7, NAVY)
    txt(s5, "📦  STOCKS & PRODUCTION — État opérationnel", 0.3, 0.18, 10, 0.4,
        size=14, bold=True, color=WHITE)
    # Stocks KPIs
    kpi_box(s5, 0.3, 0.9, 3.0, 1.2, "Valeur Totale Stock", fmtp(val_stock), INDIGO, f"{len(df_inv)} références")
    kpi_box(s5, 3.5, 0.9, 3.0, 1.2, "Stock Normal",  str((df_inv["Statut"]=="Stock normal").sum()),   GREEN, "Niveau OK")
    kpi_box(s5, 6.7, 0.9, 3.0, 1.2, "Stocks Faibles",str((df_inv["Statut"]=="Stock faible").sum()),   AMBER, "Seuil atteint")
    kpi_box(s5, 9.9, 0.9, 3.1, 1.2, "Ruptures",      str((df_inv["Statut"]=="Non disponible").sum()), RED,  "Stock = 0 !")
    # Production KPIs
    kpi_box(s5, 0.3, 2.3, 3.0, 1.2, "Rendement Moy.", f"{rend_moy:.1f}%",
            GREEN if rend_moy>=80 else RED, "Objectif : 85%")
    kpi_box(s5, 3.5, 2.3, 3.0, 1.2, "Production Totale",
            f"{df_prod['Production réelle'].sum():,.0f}", INDIGO, "Unités sur la période")
    kpi_box(s5, 6.7, 2.3, 3.0, 1.2, "Taux Rebut Moy.", f"{df_prod['Taux rebut'].mean():.2f}%",
            GREEN if df_prod['Taux rebut'].mean()<3 else RED, "Objectif : < 3%")
    kpi_box(s5, 9.9, 2.3, 3.1, 1.2, "Rebuts Totaux",
            f"{df_prod['Rebuts'].sum():,}", AMBER, "Unités perdues")
    # Alertes stock
    rect(s5, 0.3, 3.7, 12.7, 0.38, RED)
    txt(s5, f"🚨 {nb_rupt} RUPTURE(S) DE STOCK — {nb_faib} STOCK(S) FAIBLE(S) — Action requise",
        0.4, 3.75, 12, 0.28, size=9, bold=True, color=WHITE)
    ruptures = df_inv[df_inv["Statut"]=="Non disponible"]["designation"].tolist()[:6]
    faibles  = df_inv[df_inv["Statut"]=="Stock faible"]["designation"].tolist()[:6]
    for i, r in enumerate(ruptures[:3]):
        txt(s5, f"🔴 {r[:40]}", 0.4, 4.2+i*0.45, 6.0, 0.38, size=8.5, color=RGBColor(0x7F,0x1D,0x1D))
    for i, r in enumerate(faibles[:3]):
        txt(s5, f"🟡 {r[:40]}", 6.8, 4.2+i*0.45, 6.0, 0.38, size=8.5, color=RGBColor(0x78,0x35,0x0F))

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 6 — RECOMMANDATIONS & CONCLUSIONS
    # ═══════════════════════════════════════════════════════════════
    s6 = blank_slide(prs); bg(s6, NAVY)
    rect(s6, 0, 0, 13.33, 0.12, INDIGO)
    txt(s6, "💡  RECOMMANDATIONS STRATÉGIQUES — À RETENIR", 0.4, 0.25, 11, 0.55,
        size=18, bold=True, color=WHITE)
    recos = []
    if taux_rec < 80:
        recos.append(("🔴", "TRÉSORERIE", f"Lancer une campagne de relance. {len(df_f[df_f['ETAT DE PAIEMENT']=='Impayée'])} factures impayées pour {fmtp(ca_impaye)}."))
    if nb_rupt > 0:
        recos.append(("🔴", "STOCKS", f"{nb_rupt} produit(s) en rupture totale. Passer commande immédiatement."))
    if nb_faib > 3:
        recos.append(("🟡", "STOCKS", f"{nb_faib} stocks faibles. Planifier les réapprovisionnements sous 2 semaines."))
    if rend_moy < 82:
        recos.append(("🟡", "PRODUCTION", f"Rendement à {rend_moy:.1f}% (< 85%). Révision technique des machines requise."))
    top_seg = df_f.groupby("SEGMENT")["MONTANT TTC"].sum().idxmax() if len(df_f)>0 else "N/A"
    recos.append(("🟢", "COMMERCIAL", f"Renforcer la prospection sur le segment '{top_seg}' — meilleur contributeur CA."))
    recos.append(("🟢", "OBJECTIF", f"CA prévu prochain mois : {fmtp(ca_total/max(len(df_f['DATE'].dt.to_period('M').unique()),1)*1.05)}"))
    for i, (icon, cat, msg) in enumerate(recos[:6]):
        yy = 1.05 + i * 0.95
        col_dot = RED if icon=="🔴" else (AMBER if icon=="🟡" else GREEN)
        dot = s6.shapes.add_shape(1, Inches(0.35), Inches(yy+0.12), Inches(0.22), Inches(0.22))
        dot.fill.solid(); dot.fill.fore_color.rgb = col_dot; dot.line.fill.background()
        txt(s6, cat, 0.7, yy+0.06, 2.0, 0.28, size=8.5, bold=True, color=col_dot)
        txt(s6, msg, 2.85, yy+0.02, 10.0, 0.72, size=9.5, color=RGBColor(0xCB,0xD5,0xE1))
    # Footer
    rect(s6, 0, 7.1, 13.33, 0.4, RGBColor(0x1E,0x29,0x3B))
    txt(s6, f"MULTIPACK SA — Rapport Direction du {now.strftime('%d %B %Y')} · Tableau de Bord v9",
        0.3, 7.15, 10, 0.28, size=8, color=RGBColor(0x47,0x55,0x69))
    txt(s6, "Confidentiel", 11.5, 7.15, 1.5, 0.28, size=8, bold=True,
        color=RGBColor(0x64,0x74,0x8B), align=PP_ALIGN.RIGHT)

    buf = _io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf




# ══════════════════════════════════════════════


# ══════════════════════════════════════════════════════════════


# CHARGEMENT DES DONNÉES
# ══════════════════════════════════════════════════════════════
with st.spinner("Initialisation du tableau de bord MULTIPACK SA…"):
    df_fact, df_inv, df_ent, df_sor, df_prod = generer_donnees()

ALL_SEGMENTS = sorted(df_fact["SEGMENT"].unique().tolist())
ALL_ZONES    = sorted(df_fact["ZONE"].unique().tolist())
ALL_ANNEES   = sorted(df_fact["ANNEE"].unique().tolist())

# ══════════════════════════════════════════════════════════════
# CSS ADDITIONNEL — Page accueil + Header navigation
# ══════════════════════════════════════════════════════════════

# CSS merged into main block above



# ══════════════════════════════════════════════════════════════
# SIDEBAR — PANNEAU DE FILTRES UNIQUEMENT
# ══════════════════════════════════════════════════════════════
with st.sidebar:
    st.markdown(f"""
    <div style="text-align:center;padding:22px 0 14px;">
        <div style="font-size:2.5rem;margin-bottom:6px;">📦</div>
        <div style="font-size:1rem;font-weight:800;color:white;letter-spacing:-0.01em;">MULTIPACK SA</div>
        <div style="font-size:0.65rem;color:#475569;margin-top:3px;">Zone Ind. Yopougon · Abidjan, CI</div>
        <div style="display:inline-flex;align-items:center;gap:5px;background:rgba(46,204,113,0.15);
                    border:1px solid rgba(46,204,113,0.3);border-radius:20px;padding:3px 10px;
                    margin-top:10px;">
            <div style="width:6px;height:6px;border-radius:50%;background:{C['vert']};"></div>
            <span style="font-size:0.62rem;color:{C['vert']};font-weight:700;">Système actif</span>
        </div>
    </div>
    <div style="border-top:1px solid rgba(255,255,255,0.06);margin:0 0 16px;"></div>
    <div style="font-size:0.62rem;font-weight:700;color:#334155;text-transform:uppercase;
                letter-spacing:0.1em;padding:0 4px 8px;">🔧 Panneau de filtres</div>
    """, unsafe_allow_html=True)

    # Navigation (cachée visuellement, utilisée pour la logique)
    page = st.radio("Page", [
        "🏠  Accueil",
        "📊  Vue d'ensemble",
        "💰  Chiffre d'Affaires",
        "🛍️  Nos Produits",
        "👥  Nos Clients",
        "📦  Stocks & Inventaire",
        "🏭  Production",
        "💳  Paiements",
        "🔮  Prévisions",
        "📈  Comparaisons",
        "🧮  Contrôle de Gestion",
        "⚠️  Alertes",
    ], label_visibility="collapsed")

    st.markdown(f"""
    <div style="border-top:1px solid rgba(255,255,255,0.06);margin:10px 0 14px;"></div>
    <div style="font-size:0.62rem;font-weight:700;color:#334155;text-transform:uppercase;
                letter-spacing:0.1em;padding:0 4px 6px;">🗓 Période</div>
    """, unsafe_allow_html=True)
    sel_annees = st.multiselect("Années", ALL_ANNEES, default=ALL_ANNEES, label_visibility="collapsed")

    st.markdown(f"""
    <div style="font-size:0.62rem;font-weight:700;color:#334155;text-transform:uppercase;
                letter-spacing:0.1em;padding:8px 4px 6px;">🏷 Segments</div>
    """, unsafe_allow_html=True)
    sel_seg = st.multiselect("Segments", ALL_SEGMENTS, default=ALL_SEGMENTS,
                              label_visibility="collapsed",
                              help="Segments clients inclus dans toutes les pages")

    st.markdown(f"""
    <div style="font-size:0.62rem;font-weight:700;color:#334155;text-transform:uppercase;
                letter-spacing:0.1em;padding:8px 4px 6px;">🗺 Zones</div>
    """, unsafe_allow_html=True)
    sel_zones = st.multiselect("Zones", ALL_ZONES, default=ALL_ZONES,
                                label_visibility="collapsed")

    # Stats live dans sidebar
    df_sidebar = df_fact[
        df_fact["ANNEE"].isin(sel_annees) &
        df_fact["SEGMENT"].isin(sel_seg) &
        df_fact["ZONE"].isin(sel_zones)
    ]
    ca_side = df_sidebar["MONTANT TTC"].sum()
    nb_side = len(df_sidebar)
    cl_side = df_sidebar["CLIENT"].nunique()

    st.markdown(f"""
    <div style="border-top:1px solid rgba(255,255,255,0.06);margin:14px 0 12px;"></div>
    <div style="font-size:0.62rem;font-weight:700;color:#334155;text-transform:uppercase;
                letter-spacing:0.1em;padding:0 4px 10px;">📊 Sélection active</div>
    <div style="background:rgba(97,114,243,0.1);border:1px solid rgba(97,114,243,0.2);
                border-radius:10px;padding:12px 14px;">
        <div style="display:flex;justify-content:space-between;margin-bottom:7px;">
            <span style="font-size:0.72rem;color:#64748B;">CA filtré</span>
            <span style="font-size:0.78rem;font-weight:700;color:white;">{fmt(ca_side,'').strip()}</span>
        </div>
        <div style="display:flex;justify-content:space-between;margin-bottom:7px;">
            <span style="font-size:0.72rem;color:#64748B;">Factures</span>
            <span style="font-size:0.78rem;font-weight:700;color:white;">{nb_side:,}</span>
        </div>
        <div style="display:flex;justify-content:space-between;">
            <span style="font-size:0.72rem;color:#64748B;">Clients</span>
            <span style="font-size:0.78rem;font-weight:700;color:white;">{cl_side}</span>
        </div>
    </div>
    <div style="margin-top:16px;font-size:0.63rem;color:#334155;text-align:center;line-height:1.7;">
        {len(sel_seg)} segment(s) · {len(sel_zones)} zone(s)<br>
        Maj. {datetime.now().strftime('%d/%m/%Y %H:%M')}<br>
        <span style="color:#1E3A5F;">© MULTIPACK SA 2024</span>
    </div>
    """, unsafe_allow_html=True)

# ── Filtre global appliqué
df_f = df_fact[
    df_fact["ANNEE"].isin(sel_annees) &
    df_fact["SEGMENT"].isin(sel_seg) &
    df_fact["ZONE"].isin(sel_zones)
]

# ══════════════════════════════════════════════════════════════
# HEADER NAVIGATION HORIZONTAL — affiché sur toutes les pages
# ══════════════════════════════════════════════════════════════
now = datetime.now()

# Mapping page → label court pour la nav
NAV_ITEMS = [
    ("🏠  Accueil",          "🏠","Accueil"),
    ("📊  Vue d'ensemble",   "📊","Résumé"),
    ("💰  Chiffre d'Affaires","💰","CA"),
    ("🛍️  Nos Produits",     "🛍️","Produits"),
    ("👥  Nos Clients",      "👥","Clients"),
    ("📦  Stocks & Inventaire","📦","Stocks"),
    ("🏭  Production",       "🏭","Production"),
    ("💳  Paiements",        "💳","Paiements"),
    ("🔮  Prévisions",       "🔮","Prévisions"),
    ("📈  Comparaisons",     "📈","Comparaisons"),
    ("🧮  Contrôle de Gestion","🧮","Contrôle"),
    ("⚠️  Alertes",          "⚠️","Alertes"),
]

# ── HEADER NAVIGATION avec st.columns (pas de raw HTML) ──────────
_page_labels = {k: (ico, lbl) for k, ico, lbl in NAV_ITEMS}
_ico_active, _lbl_active = _page_labels.get(page, ("📊","Page"))

st.markdown(f"""
<div style="background:#0F172A;border-radius:16px;padding:0 24px;
            margin-bottom:20px;display:flex;align-items:center;
            justify-content:space-between;box-shadow:0 4px 24px rgba(0,0,0,0.18);
            border:1px solid rgba(255,255,255,0.05);min-height:72px;gap:16px;">
    <div style="display:flex;align-items:center;gap:12px;min-width:180px;
                border-right:1px solid rgba(255,255,255,0.06);padding-right:20px;padding:16px 20px 16px 0;">
        <span style="font-size:1.8rem;">📦</span>
        <div>
            <div style="font-size:0.92rem;font-weight:800;color:white;letter-spacing:-0.02em;">MULTIPACK SA</div>
            <div style="font-size:0.6rem;color:#475569;font-weight:500;text-transform:uppercase;letter-spacing:0.08em;">Tableau de Bord Direction</div>
        </div>
    </div>
    <div style="flex:1;display:flex;align-items:center;justify-content:center;">
        <div style="background:rgba(79,110,247,0.18);border:1px solid rgba(79,110,247,0.4);
                    border-radius:10px;padding:8px 18px;display:flex;align-items:center;gap:8px;">
            <span style="font-size:1.2rem;">{_ico_active}</span>
            <span style="font-size:0.9rem;font-weight:700;color:white;">{_lbl_active}</span>
            <span style="font-size:0.72rem;color:rgba(255,255,255,0.45);margin-left:4px;">— Page active</span>
        </div>
        <div style="margin-left:16px;font-size:0.72rem;color:#475569;">
            Naviguez via la <b style="color:#94A3B8;">barre latérale gauche</b>
        </div>
    </div>
    <div style="display:flex;align-items:center;gap:10px;border-left:1px solid rgba(255,255,255,0.06);
                padding-left:20px;">
        <div style="width:8px;height:8px;border-radius:50%;background:#10B981;
                    box-shadow:0 0 8px #10B981;animation:pulse 2s infinite;"></div>
        <div style="font-size:0.7rem;color:#475569;line-height:1.5;text-align:right;">
            <b style="color:#94A3B8;">{now.strftime('%d %b %Y')}</b><br>{now.strftime('%H:%M')}
        </div>
    </div>
</div>
""", unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════
# PAGE 0 : ACCUEIL — Page de bienvenue
# ══════════════════════════════════════════════════════════════


if page == "🏠  Accueil":

    # ── HERO BANNER
    ca_total_all = df_fact["MONTANT TTC"].sum()
    nb_clients_all = df_fact["CLIENT"].nunique()
    nb_facts_all = len(df_fact)
    val_stock_all = df_inv["Valeur"].sum()

    st.markdown(f"""
    <div class="welcome-hero">
        <div class="hero-badge">🇨🇮 Zone Industrielle de Yopougon · Abidjan</div>
        <div class="hero-title">Bienvenue sur votre<br>Tableau de Bord</div>
        <div class="hero-sub">
            Plateforme de pilotage commercial et reporting de <b>MULTIPACK SA</b> —
            fabricant de plastiques, sacs supermarché, gobelets et emballages industriels.
            Toutes vos données en un seul endroit, claires et actionnables.
        </div>
        <div class="hero-stats">
            <div class="hero-stat">
                <span class="hs-val">{fmt(ca_total_all,'').strip()}</span>
                <span class="hs-lbl">CA Total (FCFA)</span>
            </div>
            <div class="hero-stat">
                <span class="hs-val">{nb_facts_all:,}</span>
                <span class="hs-lbl">Factures</span>
            </div>
            <div class="hero-stat">
                <span class="hs-val">{nb_clients_all}</span>
                <span class="hs-lbl">Clients actifs</span>
            </div>
            <div class="hero-stat">
                <span class="hs-val">{len(ALL_SEGMENTS)}</span>
                <span class="hs-lbl">Segments marché</span>
            </div>
            <div class="hero-stat">
                <span class="hs-val">{len(ALL_ZONES)-1}</span>
                <span class="hs-lbl">Zones géo.</span>
            </div>
            <div class="hero-stat">
                <span class="hs-val">{fmt(val_stock_all,'').strip()}</span>
                <span class="hs-lbl">Valeur Stock</span>
            </div>
        </div>
    </div>
    """, unsafe_allow_html=True)

    # ── MODULES DU TABLEAU DE BORD
    st.markdown(f'<div class="sec">🧭  LES MODULES DU TABLEAU DE BORD — Cliquez sur un onglet dans la barre latérale gauche</div>',
                unsafe_allow_html=True)

    MODULES = [
        ("📊","Vue d'ensemble","mc-bleu",
         "Résumé exécutif complet en un coup d'œil. KPIs globaux, courbe CA mensuelle, top clients, répartition payé/impayé.",
         ["KPIs","CA mensuel","Top clients","Stocks"]),
        ("💰","Chiffre d'Affaires","mc-vert",
         "Analyse détaillée des ventes : barres mensuelles colorées, comparaison par année, CA par segment client.",
         ["CA mensuel","Par année","Par segment","Objectifs"]),
        ("🛍️","Nos Produits","mc-violet",
         "Performance de chaque famille de produits : Sacs, Gobelets, Emballages. Ventes, marges, évolution mensuelle.",
         ["Par famille","Marges","Évolution","Pie chart"]),
        ("👥","Nos Clients","mc-orange",
         "Portefeuille clients complet : classement, segmentation, taux de recouvrement, zones géographiques.",
         ["Top 10","Segments","Zones","Fidélité"]),
        ("📦","Stocks","mc-rose",
         "État de l'inventaire en temps réel. Alertes ruptures, valeur par catégorie, flux achats/ventes.",
         ["Inventaire","Alertes","Valeur","Flux"]),
        ("🏭","Production","mc-or",
         "Performance de l'usine : rendement par machine, rebuts, comparaison planifié vs réel.",
         ["Rendement","Rebuts","Par machine","Évolution"]),
        ("💳","Paiements","mc-bleu",
         "Analyse de trésorerie : modes de paiement, taux de recouvrement mensuel, clients en retard.",
         ["Recouvrement","Modes","Retards","Tendance"]),
        ("🔮","Prévisions","mc-violet",
         "Projections sur 3 à 12 mois : CA futur, ruptures stocks imminentes, trésorerie anticipée, production.",
         ["CA futur","Stocks","Trésorerie","Production"]),
        ("📈","Comparaisons","mc-vert",
         "Comparez librement n'importe quels segments entre eux : 2 à 2 ou tous ensemble. Radar chart multi-indicateurs.",
         ["2 à 2","Tous segments","Radar","Tendances"]),
        ("⚠️","Alertes","mc-rouge",
         "Tableau de bord des alertes actives : ruptures stocks, impayés critiques, recommandations stratégiques.",
         ["Urgences","Impayés","Conseils","Export"]),
        ("🧮","Contrôle de Gestion","mc-or",
         "Module dédié au contrôleur de gestion : compte de résultat, budget vs réalisé, marges nettes, BFR, DSO, point mort, analyse des écarts.",
         ["P&L","Budget","Marges","BFR","DSO","Point mort"]),
    ]

    rows = [MODULES[i:i+3] for i in range(0, len(MODULES), 3)]
    for row in rows:
        cols = st.columns(len(row))
        for j, (ico, title, cls, desc, tags) in enumerate(row):
            tags_html = "".join(f'<span class="mc-tag">{t}</span>' for t in tags)
            with cols[j]:
                st.markdown(f"""
                <div class="module-card {cls}">
                    <span class="mc-ico">{ico}</span>
                    <div class="mc-title">{title}</div>
                    <div class="mc-desc">{desc}</div>
                    <div class="mc-tags">{tags_html}</div>
                </div>""", unsafe_allow_html=True)
        st.markdown("<div style='margin-bottom:12px;'></div>", unsafe_allow_html=True)

    st.markdown("<br>", unsafe_allow_html=True)

    # ── GUIDE DE DÉMARRAGE
    col_g, col_i = st.columns([1.5, 1])

    with col_g:
        st.markdown(f'<div class="sec">🚀  GUIDE DE DÉMARRAGE RAPIDE</div>', unsafe_allow_html=True)
        steps = [
            ("1", "Utilisez le menu latéral gauche",
             "Cliquez sur le nom d'une page dans la barre de navigation à gauche pour y accéder. "
             "Chaque page est indépendante et se concentre sur un sujet précis."),
            ("2", "Ajustez les filtres selon vos besoins",
             "Dans la barre latérale, vous pouvez filtrer par <b>année</b>, <b>segment de client</b> "
             "et <b>zone géographique</b>. Les graphiques se mettent à jour automatiquement."),
            ("3", "Survolez les graphiques pour le détail",
             "Chaque point, barre ou secteur d'un graphique affiche une <b>info-bulle détaillée</b> "
             "au survol. Cliquez sur la légende pour afficher/masquer des séries."),
            ("4", "Lisez les encadrés bleus",
             "Chaque graphique est accompagné d'un <b>encadré bleu</b> qui explique "
             "en langage simple ce que vous voyez et ce que cela signifie pour MULTIPACK."),
            ("5", "Consultez les Prévisions",
             "La page <b>🔮 Prévisions</b> calcule automatiquement vos tendances futures "
             "sur 3 à 12 mois. Utilisez le slider pour ajuster l'horizon."),
            ("6", "Comparez les segments",
             "La page <b>📈 Comparaisons</b> vous permet de mettre deux segments "
             "face à face ou de tous les analyser ensemble avec un radar chart."),
        ]
        for num, title, desc in steps:
            st.markdown(f"""
            <div class="guide-step">
                <div class="gs-num">{num}</div>
                <div class="gs-content">
                    <div class="gs-title">{title}</div>
                    <div class="gs-desc">{desc}</div>
                </div>
            </div>""", unsafe_allow_html=True)

    with col_i:
        st.markdown(f'<div class="sec">🏭  INFORMATIONS MULTIPACK SA</div>', unsafe_allow_html=True)
        infos = [
            ("Entreprise",    "MULTIPACK SA"),
            ("Secteur",       "Fabrication de plastiques"),
            ("Siège social",  "Zone Ind. Yopougon, Abidjan"),
            ("Pays",          "Côte d'Ivoire 🇨🇮"),
            ("Période données","2022 – 2024"),
            ("Nb. de factures",f"{nb_facts_all:,}"),
            ("Nb. de clients", str(nb_clients_all)),
            ("Segments marché",str(len(ALL_SEGMENTS))),
            ("Zones couvertes", str(len(ALL_ZONES)-1) + " + Export"),
            ("Références stock",str(len(df_inv))),
            ("Machines usine", "5 lignes de production"),
        ]
        st.markdown(f"""
        <div style="background:{C['blanc']};border:1px solid {C['bordure']};border-radius:12px;
                    padding:18px 20px;box-shadow:0 1px 5px rgba(0,0,0,0.05);">
        """, unsafe_allow_html=True)
        for label, val in infos:
            st.markdown(f"""
            <div class="info-row">
                <span class="ir-label">{label}</span>
                <span class="ir-value">{val}</span>
            </div>""", unsafe_allow_html=True)
        st.markdown("</div>", unsafe_allow_html=True)

        st.markdown("<br>", unsafe_allow_html=True)

        # ── Familles de produits
        st.markdown(f'<div class="sec">📦  NOS FAMILLES DE PRODUITS</div>', unsafe_allow_html=True)
        produits_infos = [
            ("🛍️","Sacs Supermarché", "Sacs standard, renforcés, transparents, biosourcés"),
            ("🗑️","Sacs Poubelles",   "De 20L à 240L, vert, noir, bleu"),
            ("🥤","Gobelets",         "De 10cl à 50cl, standard et cristal"),
            ("🥡","Emballages Alim.", "Barquettes PP, film étirable, sachets zip"),
            ("🏭","Emballages Ind.",  "Big bag, fûts HDPE, jerrycans, bidons"),
            ("⚗️","Matières Premières","Granulés PE, PP, PVC, recyclés, masterbatch"),
        ]
        for ico, name, desc in produits_infos:
            st.markdown(f"""
            <div style="display:flex;gap:10px;align-items:flex-start;padding:8px 0;
                        border-bottom:1px solid {C['bordure']};">
                <span style="font-size:1.1rem;min-width:24px;">{ico}</span>
                <div>
                    <div style="font-size:0.82rem;font-weight:700;color:{C['texte']};">{name}</div>
                    <div style="font-size:0.72rem;color:{C['muted']};">{desc}</div>
                </div>
            </div>""", unsafe_allow_html=True)

    st.markdown("<br>", unsafe_allow_html=True)

    # ── ALERTE RAPIDE si problèmes détectés
    nb_rupt = (df_inv["Statut"]=="Non disponible").sum()
    nb_faib = (df_inv["Statut"]=="Stock faible").sum()
    imp_tot = df_fact[df_fact["ETAT DE PAIEMENT"]=="Impayée"]["MONTANT TTC"].sum()
    taux_r  = df_fact[df_fact["ETAT DE PAIEMENT"]=="Payée"]["MONTANT TTC"].sum() / df_fact["MONTANT TTC"].sum() * 100

    has_alert = nb_rupt > 0 or nb_faib > 3 or taux_r < 75
    if has_alert:
        st.markdown(f'<div class="sec">🚨  POINTS D\'ATTENTION DÉTECTÉS AUJOURD\'HUI</div>', unsafe_allow_html=True)
        if nb_rupt > 0:
            st.markdown(f'<div class="alert-r">🔴<div><b>{nb_rupt} produit(s) en rupture de stock</b><br><small>→ Rendez-vous sur la page <b>📦 Stocks</b> ou <b>⚠️ Alertes</b></small></div></div>', unsafe_allow_html=True)
        if nb_faib > 3:
            st.markdown(f'<div class="alert-y">⚠️<div><b>{nb_faib} produits avec stock faible</b><br><small>→ Planifiez des réapprovisionnements · Page <b>🔮 Prévisions → Stocks</b></small></div></div>', unsafe_allow_html=True)
        if taux_r < 75:
            st.markdown(f'<div class="alert-r">💸<div><b>Taux de recouvrement bas : {taux_r:.1f}%</b> (objectif : 80%)<br><small>Impayés : <b>{fmt(imp_tot)}</b> → Page <b>💳 Paiements</b> ou <b>⚠️ Alertes</b></small></div></div>', unsafe_allow_html=True)
    else:
        st.markdown('<div class="alert-g">✅<div><b>Aucune alerte critique détectée</b><br><small>Situation globale satisfaisante. Consultez les pages thématiques pour le détail.</small></div></div>', unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════
# CORRESPONDANCE ANCIENS → NOUVEAUX NOMS DE PAGES
# ══════════════════════════════════════════════════════════════
# Les pages suivantes gardent le même contenu mais avec les nouveaux noms

# PAGE 1 : VUE D'ENSEMBLE
# ══════════════════════════════════════════════════════════════
elif page == "📊  Vue d'ensemble":
    ca_total = df_f["MONTANT TTC"].sum()
    ca_paye  = df_f[df_f["ETAT DE PAIEMENT"]=="Payée"]["MONTANT TTC"].sum()
    ca_impaye= df_f[df_f["ETAT DE PAIEMENT"]=="Impayée"]["MONTANT TTC"].sum()
    taux_rec = ca_paye/ca_total*100 if ca_total else 0
    val_stock= df_inv["Valeur"].sum()
    nb_rupt  = (df_inv["Statut"]=="Non disponible").sum()
    nb_faib  = (df_inv["Statut"]=="Stock faible").sum()
    rend_moy = df_prod["Taux rendement"].mean()

    comment(f"Bonjour. Voici la situation globale de MULTIPACK SA. Sur la période sélectionnée, "
            f"nous avons réalisé <b>{fmt(ca_total)}</b> de chiffre d'affaires, dont <b>{fmt(ca_paye)}</b> encaissés. "
            f"Il reste <b>{fmt(ca_impaye)}</b> à recouvrer. Nos stocks valent <b>{fmt(val_stock)}</b> "
            f"et nos machines tournent à <b>{rend_moy:.0f}%</b> de rendement.", "📋 Résumé exécutif")

    c1,c2,c3,c4 = st.columns(4)
    with c1: kpi("💰","CA Total", fmt(ca_total), f"{len(df_f)} factures","bleu")
    with c2: kpi("✅","Encaissé", fmt(ca_paye), f"Taux : {taux_rec:.1f}%","vert")
    with c3: kpi("⏳","Impayés", fmt(ca_impaye), f"{len(df_f[df_f['ETAT DE PAIEMENT']=='Impayée'])} factures","rouge")
    with c4: kpi("👥","Clients Actifs", str(df_f["CLIENT"].nunique()), f"{df_f['SEGMENT'].nunique()} segments","violet")
    st.markdown("<br>", unsafe_allow_html=True)
    c5,c6,c7,c8 = st.columns(4)
    with c5: kpi("🏭","Rendement Usine", f"{rend_moy:.1f}%","Moyenne machines","bleu")
    with c6: kpi("📦","Valeur Stock", fmt(val_stock), f"{len(df_inv)} références","violet")
    with c7: kpi("🟡","Stocks Faibles", str(nb_faib),"À commander bientôt","orange")
    with c8: kpi("🔴","Ruptures", str(nb_rupt),"Produits épuisés — urgent","rouge" if nb_rupt>0 else "vert")
    st.markdown("<br>", unsafe_allow_html=True)

    # Courbe CA mensuelle seule
    section("📈  ÉVOLUTION DU CHIFFRE D'AFFAIRES — Mois par mois")
    comment("Chaque point de cette courbe = les ventes d'un mois. Une courbe qui monte = bonne dynamique. "
            "La ligne pointillée = notre objectif mensuel moyen à dépasser.")
    ca_m = df_f.copy(); ca_m["Période"] = ca_m["DATE"].dt.to_period("M").astype(str)
    ca_mois = ca_m.groupby("Période")["MONTANT TTC"].sum().reset_index(); ca_mois.columns = ["Période","CA"]
    objectif = ca_mois["CA"].mean() * 1.10
    fig = go.Figure()
    fig.add_trace(go.Scatter(x=ca_mois["Période"], y=ca_mois["CA"], mode="lines+markers",
        name="CA mensuel", line=dict(color=C["bleu"], width=3),
        marker=dict(size=8, color=C["bleu"], line=dict(color="white",width=2)),
        fill="tozeroy", fillcolor="rgba(97,114,243,0.10)",
        hovertemplate="<b>%{x}</b><br>CA : %{y:,.0f} FCFA<extra></extra>"))
    fig.add_hline(y=objectif, line_dash="dot", line_color=C["rose"], line_width=2,
        annotation_text=f"Objectif {fmt(objectif,'').strip()}", annotation_font_color=C["rose"],
        annotation_position="top right")
    fig = excel_style(fig, 320, False)
    fig.update_xaxes(tickangle=-35, tickfont=dict(size=9))
    fig.update_yaxes(tickformat=",.0f", ticksuffix=" F")
    st.plotly_chart(fig, use_container_width=True)
    st.markdown("<br>", unsafe_allow_html=True)

    col_a, col_b = st.columns(2)
    with col_a:
        section("🥧  PART DES PAIEMENTS REÇUS")
        comment(f"Sur <b>{fmt(ca_total)}</b> de ventes, nous avons encaissé <b>{taux_rec:.0f}%</b>. "
                f"Un bon taux se situe au-dessus de 80%.")
        stat_d = df_f.groupby("ETAT DE PAIEMENT")["MONTANT TTC"].sum().reset_index()
        clrs = {r: (C["bleu"] if r=="Payée" else C["rouge"]) for r in stat_d["ETAT DE PAIEMENT"]}
        fig2 = go.Figure(go.Pie(
            labels=stat_d["ETAT DE PAIEMENT"], values=stat_d["MONTANT TTC"], hole=0.65,
            marker=dict(colors=[clrs[r] for r in stat_d["ETAT DE PAIEMENT"]], line=dict(color="white",width=3)),
            textinfo="label+percent", textfont=dict(size=12),
            hovertemplate="%{label}<br><b>%{value:,.0f} FCFA</b> (%{percent})<extra></extra>"))
        fig2.update_layout(paper_bgcolor="#FFFFFF", height=300,
            margin=dict(t=10,b=30,l=10,r=10),
            legend=dict(orientation="h", y=-0.08, xanchor="center", x=0.5, font=dict(color=C["texte"],size=11)),
            annotations=[dict(text=f"<b>{taux_rec:.0f}%</b>", x=0.5, y=0.5, showarrow=False,
                              font=dict(size=22, color=C["bleu"]), xanchor="center")])
        st.plotly_chart(fig2, use_container_width=True)

    with col_b:
        section("🏆  TOP 5 CLIENTS")
        comment("Nos 5 meilleurs clients par chiffre d'affaires. Ce sont eux qu'il faut fidéliser en priorité.")
        top5 = df_f.groupby("CLIENT")["MONTANT TTC"].sum().sort_values(ascending=False).head(5).reset_index()
        fig3 = go.Figure(go.Bar(
            x=top5["CLIENT"], y=top5["MONTANT TTC"],
            marker_color=[C["bleu"],C["bleu_clair"],C["violet"],C["rose"],C["bleu_fonce"]],
            text=top5["MONTANT TTC"].apply(lambda v: fmt(v,"").strip()),
            textposition="outside", textfont=dict(size=11),
            hovertemplate="<b>%{x}</b><br>CA : %{y:,.0f} FCFA<extra></extra>"))
        fig3 = excel_style(fig3, 300, False)
        fig3.update_xaxes(tickangle=-20, tickfont=dict(size=10))
        fig3.update_yaxes(tickformat=",.0f", ticksuffix=" F")
        st.plotly_chart(fig3, use_container_width=True)


# ══════════════════════════════════════════════════════════════
# PAGE 2 : CHIFFRE D'AFFAIRES
# ══════════════════════════════════════════════════════════════
elif page == "💰  Chiffre d'Affaires":
    ca = df_f["MONTANT TTC"].sum()
    paye = df_f[df_f["ETAT DE PAIEMENT"]=="Payée"]["MONTANT TTC"].sum()
    taux = paye/ca*100 if ca else 0
    ca21 = df_fact[df_fact["ANNEE"]==2022]["MONTANT TTC"].sum()
    ca22 = df_fact[df_fact["ANNEE"]==2023]["MONTANT TTC"].sum()
    crois = (ca22-ca21)/ca21*100 if ca21 else 0

    comment(f"CA total : <b>{fmt(ca)}</b>. Croissance 2022→2023 : <b>{crois:+.1f}%</b>. "
            f"Analysez les mois forts et faibles pour mieux planifier.", "📊 Guide de lecture")

    c1,c2,c3,c4,c5 = st.columns(5)
    with c1: kpi("💰","CA Total", fmt(ca), f"{len(df_f)} factures","bleu")
    with c2: kpi("✅","Encaissé", fmt(paye), f"{taux:.1f}%","vert")
    with c3: kpi("🛒","Commande Moy.", fmt(df_f["MONTANT TTC"].mean()),"Par facture","violet")
    with c4: kpi("👥","Nb Clients", str(df_f["CLIENT"].nunique()),"Actifs","orange")
    with c5: kpi("📈","Croissance 22→23", f"{crois:+.1f}%","Variation annuelle","vert" if crois>0 else "rouge")
    st.markdown("<br>", unsafe_allow_html=True)

    # BARRES mensuelles colorées
    section("📊  CA PAR MOIS — Chaque barre = un mois de ventes")
    comment("🟢 Vert = mois fort &nbsp; 🟡 Orange = mois moyen &nbsp; 🔴 Rouge = mois faible. "
            "Survolez une barre pour voir le montant exact.")
    ca_mensuel = df_f.copy(); ca_mensuel["Période"] = ca_mensuel["DATE"].dt.to_period("M").astype(str)
    ca_mb = ca_mensuel.groupby("Période")["MONTANT TTC"].sum().reset_index()
    q33 = ca_mb["MONTANT TTC"].quantile(0.33); q66 = ca_mb["MONTANT TTC"].quantile(0.66)
    ca_mb["Couleur"] = ca_mb["MONTANT TTC"].apply(lambda v: C["vert"] if v>=q66 else (C["orange"] if v>=q33 else C["rouge"]))
    fig = go.Figure(go.Bar(x=ca_mb["Période"], y=ca_mb["MONTANT TTC"],
        marker_color=ca_mb["Couleur"].tolist(),
        text=ca_mb["MONTANT TTC"].apply(lambda v: fmt(v,"").strip()),
        textposition="outside", textfont=dict(size=9),
        hovertemplate="<b>%{x}</b><br>CA : %{y:,.0f} FCFA<extra></extra>"))
    for lbl, clr in [("🟢 Mois fort",C["vert"]),("🟡 Mois moyen",C["orange"]),("🔴 Mois faible",C["rouge"])]:
        fig.add_trace(go.Scatter(x=[None],y=[None],mode="markers",
            marker=dict(color=clr,size=10,symbol="square"),name=lbl))
    fig = excel_style(fig, 360)
    fig.update_xaxes(tickangle=-40, tickfont=dict(size=8))
    fig.update_yaxes(tickformat=",.0f", ticksuffix=" F")
    st.plotly_chart(fig, use_container_width=True)
    st.markdown("<br>", unsafe_allow_html=True)

    # COURBES séparées par année
    section("📈  COMPARAISON ANNÉE PAR ANNÉE — Une courbe = une année")
    comment("Si la courbe 2024 est au-dessus de 2023 → nous progressons. Chaque couleur = une année.")
    ca_ann = df_f.groupby(["ANNEE","MOIS"])["MONTANT TTC"].sum().reset_index()
    mois_l = {1:"Jan",2:"Fév",3:"Mar",4:"Avr",5:"Mai",6:"Jun",7:"Jul",8:"Aoû",9:"Sep",10:"Oct",11:"Nov",12:"Déc"}
    ca_ann["Mois_L"] = ca_ann["MOIS"].map(mois_l)
    ann_colors = {2022:C["violet"],2023:C["rose"],2024:C["bleu"]}
    fig2 = go.Figure()
    for an in sorted(ca_ann["ANNEE"].unique()):
        d = ca_ann[ca_ann["ANNEE"]==an].sort_values("MOIS")
        col_an = ann_colors.get(an, C["bleu_fonce"])
        fig2.add_trace(go.Scatter(x=d["Mois_L"], y=d["MONTANT TTC"], mode="lines+markers", name=str(an),
            line=dict(color=col_an, width=3), marker=dict(size=9, color=col_an, line=dict(color="white",width=2)),
            hovertemplate=f"<b>%{{x}} {an}</b><br>CA : %{{y:,.0f}} FCFA<extra></extra>"))
    fig2 = excel_style(fig2, 320)
    fig2.update_yaxes(tickformat=",.0f", ticksuffix=" F")
    st.plotly_chart(fig2, use_container_width=True)
    st.markdown("<br>", unsafe_allow_html=True)

    # BARRES horizontales CA par segment
    section("🏷️  CA PAR TYPE DE CLIENT — Qui nous rapporte le plus ?")
    comment("La Grande Distribution et les Industriels sont généralement nos premiers acheteurs. "
            "Les distributeurs locaux et l'Export sont des relais de croissance.")
    seg_ca = df_f.groupby("SEGMENT")["MONTANT TTC"].sum().sort_values(ascending=True).reset_index()
    fig3 = go.Figure(go.Bar(x=seg_ca["MONTANT TTC"], y=seg_ca["SEGMENT"], orientation="h",
        marker_color=PIE[:len(seg_ca)],
        text=seg_ca["MONTANT TTC"].apply(lambda v: fmt(v,"").strip()),
        textposition="outside", textfont=dict(size=11),
        hovertemplate="<b>%{y}</b><br>CA : %{x:,.0f} FCFA<extra></extra>"))
    fig3 = excel_style(fig3, 300, False)
    fig3.update_xaxes(tickformat=",.0f")
    st.plotly_chart(fig3, use_container_width=True)


# ══════════════════════════════════════════════════════════════
# PAGE 3 : NOS PRODUITS
# ══════════════════════════════════════════════════════════════
elif page == "🛍️  Nos Produits":
    sor_agg = df_sor.groupby("Catégorie").agg(CA=("Total","sum"),Qte=("Quantité","sum")).reset_index()
    ent_agg = df_ent.groupby("Catégorie").agg(Achats=("Total","sum")).reset_index()
    perf = sor_agg.merge(ent_agg, on="Catégorie", how="left").fillna(0)
    perf["Marge"] = perf["CA"] - perf["Achats"]
    perf["Taux_m"] = np.where(perf["CA"]>0, perf["Marge"]/perf["CA"]*100, 0).round(1)
    perf = perf.sort_values("CA", ascending=False)
    tm = perf["Marge"].sum()/sor_agg["CA"].sum()*100 if sor_agg["CA"].sum()>0 else 0

    comment("Cette page montre la performance de chaque famille de produits : "
            "Sacs supermarché, Sacs poubelles, Gobelets, Emballages. "
            "La marge = ce qu'on gagne réellement après les matières premières.", "🛍️ Guide de lecture")

    c1,c2,c3,c4 = st.columns(4)
    with c1: kpi("📦","Total Ventes", fmt(sor_agg["CA"].sum()),"Toutes familles","bleu")
    with c2: kpi("🛒","Coût Matières", fmt(ent_agg["Achats"].sum()),"Achats production","orange")
    with c3: kpi("💚","Marge Brute", fmt(perf["Marge"].sum()),"Ventes - Achats","vert")
    with c4: kpi("💹","Taux de Marge", f"{tm:.1f}%","Rentabilité globale","vert" if tm>25 else "orange")
    st.markdown("<br>", unsafe_allow_html=True)

    # Barres CA par famille
    section("📊  VENTES PAR FAMILLE DE PRODUIT")
    comment("Chaque barre = le total des ventes d'une famille. "
            "Les Sacs et Gobelets sont nos produits phares. Un bon CA avec une mauvaise marge = problème à corriger.")
    fig = go.Figure(go.Bar(x=perf["Catégorie"], y=perf["CA"], marker_color=PIE[:len(perf)],
        text=perf["CA"].apply(lambda v: fmt(v,"").strip()),
        textposition="outside", textfont=dict(size=10),
        hovertemplate="<b>%{x}</b><br>Ventes : %{y:,.0f} FCFA<extra></extra>"))
    fig = excel_style(fig, 320, False)
    fig.update_xaxes(tickangle=-15, tickfont=dict(size=10))
    fig.update_yaxes(tickformat=",.0f", ticksuffix=" F")
    st.plotly_chart(fig, use_container_width=True)
    st.markdown("<br>", unsafe_allow_html=True)

    col_a, col_b = st.columns(2)
    with col_a:
        section("💹  MARGE BRUTE PAR FAMILLE")
        comment("La marge = ce que chaque famille nous rapporte réellement. "
                "🟢 Positif = rentable. 🔴 Négatif = on perd de l'argent sur ce produit.")
        perf_s = perf.sort_values("Marge", ascending=True)
        fig2 = go.Figure(go.Bar(x=perf_s["Marge"], y=perf_s["Catégorie"], orientation="h",
            marker_color=[C["vert"] if v>0 else C["rouge"] for v in perf_s["Marge"]],
            text=perf_s["Marge"].apply(lambda v: fmt(v,"").strip()),
            textposition="outside", textfont=dict(size=10),
            hovertemplate="<b>%{y}</b><br>Marge : %{x:,.0f} FCFA<extra></extra>"))
        fig2 = excel_style(fig2, 300, False)
        fig2.update_xaxes(tickformat=",.0f")
        st.plotly_chart(fig2, use_container_width=True)

    with col_b:
        section("🥧  PART DE CHAQUE FAMILLE DANS LES VENTES")
        comment("Ce graphique rond montre la part (%) de chaque famille dans notre CA total. "
                "Un seul produit dominant = risque si la demande chute. Mieux vaut diversifier.")
        fig3 = go.Figure(go.Pie(labels=perf["Catégorie"], values=perf["CA"], hole=0,
            marker=dict(colors=PIE[:len(perf)], line=dict(color="white",width=2)),
            textinfo="label+percent", textfont=dict(size=11),
            hovertemplate="<b>%{label}</b><br>%{value:,.0f} FCFA (%{percent})<extra></extra>"))
        fig3.update_layout(paper_bgcolor="#FFFFFF", height=300, margin=dict(t=10,b=30,l=10,r=10),
            legend=dict(orientation="v", x=1.02, y=0.5, font=dict(color=C["texte"],size=10)))
        st.plotly_chart(fig3, use_container_width=True)

    st.markdown("<br>", unsafe_allow_html=True)

    # Courbes séparées par famille
    section("📈  ÉVOLUTION MENSUELLE PAR FAMILLE — Une courbe par gamme de produit")
    comment("Chaque courbe = une famille de produits. Survolez un point pour voir le montant exact. "
            "Une courbe qui monte régulièrement = la demande augmente. Une chute = risque à surveiller.")
    df_sor_m = df_sor.copy()
    df_sor_m["Mois"] = pd.to_datetime(df_sor_m["Date"]).dt.to_period("M").astype(str)
    sm = df_sor_m.groupby(["Mois","Catégorie"])["Total"].sum().reset_index()
    fig4 = go.Figure()
    for j, cat in enumerate(sm["Catégorie"].unique()):
        d = sm[sm["Catégorie"]==cat]; c = PIE[j % len(PIE)]
        fig4.add_trace(go.Scatter(x=d["Mois"], y=d["Total"], name=cat, mode="lines+markers",
            line=dict(color=c, width=2.5), marker=dict(size=6, color=c, line=dict(color="white",width=1.5)),
            hovertemplate=f"<b>{cat}</b><br>%{{x}}<br>%{{y:,.0f}} FCFA<extra></extra>"))
    fig4 = excel_style(fig4, 340)
    fig4.update_xaxes(tickangle=-35, tickfont=dict(size=8))
    fig4.update_yaxes(tickformat=",.0f", ticksuffix=" F")
    st.plotly_chart(fig4, use_container_width=True)


# ══════════════════════════════════════════════════════════════
# PAGE 4 : NOS CLIENTS
# ══════════════════════════════════════════════════════════════
elif page == "👥  Nos Clients":
    cl_stats = df_f.groupby(["CLIENT","SEGMENT","ZONE"]).agg(
        CA=("MONTANT TTC","sum"), Nb_Cmd=("MONTANT TTC","count"), Panier=("MONTANT TTC","mean"),
        CA_P=("MONTANT TTC", lambda x: x[df_f.loc[x.index,"ETAT DE PAIEMENT"]=="Payée"].sum()),
        CA_I=("MONTANT TTC", lambda x: x[df_f.loc[x.index,"ETAT DE PAIEMENT"]=="Impayée"].sum()),
    ).reset_index()
    cl_stats["Taux_R"] = np.where(cl_stats["CA"]>0, cl_stats["CA_P"]/cl_stats["CA"]*100, 0).round(1)
    cl_stats = cl_stats.sort_values("CA", ascending=False)

    comment(f"Nous avons <b>{len(cl_stats)}</b> clients actifs. Un client fidèle coûte 5x moins cher "
            f"à garder qu'à en trouver un nouveau. Concentrons nos efforts sur les meilleurs.", "👥 Pourquoi analyser nos clients ?")

    c1,c2,c3,c4 = st.columns(4)
    with c1: kpi("👥","Total Clients", str(len(cl_stats)),"Au moins 1 commande","bleu")
    with c2: kpi("⭐","Clients Fidèles", str((cl_stats["Nb_Cmd"]>=5).sum()),"5 commandes ou +","vert")
    with c3: kpi("⚠️","Clients Impayés", str((cl_stats["CA_I"]>0).sum()),"À relancer","rouge")
    with c4: kpi("🏆","CA Moyen/Client", fmt(cl_stats["CA"].mean()),"Par client actif","violet")
    st.markdown("<br>", unsafe_allow_html=True)

    # Top 10 clients
    section("🏆  TOP 10 CLIENTS — Classement par chiffre d'affaires")
    comment("Nos 10 meilleurs clients. Une perte d'un de ces clients aurait un impact fort sur notre CA. "
            "Visitez-les régulièrement et soignez la relation.")
    top10 = cl_stats.head(10).sort_values("CA", ascending=True)
    fig = go.Figure(go.Bar(x=top10["CA"], y=top10["CLIENT"], orientation="h", marker_color=C["bleu"],
        text=top10["CA"].apply(lambda v: fmt(v,"").strip()),
        textposition="outside", textfont=dict(size=10),
        hovertemplate="<b>%{y}</b><br>CA : %{x:,.0f} FCFA<extra></extra>"))
    fig = excel_style(fig, 360, False)
    fig.update_xaxes(tickformat=",.0f")
    st.plotly_chart(fig, use_container_width=True)
    st.markdown("<br>", unsafe_allow_html=True)

    col_a, col_b = st.columns(2)
    with col_a:
        section("🗂️  CA PAR TYPE DE CLIENT")
        comment("Répartition des ventes par catégorie de clients. "
                "Trop dépendre d'un seul type = risque. Mieux vaut diversifier.")
        seg_ca = cl_stats.groupby("SEGMENT")["CA"].sum().reset_index()
        fig2 = go.Figure(go.Pie(labels=seg_ca["SEGMENT"], values=seg_ca["CA"], hole=0.60,
            marker=dict(colors=PIE[:len(seg_ca)], line=dict(color="white",width=2)),
            textinfo="label+percent", textfont=dict(size=11),
            hovertemplate="<b>%{label}</b><br>%{value:,.0f} FCFA (%{percent})<extra></extra>"))
        fig2.update_layout(paper_bgcolor="#FFFFFF", height=320, margin=dict(t=10,b=30,l=10,r=10),
            legend=dict(orientation="h", y=-0.08, xanchor="center", x=0.5, font=dict(color=C["texte"],size=10)))
        st.plotly_chart(fig2, use_container_width=True)

    with col_b:
        section("🗺️  CA PAYÉ vs IMPAYÉ PAR ZONE")
        comment("Bleu = encaissé. Rose = factures non payées. "
                "Une zone avec beaucoup de rose = relancer les clients de cette région en priorité.")
        zone_ca = cl_stats.groupby("ZONE")[["CA_P","CA_I"]].sum().reset_index()
        fig3 = go.Figure()
        fig3.add_trace(go.Bar(x=zone_ca["ZONE"], y=zone_ca["CA_P"], name="Payé ✅", marker_color=C["bleu"]))
        fig3.add_trace(go.Bar(x=zone_ca["ZONE"], y=zone_ca["CA_I"], name="Impayé ⏳", marker_color=C["rose"]))
        fig3.update_layout(barmode="stack")
        fig3 = excel_style(fig3, 320)
        fig3.update_xaxes(tickangle=-15, tickfont=dict(size=10))
        fig3.update_yaxes(tickformat=",.0f", ticksuffix=" F")
        st.plotly_chart(fig3, use_container_width=True)

    section("📋  FICHE COMPLÈTE NOS CLIENTS")
    disp = cl_stats[["CLIENT","SEGMENT","ZONE","CA","Nb_Cmd","Panier","CA_P","CA_I","Taux_R"]].copy()
    for col_n in ["CA","CA_P","CA_I","Panier"]:
        disp[col_n] = disp[col_n].apply(lambda x: f"{x:,.0f} FCFA")
    disp["Taux_R"] = cl_stats["Taux_R"].apply(lambda x: f"{x:.1f}%")
    disp = disp.rename(columns={"CLIENT":"Client","SEGMENT":"Segment","ZONE":"Zone",
        "CA":"CA Total","Nb_Cmd":"Commandes","Panier":"Panier Moy.","CA_P":"Payé","CA_I":"Impayé","Taux_R":"% Payé"})
    st.dataframe(disp, use_container_width=True, height=350)

# ══════════════════════════════════════════════════════════════
# PAGE 5 : STOCKS & INVENTAIRE
# ══════════════════════════════════════════════════════════════
elif page == "📦  Stocks & Inventaire":
    val_stock = df_inv["Valeur"].sum()
    nb_rupt   = (df_inv["Statut"]=="Non disponible").sum()
    nb_faib   = (df_inv["Statut"]=="Stock faible").sum()
    nb_norm   = (df_inv["Statut"]=="Stock normal").sum()

    comment(f"Notre stock total vaut <b>{fmt(val_stock)}</b>. "
            f"<b>{nb_rupt}</b> produit(s) en rupture totale et <b>{nb_faib}</b> approchent de la limite. "
            f"Un bon stock = on ne rate aucune commande client.", "📦 Pourquoi surveiller les stocks ?")

    c1,c2,c3,c4,c5 = st.columns(5)
    with c1: kpi("💰","Valeur du Stock", fmt(val_stock), f"{len(df_inv)} références","bleu")
    with c2: kpi("✅","Stock Normal", str(nb_norm),"Niveau satisfaisant","vert")
    with c3: kpi("🟡","Stock Faible", str(nb_faib),"Commander bientôt","orange")
    with c4: kpi("🔴","Rupture Totale", str(nb_rupt),"Commander d'urgence","rouge")
    with c5:
        rot = df_sor["Quantité"].sum()/max(df_inv["Stock final"].sum(),1)
        kpi("🔄","Rotation Stock", f"{rot:.1f}×","Fréquence renouvellement","violet")
    st.markdown("<br>", unsafe_allow_html=True)

    # Barres stock par catégorie colorées
    section("📊  STOCK DISPONIBLE PAR FAMILLE")
    comment("🟢 Vert = stock suffisant &nbsp; 🟡 Orange = à surveiller &nbsp; 🔴 Rouge = rupture imminente. "
            "La ligne pointillée = le seuil minimum acceptable.")
    cat_s = df_inv.groupby("categorie").agg(Stock=("Stock final","sum"), Seuil=("seuil","sum")).reset_index()
    clr_s = [C["vert"] if row["Stock"]>=row["Seuil"] else (C["orange"] if row["Stock"]>0 else C["rouge"])
              for _,row in cat_s.iterrows()]
    fig = go.Figure(go.Bar(x=cat_s["categorie"], y=cat_s["Stock"], marker_color=clr_s,
        text=cat_s["Stock"].apply(lambda v: f"{v:,} u."),
        textposition="outside", textfont=dict(size=11),
        hovertemplate="<b>%{x}</b><br>Stock : %{y:,} unités<extra></extra>"))
    fig = excel_style(fig, 320, False)
    fig.update_xaxes(tickangle=-15, tickfont=dict(size=10))
    fig.update_yaxes(tickformat=",")
    st.plotly_chart(fig, use_container_width=True)
    st.markdown("<br>", unsafe_allow_html=True)

    col_a, col_b = st.columns(2)
    with col_a:
        section("💰  VALEUR DU STOCK PAR FAMILLE")
        comment("Combien vaut en FCFA chaque famille en stock. "
                "Trop de stock coûte de l'argent (entrepôt, risque) — pas assez = ventes ratées.")
        cat_v = df_inv.groupby("categorie")["Valeur"].sum().reset_index()
        fig2 = go.Figure(go.Pie(labels=cat_v["categorie"], values=cat_v["Valeur"], hole=0.55,
            marker=dict(colors=PIE[:len(cat_v)], line=dict(color="white",width=2)),
            textinfo="label+percent", textfont=dict(size=11),
            hovertemplate="<b>%{label}</b><br>%{value:,.0f} FCFA (%{percent})<extra></extra>"))
        fig2.update_layout(paper_bgcolor="#FFFFFF", height=300, margin=dict(t=10,b=30,l=10,r=10),
            legend=dict(orientation="v", x=1.02, y=0.5, font=dict(color=C["texte"],size=10)))
        st.plotly_chart(fig2, use_container_width=True)

    with col_b:
        section("🚦  ÉTAT DE SANTÉ DU STOCK")
        comment("🟢 Normal : tout va bien. 🟡 Faible : commander prochainement. "
                "🔴 Non disponible : rupture totale, risque de perdre des ventes.")
        statuts = ["Stock normal","Stock faible","Non disponible"]
        nbs = [(df_inv["Statut"]==s).sum() for s in statuts]
        fig3 = go.Figure(go.Bar(x=nbs, y=statuts, orientation="h",
            marker_color=[C["vert"],C["orange"],C["rouge"]],
            text=[f"{n} référence(s)" for n in nbs],
            textposition="outside", textfont=dict(size=12),
            hovertemplate="<b>%{y}</b> : %{x} référence(s)<extra></extra>"))
        fig3 = excel_style(fig3, 300, False)
        fig3.update_xaxes(range=[0, max(nbs)*1.4])
        st.plotly_chart(fig3, use_container_width=True)

    st.markdown("<br>", unsafe_allow_html=True)

    # Courbe achats mensuels
    section("📥  ACHATS DE MATIÈRES — Évolution mensuelle")
    comment("Combien nous avons dépensé chaque mois pour acheter nos matières premières. "
            "Des dépenses élevées sans ventes correspondantes = attention à la trésorerie.")
    df_em = df_ent.copy(); df_em["Mois"] = pd.to_datetime(df_em["Date"]).dt.to_period("M").astype(str)
    em_mois = df_em.groupby("Mois")["Total"].sum().reset_index()
    fig4 = go.Figure(go.Scatter(x=em_mois["Mois"], y=em_mois["Total"], mode="lines+markers",
        line=dict(color=C["violet"], width=3),
        marker=dict(size=8, color=C["violet"], line=dict(color="white",width=2)),
        fill="tozeroy", fillcolor="rgba(81,69,158,0.10)",
        hovertemplate="<b>%{x}</b><br>Achats : %{y:,.0f} FCFA<extra></extra>"))
    fig4 = excel_style(fig4, 280, False)
    fig4.update_xaxes(tickangle=-35, tickfont=dict(size=8))
    fig4.update_yaxes(tickformat=",.0f", ticksuffix=" F")
    st.plotly_chart(fig4, use_container_width=True)

    section("📋  INVENTAIRE COMPLET")
    inv_d = df_inv[["ref","designation","categorie","seuil","Stock initial","Entrées","Sorties","Stock final","Valeur","Statut"]].copy()
    inv_d["Valeur"] = inv_d["Valeur"].apply(lambda x: f"{x:,.0f} FCFA")
    inv_d = inv_d.rename(columns={"ref":"Réf.","designation":"Désignation","categorie":"Catégorie","seuil":"Seuil"})
    st.dataframe(inv_d, use_container_width=True, height=340)


# ══════════════════════════════════════════════════════════════
# PAGE 6 : PRODUCTION
# ══════════════════════════════════════════════════════════════
elif page == "🏭  Production":
    rend_moy = df_prod["Taux rendement"].mean()
    rebut_moy= df_prod["Taux rebut"].mean()
    prod_tot = df_prod["Production réelle"].sum()
    reb_tot  = df_prod["Rebuts"].sum()

    comment(f"Nos machines ont produit <b>{prod_tot:,.0f} unités</b> avec un rendement de <b>{rend_moy:.1f}%</b>. "
            f"Taux de rebut (produits ratés) : <b>{rebut_moy:.2f}%</b>. Objectif : atteindre 85%.", "🏭 Comment lire la production ?")

    c1,c2,c3,c4 = st.columns(4)
    with c1: kpi("⚙️","Production Totale", f"{prod_tot:,.0f}","unités fabriquées","bleu")
    with c2: kpi("📈","Rendement Moyen", f"{rend_moy:.1f}%","Objectif : 85%","vert" if rend_moy>=85 else "orange")
    with c3: kpi("🗑️","Taux Rebut Moyen", f"{rebut_moy:.2f}%","Objectif < 3%","rouge" if rebut_moy>=3 else "vert")
    with c4: kpi("🔧","Unités Perdues", f"{reb_tot:,.0f}","Rebuts cumulés — coût caché","orange")
    st.markdown("<br>", unsafe_allow_html=True)

    # Courbe production mensuelle globale
    section("📈  PRODUCTION MENSUELLE — Évolution du volume produit")
    comment("Combien d'unités nos machines ont fabriqué chaque mois. "
            "Les pics = forte demande. Les creux = panne, maintenance ou manque de commandes.")
    pm = df_prod.groupby("Mois_Label").agg(Réelle=("Production réelle","sum")).reset_index()
    try:
        pm["sk"] = pd.to_datetime(pm["Mois_Label"], format="%b %Y"); pm = pm.sort_values("sk")
    except: pass
    fig = go.Figure(go.Scatter(x=pm["Mois_Label"], y=pm["Réelle"], mode="lines+markers",
        line=dict(color=C["bleu"], width=3), marker=dict(size=8, color=C["bleu"], line=dict(color="white",width=2)),
        fill="tozeroy", fillcolor="rgba(97,114,243,0.10)",
        hovertemplate="<b>%{x}</b><br>Production : %{y:,.0f} unités<extra></extra>"))
    fig = excel_style(fig, 300, False)
    fig.update_xaxes(tickangle=-35, tickfont=dict(size=8))
    fig.update_yaxes(tickformat=",")
    st.plotly_chart(fig, use_container_width=True)
    st.markdown("<br>", unsafe_allow_html=True)

    col_a, col_b = st.columns(2)
    with col_a:
        section("⚙️  RENDEMENT PAR MACHINE")
        comment("🟢 Au-dessus de 85% = excellent. 🟡 75-85% = acceptable. 🔴 Sous 75% = maintenance nécessaire. "
                "La ligne rouge = notre objectif.")
        rm = df_prod.groupby("Machine")["Taux rendement"].mean().reset_index().sort_values("Taux rendement")
        clr_m = [C["vert"] if v>=85 else C["orange"] if v>=75 else C["rouge"] for v in rm["Taux rendement"]]
        fig2 = go.Figure(go.Bar(x=rm["Taux rendement"], y=rm["Machine"], orientation="h",
            marker_color=clr_m, text=rm["Taux rendement"].apply(lambda v: f"{v:.1f}%"),
            textposition="outside", textfont=dict(size=12),
            hovertemplate="<b>%{y}</b><br>Rendement : %{x:.1f}%<extra></extra>"))
        fig2.add_vline(x=85, line_dash="dot", line_color=C["rouge"], line_width=2,
            annotation_text="Objectif 85%", annotation_font=dict(color=C["rouge"],size=11))
        fig2 = excel_style(fig2, 300, False)
        fig2.update_xaxes(ticksuffix="%", range=[60,108])
        st.plotly_chart(fig2, use_container_width=True)

    with col_b:
        section("🗑️  REBUTS PAR MACHINE")
        comment("Les rebuts = produits fabriqués mais inutilisables. "
                "Moins il y en a, mieux c'est. Une machine avec beaucoup de rebuts doit être révisée.")
        rb = df_prod.groupby("Machine")["Rebuts"].sum().reset_index().sort_values("Rebuts", ascending=True)
        fig3 = go.Figure(go.Bar(x=rb["Rebuts"], y=rb["Machine"], orientation="h",
            marker_color=C["rose"], text=rb["Rebuts"].apply(lambda v: f"{v:,.0f}"),
            textposition="outside", textfont=dict(size=11),
            hovertemplate="<b>%{y}</b><br>Rebuts : %{x:,.0f} unités<extra></extra>"))
        fig3 = excel_style(fig3, 300, False)
        fig3.update_xaxes(tickformat=",")
        st.plotly_chart(fig3, use_container_width=True)

    st.markdown("<br>", unsafe_allow_html=True)

    # Courbe taux rebut mensuel
    section("📉  ÉVOLUTION DU TAUX DE REBUT — Est-ce qu'on s'améliore ?")
    comment("Une courbe qui descend = on s'améliore. Une courbe qui monte = problème de qualité. "
            "L'objectif est de maintenir ce taux sous 3% (ligne pointillée).")
    rb_m = df_prod.groupby("Mois_Label")["Taux rebut"].mean().reset_index()
    try:
        rb_m["sk"] = pd.to_datetime(rb_m["Mois_Label"], format="%b %Y"); rb_m = rb_m.sort_values("sk")
    except: pass
    fig4 = go.Figure(go.Scatter(x=rb_m["Mois_Label"], y=rb_m["Taux rebut"], mode="lines+markers",
        line=dict(color=C["rose"], width=3), marker=dict(size=7, color=C["rose"], line=dict(color="white",width=2)),
        fill="tozeroy", fillcolor="rgba(255,107,167,0.10)",
        hovertemplate="<b>%{x}</b><br>Taux rebut : %{y:.2f}%<extra></extra>"))
    fig4.add_hline(y=3, line_dash="dot", line_color=C["rouge"], line_width=1.5,
        annotation_text="Seuil max 3%", annotation_font=dict(color=C["rouge"],size=10))
    fig4 = excel_style(fig4, 280, False)
    fig4.update_xaxes(tickangle=-35, tickfont=dict(size=8))
    fig4.update_yaxes(ticksuffix="%")
    st.plotly_chart(fig4, use_container_width=True)


# ══════════════════════════════════════════════════════════════
# PAGE 7 : PAIEMENTS
# ══════════════════════════════════════════════════════════════
elif page == "💳  Paiements":
    df_p = df_f[df_f["ETAT DE PAIEMENT"]=="Payée"]
    df_i = df_f[df_f["ETAT DE PAIEMENT"]=="Impayée"]
    taux  = len(df_p)/len(df_f)*100 if len(df_f) else 0
    mode_top = df_p["MODE DE PAIEMENT"].value_counts().idxmax() if len(df_p)>0 else "N/A"
    mode_colors = [C["bleu"],C["rose"],C["violet"],C["bleu_clair"],C["bleu_fonce"]]

    comment(f"Sur <b>{len(df_f)}</b> factures : <b>{len(df_p)}</b> payées ({taux:.1f}%) et "
            f"<b>{len(df_i)}</b> impayées. Mode dominant : <b>{mode_top}</b>. "
            f"Un bon taux de recouvrement = moins de risques financiers.", "💳 Comprendre nos paiements")

    c1,c2,c3,c4,c5 = st.columns(5)
    with c1: kpi("💰","Total Encaissé", fmt(df_p["MONTANT TTC"].sum()),"Reçus","vert")
    with c2: kpi("⏳","Total Impayé", fmt(df_i["MONTANT TTC"].sum()), f"{len(df_i)} fact.","rouge")
    with c3: kpi("📊","Taux Recouvrement", f"{taux:.1f}%","% factures payées","bleu")
    with c4: kpi("🧾","Montant Moy. Payé", fmt(df_p["MONTANT TTC"].mean()),"Par facture réglée","violet")
    with c5: kpi("🏆","Mode Dominant", mode_top.split()[0],"Le + utilisé","orange")
    st.markdown("<br>", unsafe_allow_html=True)

    # Courbe taux recouvrement mensuel
    section("📈  TAUX DE RECOUVREMENT MENSUEL — Récupère-t-on bien notre argent ?")
    comment("Au-dessus de 80% (ligne verte) = bonne santé financière. "
            "En dessous = relancer les clients. Les points rouges = mois en dessous de l'objectif.")
    df_fm = df_f.copy(); df_fm["Période"] = df_fm["DATE"].dt.to_period("M").astype(str)
    rec_m = df_fm.groupby("Période").apply(
        lambda g: g[g["ETAT DE PAIEMENT"]=="Payée"]["MONTANT TTC"].sum()/g["MONTANT TTC"].sum()*100
        if g["MONTANT TTC"].sum()>0 else 0).reset_index()
    rec_m.columns = ["Période","Taux"]
    fig = go.Figure(go.Scatter(x=rec_m["Période"], y=rec_m["Taux"], mode="lines+markers",
        line=dict(color=C["bleu"], width=3),
        marker=dict(size=9, color=[C["vert"] if v>=80 else C["rouge"] for v in rec_m["Taux"]],
                    line=dict(color="white",width=2)),
        fill="tozeroy", fillcolor="rgba(97,114,243,0.08)",
        hovertemplate="<b>%{x}</b><br>Recouvrement : %{y:.1f}%<extra></extra>"))
    fig.add_hline(y=80, line_dash="dot", line_color=C["vert"], line_width=2,
        annotation_text="Objectif 80% ✅", annotation_font=dict(color=C["vert"],size=11),
        annotation_position="top right")
    fig = excel_style(fig, 300, False)
    fig.update_xaxes(tickangle=-35, tickfont=dict(size=8))
    fig.update_yaxes(ticksuffix="%", range=[40,110])
    st.plotly_chart(fig, use_container_width=True)
    st.markdown("<br>", unsafe_allow_html=True)

    col_a, col_b = st.columns(2)
    with col_a:
        section("🥧  COMMENT NOS CLIENTS PAIENT-ILS ?")
        comment("Le virement bancaire est le plus sécurisé. "
                "Un fort pourcentage d'espèces peut rendre le suivi plus difficile.")
        md = df_p.groupby("MODE DE PAIEMENT")["MONTANT TTC"].sum().reset_index().sort_values("MONTANT TTC",ascending=False)
        fig2 = go.Figure(go.Pie(labels=md["MODE DE PAIEMENT"], values=md["MONTANT TTC"], hole=0.60,
            marker=dict(colors=mode_colors[:len(md)], line=dict(color="white",width=2)),
            textinfo="label+percent", textfont=dict(size=11),
            hovertemplate="<b>%{label}</b><br>%{value:,.0f} FCFA (%{percent})<extra></extra>"))
        fig2.update_layout(paper_bgcolor="#FFFFFF", height=300, margin=dict(t=10,b=30,l=10,r=10),
            legend=dict(orientation="h", y=-0.08, xanchor="center", x=0.5, font=dict(color=C["texte"],size=10)))
        st.plotly_chart(fig2, use_container_width=True)

    with col_b:
        section("🔴  CLIENTS LES PLUS EN RETARD")
        comment("Ces clients ont les montants impayés les plus élevés. "
                "À contacter en priorité. Plus une facture est vieille, plus elle risque d'être perdue.")
        top_i = df_i.groupby("CLIENT")["MONTANT TTC"].sum().sort_values(ascending=True).tail(8).reset_index()
        fig3 = go.Figure(go.Bar(x=top_i["MONTANT TTC"], y=top_i["CLIENT"], orientation="h",
            marker_color=C["rouge"], text=top_i["MONTANT TTC"].apply(lambda v: fmt(v,"").strip()),
            textposition="outside", textfont=dict(size=10),
            hovertemplate="<b>%{y}</b><br>Impayé : %{x:,.0f} FCFA<extra></extra>"))
        fig3 = excel_style(fig3, 300, False)
        fig3.update_xaxes(tickformat=",.0f")
        st.plotly_chart(fig3, use_container_width=True)

    st.markdown("<br>", unsafe_allow_html=True)

    # Barres groupées (non empilées) par mode et mois
    section("📅  ENCAISSEMENTS PAR MODE DE PAIEMENT — Barres séparées par mois")
    comment("Chaque couleur = un mode de paiement. Les barres sont séparées pour faciliter la lecture. "
            "Permettez-vous de voir si un mode de paiement progresse ou régresse dans le temps.")
    df_pm = df_p.copy(); df_pm["Période"] = df_pm["DATE"].dt.to_period("M").astype(str)
    mm = df_pm.groupby(["Période","MODE DE PAIEMENT"])["MONTANT TTC"].sum().reset_index()
    fig4 = go.Figure()
    for j, m in enumerate(mm["MODE DE PAIEMENT"].unique()):
        d = mm[mm["MODE DE PAIEMENT"]==m]
        fig4.add_trace(go.Bar(x=d["Période"], y=d["MONTANT TTC"], name=m,
            marker_color=mode_colors[j % len(mode_colors)],
            hovertemplate=f"<b>{m}</b><br>%{{x}}<br>%{{y:,.0f}} FCFA<extra></extra>"))
    fig4.update_layout(barmode="group")
    fig4 = excel_style(fig4, 320)
    fig4.update_xaxes(tickangle=-35, tickfont=dict(size=8))
    fig4.update_yaxes(tickformat=",.0f", ticksuffix=" F")
    st.plotly_chart(fig4, use_container_width=True)


# ══════════════════════════════════════════════════════════════


# ══════════════════════════════════════════════════════════════
# PAGE 8 : PRÉVISIONS & ANTICIPATIONS — v6 CORRIGÉE & AMÉLIORÉE
# ══════════════════════════════════════════════════════════════
elif page == "🔮  Prévisions":

    # Palette étendue pour 13 segments
    SEG_PALETTE = [
        "#6172F3","#FF6BA7","#51459E","#F39C12","#2ECC71","#3C41CD",
        "#C5CBFB","#E74C3C","#1ABC9C","#8E44AD","#2980B9","#D35400","#27AE60",
    ]

    st.markdown("""
    <div style="background:linear-gradient(135deg,#1E2A4A,#2D3A6B);border-radius:14px;
                padding:22px 28px;margin-bottom:18px;border:1px solid rgba(97,114,243,0.3);">
        <div style="font-size:1.15rem;font-weight:800;color:white;margin-bottom:4px;">
            🔮 Prévisions & Anticipations — MULTIPACK SA
        </div>
        <div style="font-size:0.8rem;color:rgba(255,255,255,0.65);">
            Modèle de régression polynomiale sur tendances historiques 2022–2024
            · Sélectionnez un onglet pour chaque module de prévision
        </div>
    </div>""", unsafe_allow_html=True)

    prevision_comment(
        "Ces prévisions sont calculées automatiquement à partir de vos données historiques. "
        "Elles <b>ne sont pas des certitudes</b>, mais des estimations basées sur vos tendances passées. "
        "La zone ombrée = fourchette basse / haute. Utilisez-les pour <b>planifier à l'avance</b>, "
        "pas pour remplacer votre jugement commercial.",
        "🔮 Comment lire ces prévisions ?"
    )

    # Slider horizon (commun à tous les onglets)
    col_sl, col_info = st.columns([1, 3])
    with col_sl:
        n_mois_prev = st.slider("Horizon de prévision (mois)", 3, 12, 6, 1)
    with col_info:
        st.markdown(f"""
        <div style="background:#F0FFF4;border-left:4px solid #2ECC71;border-radius:0 8px 8px 0;
                    padding:10px 14px;margin-top:8px;font-size:0.82rem;color:#1a5c2e;">
            <b>🗓 Horizon sélectionné : {n_mois_prev} mois</b> — 
            Les prévisions couvrent les {n_mois_prev} prochains mois à partir de janvier 2025.
            Plus l'horizon est long, plus l'incertitude augmente.
        </div>""", unsafe_allow_html=True)

    st.markdown("<br>", unsafe_allow_html=True)

    # ── TABS pour séparer les modules proprement
    tab1, tab2, tab3, tab4, tab5 = st.tabs([
        "📈 CA Global",
        "🏷️ Par Segment",
        "📦 Stocks & Ruptures",
        "💰 Trésorerie",
        "🏭 Production",
    ])

    # ════════════════════════════════════════════════════════
    # TAB 1 — PRÉVISION CA GLOBAL
    # ════════════════════════════════════════════════════════
    with tab1:
        st.markdown("<br>", unsafe_allow_html=True)
        prevision_comment(
            f"La courbe bleue = votre CA réel mois par mois depuis 2022. "
            f"La courbe orange pointillée = notre projection pour les <b>{n_mois_prev} prochains mois</b>. "
            f"La zone jaune = fourchette d'incertitude (scénario bas / scénario haut). "
            f"La ligne verticale sépare le passé du futur.",
            "📈 Comment lire ce graphique ?"
        )

        # Données complètes (pas filtrées par sidebar)
        # CA Global: utilise tout l'historique pour la régression (meilleure précision)
        df_all = df_f.copy()  # Filtré par sidebar
        df_all["Période"] = df_all["DATE"].dt.to_period("M").astype(str)
        serie_hist = df_all.groupby("Période")["MONTANT TTC"].sum().sort_index()

        prev_df = prevoir_ca(serie_hist, n_mois_prev)

        hist_x = serie_hist.index.tolist()
        hist_y = serie_hist.values.tolist()

        # Calculer confiance du modèle (R² approché)
        y_arr = np.array(hist_y)
        X_arr = np.arange(len(y_arr)).reshape(-1,1)
        from sklearn.preprocessing import PolynomialFeatures
        from sklearn.linear_model import LinearRegression
        poly = PolynomialFeatures(degree=2)
        Xp   = poly.fit_transform(X_arr)
        mdl  = LinearRegression().fit(Xp, y_arr)
        y_pred_in = mdl.predict(Xp)
        ss_res = np.sum((y_arr - y_pred_in)**2)
        ss_tot = np.sum((y_arr - y_arr.mean())**2)
        r2 = max(0, 1 - ss_res/ss_tot) if ss_tot>0 else 0
        confiance = int(r2 * 100)

        # Détection saisonnalité (mois le plus fort en historique)
        df_saison = df_all.copy()
        df_saison["MOIS"] = df_saison["DATE"].dt.month
        mois_labels = {1:"Jan",2:"Fév",3:"Mar",4:"Avr",5:"Mai",6:"Jun",
                       7:"Jul",8:"Aoû",9:"Sep",10:"Oct",11:"Nov",12:"Déc"}
        mois_fort = df_saison.groupby("MOIS")["MONTANT TTC"].mean().idxmax()
        mois_faible = df_saison.groupby("MOIS")["MONTANT TTC"].mean().idxmin()

        # Indicateurs modèle
        cm1, cm2, cm3, cm4 = st.columns(4)
        with cm1:
            couleur_r2 = C["vert"] if confiance>=70 else (C["orange"] if confiance>=50 else C["rouge"])
            st.markdown(f"""
            <div style="background:#FFFFFF;border-radius:10px;padding:14px 16px;
                        border:1px solid {C['bordure']};border-top:4px solid {couleur_r2};">
                <div style="font-size:0.68rem;color:{C['muted']};font-weight:700;text-transform:uppercase;">Fiabilité modèle</div>
                <div style="font-size:1.5rem;font-weight:800;color:{couleur_r2};">{confiance}%</div>
                <div style="font-size:0.72rem;color:{C['muted']};margin-top:2px;">
                    {"Très fiable" if confiance>=70 else ("Fiabilité correcte" if confiance>=50 else "Tendance incertaine")}
                </div>
            </div>""", unsafe_allow_html=True)
        with cm2:
            prev_kpi("🎯","CA Prévu Total", fmt(prev_df['CA_Prévu'].sum()), f"Sur {n_mois_prev} mois")
        with cm3:
            prev_kpi("📅","Mois le + fort", mois_labels[mois_fort], "Historiquement")
        with cm4:
            ca_moy = serie_hist.mean()
            evol = (prev_df["CA_Prévu"].mean() - ca_moy)/ca_moy*100 if ca_moy else 0
            prev_kpi("📊","Tendance prévue", f"{evol:+.1f}%", "vs moyenne historique")

        st.markdown("<br>", unsafe_allow_html=True)

        # Graphique principal
        fig = go.Figure()
        # Zone confiance
        fig.add_trace(go.Scatter(
            x=list(prev_df["Période"]) + list(prev_df["Période"])[::-1],
            y=list(prev_df["Haut"])    + list(prev_df["Bas"])[::-1],
            fill="toself", fillcolor="rgba(247,179,49,0.18)",
            line=dict(color="rgba(0,0,0,0)"), name="Fourchette basse–haute",
            hoverinfo="skip", showlegend=True))
        # Historique réel
        fig.add_trace(go.Scatter(
            x=hist_x, y=hist_y, name="CA Réel (historique)",
            mode="lines+markers",
            line=dict(color=C["bleu"], width=2.5),
            marker=dict(size=5, color=C["bleu"], line=dict(color="white",width=1.5)),
            fill="tozeroy", fillcolor="rgba(97,114,243,0.07)",
            hovertemplate="<b>%{x}</b><br>CA réel : <b>%{y:,.0f} FCFA</b><extra></extra>"))
        # Scénario bas
        fig.add_trace(go.Scatter(
            x=prev_df["Période"], y=prev_df["Bas"],
            name="Scénario pessimiste", mode="lines",
            line=dict(color=C["rouge"], width=1.5, dash="dot"),
            hovertemplate="<b>%{x}</b><br>Scénario bas : %{y:,.0f} FCFA<extra></extra>"))
        # Prévision centrale
        fig.add_trace(go.Scatter(
            x=prev_df["Période"], y=prev_df["CA_Prévu"],
            name="Prévision centrale", mode="lines+markers",
            line=dict(color=C["orange"], width=3, dash="dash"),
            marker=dict(size=10, color=C["orange"], symbol="diamond", line=dict(color="white",width=2)),
            hovertemplate="<b>Prévision %{x}</b><br>CA estimé : <b>%{y:,.0f} FCFA</b><extra></extra>"))
        # Scénario haut
        fig.add_trace(go.Scatter(
            x=prev_df["Période"], y=prev_df["Haut"],
            name="Scénario optimiste", mode="lines",
            line=dict(color=C["vert"], width=1.5, dash="dot"),
            hovertemplate="<b>%{x}</b><br>Scénario haut : %{y:,.0f} FCFA<extra></extra>"))

        # Ligne de séparation passé/futur (sur index numérique, pas catégoriel)
        n_hist = len(hist_x)
        all_x = hist_x + list(prev_df["Période"])
        fig.add_vrect(
            x0=hist_x[-1], x1=prev_df["Période"].iloc[-1],
            fillcolor="rgba(247,179,49,0.05)", layer="below",
            line_width=0, annotation_text="← Historique | Prévision →",
            annotation_position="top left",
            annotation_font=dict(size=10, color=C["muted"]))

        fig = excel_style(fig, 420)
        fig.update_xaxes(tickangle=-35, tickfont=dict(size=8))
        fig.update_yaxes(tickformat=",.0f", ticksuffix=" F")
        fig.update_layout(legend=dict(orientation="h", y=-0.18, x=0.5, xanchor="center"))
        st.plotly_chart(fig, use_container_width=True)

        # Tableau mois par mois
        with st.expander("📋 Voir le détail mois par mois"):
            prev_disp = prev_df.copy()
            prev_disp["CA_Prévu"] = prev_disp["CA_Prévu"].apply(lambda v: f"{v:,.0f} FCFA")
            prev_disp["Bas"]      = prev_disp["Bas"].apply(lambda v: f"{v:,.0f} FCFA")
            prev_disp["Haut"]     = prev_disp["Haut"].apply(lambda v: f"{v:,.0f} FCFA")
            prev_disp = prev_disp.rename(columns={"Période":"Mois","CA_Prévu":"CA Central",
                                                    "Bas":"Scénario Pessimiste","Haut":"Scénario Optimiste"})
            st.dataframe(prev_disp, use_container_width=True)

        # Export Excel des prévisions
        buf_prev = io.BytesIO()
        with pd.ExcelWriter(buf_prev, engine="openpyxl") as w:
            prev_df.to_excel(w, sheet_name="Prévisions CA", index=False)
        buf_prev.seek(0)
        st.download_button("⬇️ Exporter les prévisions CA (.xlsx)", buf_prev,
            file_name=f"MULTIPACK_Previsions_CA_{n_mois_prev}mois.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")


    # ════════════════════════════════════════════════════════
    # TAB 2 — PRÉVISION PAR SEGMENT
    # ════════════════════════════════════════════════════════
    with tab2:
        st.markdown("<br>", unsafe_allow_html=True)
        prevision_comment(
            "Chaque courbe = la prévision d'un segment. "
            "Cliquez sur un segment dans la légende pour l'afficher seul. "
            "Les KPIs en bas montrent qui va progresser (🟢) et qui va reculer (🔴).",
            "🏷️ Comment lire ce graphique ?"
        )

        df_fa = df_f.copy()  # Respecte les filtres sidebar
        df_fa["Période"] = df_fa["DATE"].dt.to_period("M").astype(str)
        all_segs = sorted(df_fa["SEGMENT"].unique().tolist())

        # Filtre segments pour prévisions
        segs_prev = st.multiselect(
            "Sélectionnez les segments à projeter",
            all_segs, default=all_segs[:6],
            help="Choisissez jusqu'à 8 segments. Trop de courbes = illisible.")

        if not segs_prev:
            st.warning("Sélectionnez au moins un segment.")
        else:
            seg_color_map = {s: SEG_PALETTE[i % len(SEG_PALETTE)] for i, s in enumerate(all_segs)}

            # GRAPHIQUE 1 — Courbes prévision par segment (séparées)
            fig2 = go.Figure()
            seg_summary = []
            for seg in segs_prev:
                serie_seg = df_fa[df_fa["SEGMENT"]==seg].groupby("Période")["MONTANT TTC"].sum().sort_index()
                if len(serie_seg) < 4:
                    continue
                p = prevoir_ca(serie_seg, n_mois_prev)
                c = seg_color_map.get(seg, C["bleu"])
                fig2.add_trace(go.Scatter(
                    x=p["Période"], y=p["CA_Prévu"], name=seg,
                    mode="lines+markers",
                    line=dict(color=c, width=2.5, dash="dash"),
                    marker=dict(size=8, symbol="diamond", color=c, line=dict(color="white",width=1.5)),
                    hovertemplate=f"<b>{seg}</b><br>%{{x}}<br>Prévision : <b>%{{y:,.0f}} FCFA</b><extra></extra>"))
                moy_h = serie_seg.mean()
                moy_p = p["CA_Prévu"].mean()
                tend  = (moy_p - moy_h) / moy_h * 100 if moy_h else 0
                seg_summary.append({"seg": seg, "tend": tend, "moy_p": moy_p,
                                     "total_p": p["CA_Prévu"].sum(), "color": c})

            fig2 = excel_style(fig2, 360)
            fig2.update_xaxes(tickangle=-30, tickfont=dict(size=9))
            fig2.update_yaxes(tickformat=",.0f", ticksuffix=" F")
            fig2.update_layout(legend=dict(orientation="h", y=-0.22, x=0.5, xanchor="center"))
            st.plotly_chart(fig2, use_container_width=True)

            # KPIs tendance — MAX 4 par ligne pour éviter crash
            if seg_summary:
                st.markdown("<br>", unsafe_allow_html=True)
                section("📊  TENDANCE PRÉVUE PAR SEGMENT")
                comment("🟢 Croissance prévue · 🔴 Recul prévu · ➡️ Stable. "
                        "Ces tendances sont basées sur votre historique de ventes.")
                # Afficher en rangées de 4
                chunk = 4
                for row_start in range(0, len(seg_summary), chunk):
                    row_segs = seg_summary[row_start:row_start+chunk]
                    cols = st.columns(len(row_segs))
                    for i, s in enumerate(row_segs):
                        icon  = "📈" if s["tend"]>3 else ("📉" if s["tend"]<-3 else "➡️")
                        color = "vert" if s["tend"]>3 else ("rouge" if s["tend"]<-3 else "orange")
                        with cols[i]:
                            kpi(icon, s["seg"][:22], f"{s['tend']:+.1f}%",
                                f"CA prévu : {fmt(s['total_p'],'').strip()} sur {n_mois_prev} mois", color)
                    st.markdown("<div style='margin:6px 0'></div>", unsafe_allow_html=True)

                st.markdown("<br>", unsafe_allow_html=True)

                # GRAPHIQUE 2 — Barres CA prévu par segment (classement)
                section("🏆  CLASSEMENT PRÉVISIONNEL — Quel segment va générer le plus ?")
                seg_sum_df = pd.DataFrame(seg_summary).sort_values("total_p", ascending=True)
                clr_bar = [s["color"] for s in seg_summary]
                clr_bar_sorted = [seg_color_map.get(s, C["bleu"]) for s in seg_sum_df["seg"]]
                fig_bar = go.Figure(go.Bar(
                    x=seg_sum_df["total_p"], y=seg_sum_df["seg"],
                    orientation="h",
                    marker_color=clr_bar_sorted,
                    text=seg_sum_df["total_p"].apply(lambda v: fmt(v,"").strip()),
                    textposition="outside", textfont=dict(size=10),
                    hovertemplate="<b>%{y}</b><br>CA prévu : %{x:,.0f} FCFA<extra></extra>"))
                fig_bar = excel_style(fig_bar, max(280, len(segs_prev)*45), False)
                fig_bar.update_xaxes(tickformat=",.0f")
                st.plotly_chart(fig_bar, use_container_width=True)


    # ════════════════════════════════════════════════════════
    # TAB 3 — STOCKS & RUPTURES
    # ════════════════════════════════════════════════════════
    with tab3:
        st.markdown("<br>", unsafe_allow_html=True)
        prevision_comment(
            "Basé sur la consommation mensuelle moyenne de chaque produit, "
            "nous calculons combien de mois de stock il reste avant d'atteindre le seuil d'alerte. "
            "🔴 < 2 mois = commander maintenant · 🟡 2–4 mois = planifier · 🟢 > 4 mois = confortable.",
            "📦 Comment lire ce graphique ?"
        )

        stock_rows = []
        for _, row in df_inv.iterrows():
            sf  = float(row.get("Stock final", 0))
            sor = float(row.get("Sorties_moy_mois", 0))
            if pd.isna(sor) or sor <= 0:
                sor = 1.0
            seuil = float(row.get("seuil", 0))
            mois_r = (sf - seuil) / sor if sor > 0 else 99.0
            mois_r = max(0.0, round(float(mois_r), 1))
            stock_rows.append({
                "Produit":       str(row.get("designation",""))[:38],
                "Catégorie":     str(row.get("categorie","")),
                "Stock actuel":  int(sf),
                "Conso/mois":    round(sor, 1),
                "Mois restants": mois_r,
                "Seuil alerte":  int(seuil),
                "Statut":        str(row.get("Statut","Stock normal")),
            })

        df_sa = pd.DataFrame(stock_rows).sort_values("Mois restants")

        # GRAPHIQUE — barres horizontales mois restants
        df_graph = df_sa[df_sa["Mois restants"] < 12].head(20).sort_values("Mois restants", ascending=True)

        if df_graph.empty:
            st.success("✅ Tous les stocks sont largement suffisants (plus de 12 mois de stock).")
        else:
            clr_s = [C["rouge"] if v < 2 else (C["orange"] if v < 4 else C["vert"])
                     for v in df_graph["Mois restants"]]
            fig3 = go.Figure(go.Bar(
                x=df_graph["Mois restants"], y=df_graph["Produit"],
                orientation="h", marker_color=clr_s,
                text=[f"{v:.1f} mois" for v in df_graph["Mois restants"]],
                textposition="outside", textfont=dict(size=11),
                hovertemplate="<b>%{y}</b><br>Stock restant : %{x:.1f} mois<br>"
                              "Conso/mois : " + df_graph["Conso/mois"].astype(str) + " u.<extra></extra>"))
            max_x = max(df_graph["Mois restants"].max() * 1.3, 6)
            fig3.add_vline(x=2, line_dash="dot", line_color=C["rouge"], line_width=2)
            fig3.add_vline(x=4, line_dash="dot", line_color=C["orange"], line_width=2)
            # Annotations zones
            fig3.add_annotation(x=1, y=len(df_graph)-0.5, text="🔴 Urgent", showarrow=False,
                font=dict(size=10, color=C["rouge"]), bgcolor="rgba(255,255,255,0.8)")
            fig3.add_annotation(x=3, y=len(df_graph)-0.5, text="🟡 Planifier", showarrow=False,
                font=dict(size=10, color=C["orange"]), bgcolor="rgba(255,255,255,0.8)")
            fig3.add_annotation(x=max(6, max_x*0.7), y=len(df_graph)-0.5, text="🟢 Confortable", showarrow=False,
                font=dict(size=10, color=C["vert"]), bgcolor="rgba(255,255,255,0.8)")
            fig3 = excel_style(fig3, max(360, len(df_graph)*32), False)
            fig3.update_xaxes(ticksuffix=" mois", range=[0, max_x])
            st.plotly_chart(fig3, use_container_width=True)

        st.markdown("<br>", unsafe_allow_html=True)

        # KPIs stock
        urgents    = df_sa[df_sa["Mois restants"] < 2]
        attentions = df_sa[(df_sa["Mois restants"] >= 2) & (df_sa["Mois restants"] < 4)]
        confort    = df_sa[df_sa["Mois restants"] >= 4]
        ck1, ck2, ck3 = st.columns(3)
        with ck1: kpi("🔴","Commander maintenant", str(len(urgents)), "Produits critiques (< 2 mois)","rouge")
        with ck2: kpi("🟡","Planifier bientôt", str(len(attentions)), "Produits à risque (2–4 mois)","orange")
        with ck3: kpi("✅","Stock confortable", str(len(confort)), "Produits bien approvisionnés","vert")

        st.markdown("<br>", unsafe_allow_html=True)

        # Alertes textuelles
        if len(urgents) > 0:
            st.markdown(f'<div class="sec">🚨  PRODUITS À COMMANDER MAINTENANT</div>', unsafe_allow_html=True)
            for _, r in urgents.iterrows():
                st.markdown(f"""
                <div class="alert-r">🔴<div>
                    <b>{r['Produit']}</b> <span style="color:#999;font-size:0.8em">({r['Catégorie']})</span><br>
                    <small>Stock actuel : <b>{r['Stock actuel']:,} u.</b> · Consommation mensuelle : ~{r['Conso/mois']:.0f} u./mois
                    · Épuisement dans <b>{r['Mois restants']:.1f} mois</b>
                    → <b>Action : Lancer commande dans les 48h</b></small>
                </div></div>""", unsafe_allow_html=True)

        if len(attentions) > 0:
            st.markdown(f'<div class="sec">⚠️  PRODUITS À RÉAPPROVISIONNER BIENTÔT</div>', unsafe_allow_html=True)
            for _, r in attentions.head(6).iterrows():
                st.markdown(f"""
                <div class="alert-y">📅<div>
                    <b>{r['Produit']}</b> <span style="color:#999;font-size:0.8em">({r['Catégorie']})</span><br>
                    <small>Stock pour <b>{r['Mois restants']:.1f} mois</b>
                    · Conso. : ~{r['Conso/mois']:.0f} u./mois
                    → Planifier commande dans les 2–4 semaines</small>
                </div></div>""", unsafe_allow_html=True)

        # Tableau complet
        with st.expander("📋 Voir le tableau complet des stocks prévisionnels"):
            st.dataframe(df_sa, use_container_width=True, height=400)


    # ════════════════════════════════════════════════════════
    # TAB 4 — TRÉSORERIE PRÉVISIONNELLE
    # ════════════════════════════════════════════════════════
    with tab4:
        st.markdown("<br>", unsafe_allow_html=True)

        # Utiliser les prévisions CA calculées en tab1 (recalcul propre)
        df_all2 = df_fact.copy()
        df_all2["Période"] = df_all2["DATE"].dt.to_period("M").astype(str)
        serie_h2 = df_all2.groupby("Période")["MONTANT TTC"].sum().sort_index()
        prev_df2 = prevoir_ca(serie_h2, n_mois_prev)

        taux_hist = (df_fact[df_fact["ETAT DE PAIEMENT"]=="Payée"]["MONTANT TTC"].sum() /
                     df_fact["MONTANT TTC"].sum())

        prev_enc = prev_df2["CA_Prévu"] * taux_hist
        prev_imp = prev_df2["CA_Prévu"] * (1 - taux_hist)

        prevision_comment(
            f"En appliquant votre taux de recouvrement historique ({taux_hist*100:.0f}%) au CA prévu, "
            f"voici ce que vous pouvez anticiper comme <b>entrées de trésorerie</b>. "
            f"Ces chiffres permettent de planifier vos achats, salaires et investissements.",
            "💰 Comment lire la trésorerie prévisionnelle ?"
        )

        enc_total = prev_enc.sum()
        imp_total = prev_imp.sum()

        # KPIs
        tk1, tk2, tk3, tk4 = st.columns(4)
        with tk1: prev_kpi("💚","Encaissements Prévus", fmt(enc_total), f"Sur {n_mois_prev} mois")
        with tk2: prev_kpi("⏳","Risque Impayé Estimé", fmt(imp_total), f"À relancer activement")
        with tk3: prev_kpi("📊","CA Prévu Total", fmt(prev_df2['CA_Prévu'].sum()), "Ventes estimées")
        with tk4: prev_kpi("💹","Taux de Recouvrement", f"{taux_hist*100:.1f}%", "Taux historique appliqué")

        st.markdown("<br>", unsafe_allow_html=True)

        # GRAPHIQUE 1 — Barres empilées encaissé / risque impayé
        section("💰  ENCAISSEMENTS PRÉVUS VS RISQUE IMPAYÉ — Mois par mois")
        comment("Bleu = ce qu'on devrait encaisser · Rose = risque d'impayé. "
                "Ces chiffres sont basés sur votre taux historique de paiement.")
        fig4a = go.Figure()
        fig4a.add_trace(go.Bar(x=prev_df2["Période"], y=prev_enc,
            name="Encaissements prévus", marker_color=C["bleu"], opacity=0.9,
            text=prev_enc.apply(lambda v: fmt(v,"").strip()),
            textposition="inside", textfont=dict(size=9, color="white"),
            hovertemplate="<b>%{x}</b><br>Encaissé prévu : <b>%{y:,.0f} FCFA</b><extra></extra>"))
        fig4a.add_trace(go.Bar(x=prev_df2["Période"], y=prev_imp,
            name="Risque impayé", marker_color=C["rose"], opacity=0.85,
            text=prev_imp.apply(lambda v: fmt(v,"").strip()),
            textposition="inside", textfont=dict(size=9, color="white"),
            hovertemplate="<b>%{x}</b><br>Risque impayé : %{y:,.0f} FCFA<extra></extra>"))
        fig4a.update_layout(barmode="stack")
        fig4a = excel_style(fig4a, 320)
        fig4a.update_xaxes(tickangle=-30, tickfont=dict(size=9))
        fig4a.update_yaxes(tickformat=",.0f", ticksuffix=" F")
        st.plotly_chart(fig4a, use_container_width=True)

        st.markdown("<br>", unsafe_allow_html=True)

        # GRAPHIQUE 2 — Courbe cumulative trésorerie
        section("📈  TRÉSORERIE CUMULÉE PRÉVUE — Évolution du total encaissé")
        comment("Cette courbe montre le total cumulé des encaissements sur la période. "
                "Elle permet de savoir combien vous aurez encaissé au total à la fin de chaque mois.")
        cum_enc = np.cumsum(prev_enc.values)
        fig4b = go.Figure(go.Scatter(
            x=prev_df2["Période"], y=cum_enc,
            mode="lines+markers+text",
            line=dict(color=C["vert"], width=3),
            marker=dict(size=10, color=C["vert"], line=dict(color="white",width=2)),
            fill="tozeroy", fillcolor="rgba(46,204,113,0.10)",
            text=[fmt(v,"").strip() for v in cum_enc],
            textposition="top center", textfont=dict(size=9, color=C["texte"]),
            hovertemplate="<b>%{x}</b><br>Total cumulé encaissé : <b>%{y:,.0f} FCFA</b><extra></extra>"))
        fig4b = excel_style(fig4b, 280, False)
        fig4b.update_xaxes(tickangle=-30, tickfont=dict(size=9))
        fig4b.update_yaxes(tickformat=",.0f", ticksuffix=" F")
        st.plotly_chart(fig4b, use_container_width=True)

        st.markdown("<br>", unsafe_allow_html=True)

        # Top clients impayés à risque
        section("⚠️  CLIENTS LES PLUS À RISQUE — Impayés actuels à recouvrer en priorité")
        comment("Ces clients ont des factures impayées actuellement. "
                "Commencez par les relancer pour améliorer votre trésorerie avant même les prévisions futures.")
        imp_clients = df_fact[df_fact["ETAT DE PAIEMENT"]=="Impayée"].groupby("CLIENT").agg(
            Montant=("MONTANT TTC","sum"), Nb=("MONTANT TTC","count")).reset_index()
        imp_clients["SEGMENT"] = df_fact[df_fact["ETAT DE PAIEMENT"]=="Impayée"].groupby("CLIENT")["SEGMENT"].first().values
        imp_clients = imp_clients.sort_values("Montant", ascending=False).head(10)
        if not imp_clients.empty:
            fig4c = go.Figure(go.Bar(
                x=imp_clients["Montant"], y=imp_clients["CLIENT"],
                orientation="h",
                marker_color=[C["rouge"] if v==imp_clients["Montant"].max()
                               else (C["rose"] if v>=imp_clients["Montant"].quantile(0.6) else C["orange"])
                               for v in imp_clients["Montant"]],
                text=imp_clients["Montant"].apply(lambda v: fmt(v,"").strip()),
                textposition="outside", textfont=dict(size=10),
                hovertemplate="<b>%{y}</b><br>Impayé : %{x:,.0f} FCFA<extra></extra>"))
            fig4c = excel_style(fig4c, 340, False)
            fig4c.update_xaxes(tickformat=",.0f")
            st.plotly_chart(fig4c, use_container_width=True)

        prevision_comment(
            f"Pour améliorer la trésorerie : "
            f"(1) Relancez les {len(imp_clients)} clients avec impayés — potentiel : <b>{fmt(imp_clients['Montant'].sum())}</b>. "
            f"(2) Exigez des acomptes sur les nouvelles grandes commandes. "
            f"(3) Réduisez les délais de paiement dans vos contrats (de 60j à 30j).",
            "💡 Actions pour améliorer la trésorerie"
        )


    # ════════════════════════════════════════════════════════
    # TAB 5 — PRODUCTION PRÉVISIONNELLE
    # ════════════════════════════════════════════════════════
    with tab5:
        st.markdown("<br>", unsafe_allow_html=True)
        prevision_comment(
            "La courbe violette = production réelle historique. "
            "La courbe orange pointillée = prévision. "
            "Si la prévision monte → préparer plus de matières premières et de personnel. "
            "Si elle descend → revoir les objectifs de production.",
            "🏭 Comment lire la prévision de production ?"
        )

        # Tri chronologique correct des mois
        prod_all = df_prod.groupby("Mois").agg(
            Production=("Production réelle","sum"),
            Mois_Label=("Mois_Label","first")).reset_index()
        prod_all = prod_all.sort_values("Mois")  # tri par "2022-01" etc — correct

        serie_prod = pd.Series(prod_all["Production"].values, index=prod_all["Mois_Label"])
        prev_prod  = prevoir_ca(serie_prod, n_mois_prev)

        prod_prev_moy  = prev_prod["CA_Prévu"].mean()
        prod_hist_moy  = prod_all["Production"].mean()
        evol_prod = (prod_prev_moy - prod_hist_moy) / prod_hist_moy * 100 if prod_hist_moy else 0

        pk1, pk2, pk3, pk4 = st.columns(4)
        with pk1: prev_kpi("⚙️","Prod. Totale Prévue", f"{prev_prod['CA_Prévu'].sum():,.0f}", f"Unités sur {n_mois_prev} mois")
        with pk2: prev_kpi("📊","Tendance", f"{evol_prod:+.1f}%", "vs moyenne historique")
        with pk3: prev_kpi("🏭","Moy. Mensuelle Prévue", f"{prod_prev_moy:,.0f}", "Unités par mois")
        with pk4: prev_kpi("🗑️","Rebuts Estimés", f"{int(prod_prev_moy*df_prod['Taux rebut'].mean()/100):,}", "Unités / mois (taux actuel)")

        st.markdown("<br>", unsafe_allow_html=True)

        # GRAPHIQUE 1 — Courbe production + prévision avec fourchette
        section("🏭  PRODUCTION RÉELLE ET PRÉVISION — Évolution des volumes")
        fig5a = go.Figure()
        fig5a.add_trace(go.Scatter(
            x=list(prev_prod["Période"]) + list(prev_prod["Période"])[::-1],
            y=list(prev_prod["Haut"])    + list(prev_prod["Bas"])[::-1],
            fill="toself", fillcolor="rgba(247,179,49,0.15)",
            line=dict(color="rgba(0,0,0,0)"), name="Fourchette probable",
            hoverinfo="skip"))
        fig5a.add_trace(go.Scatter(
            x=prod_all["Mois_Label"], y=prod_all["Production"],
            name="Production réelle", mode="lines+markers",
            line=dict(color=C["violet"], width=2.5),
            marker=dict(size=5, color=C["violet"], line=dict(color="white",width=1.5)),
            fill="tozeroy", fillcolor="rgba(81,69,158,0.07)",
            hovertemplate="<b>%{x}</b><br>Production : <b>%{y:,.0f} u.</b><extra></extra>"))
        fig5a.add_trace(go.Scatter(
            x=prev_prod["Période"], y=prev_prod["CA_Prévu"],
            name="Prévision", mode="lines+markers",
            line=dict(color=C["orange"], width=3, dash="dash"),
            marker=dict(size=9, symbol="diamond", color=C["orange"], line=dict(color="white",width=2)),
            hovertemplate="<b>Prévision %{x}</b><br>%{y:,.0f} unités prévues<extra></extra>"))
        fig5a = excel_style(fig5a, 380)
        fig5a.update_xaxes(tickangle=-35, tickfont=dict(size=8))
        fig5a.update_yaxes(tickformat=",")
        fig5a.update_layout(legend=dict(orientation="h", y=-0.18, x=0.5, xanchor="center"))
        st.plotly_chart(fig5a, use_container_width=True)

        st.markdown("<br>", unsafe_allow_html=True)

        # GRAPHIQUE 2 — Rendement prévu par machine (basé sur tendance)
        section("⚙️  RENDEMENT PRÉVU PAR MACHINE — Objectifs à tenir")
        comment("Basé sur les tendances de rendement historiques par machine. "
                "Les machines proches de 85% sont dans la norme. "
                "En dessous = planifier une maintenance préventive avant la période de forte production.")
        rm_hist = df_prod.groupby("Machine")["Taux rendement"].agg(["mean","std","min","max"]).reset_index()
        rm_hist.columns = ["Machine","Moy","Std","Min","Max"]
        rm_hist = rm_hist.sort_values("Moy")

        fig5b = go.Figure()
        fig5b.add_trace(go.Bar(
            x=rm_hist["Moy"], y=rm_hist["Machine"],
            orientation="h",
            marker_color=[C["vert"] if v>=85 else (C["orange"] if v>=75 else C["rouge"])
                          for v in rm_hist["Moy"]],
            error_x=dict(type="data", array=rm_hist["Std"].tolist(), visible=True,
                         color=C["muted"], thickness=1.5, width=4),
            text=rm_hist["Moy"].apply(lambda v: f"{v:.1f}%"),
            textposition="outside", textfont=dict(size=11),
            hovertemplate="<b>%{y}</b><br>Rendement moy. : %{x:.1f}% ± " +
                          rm_hist["Std"].apply(lambda v: f"{v:.1f}").astype(str) + "%<extra></extra>"))
        fig5b.add_vline(x=85, line_dash="dot", line_color=C["rouge"], line_width=2,
            annotation_text="Objectif 85%", annotation_font=dict(color=C["rouge"],size=10))
        fig5b = excel_style(fig5b, 300, False)
        fig5b.update_xaxes(ticksuffix="%", range=[55, 105])
        st.plotly_chart(fig5b, use_container_width=True)

        prevision_comment(
            f"Pour la période prévue, nous estimons une production de <b>{prev_prod['CA_Prévu'].sum():,.0f} unités</b>. "
            f"Préparez les matières premières en conséquence et planifiez les maintenances préventives "
            f"sur les machines dont le rendement est sous 80%.",
            "💡 Actions recommandées pour la production"
        )


# PAGE COMPARAISONS & ANALYSES  (insérée avant Alertes)
# ══════════════════════════════════════════════════════════════
elif page == "📈  Comparaisons":

    st.markdown(f"""
    <div style="background:linear-gradient(135deg,#1E2A4A,#2D3A6B);border-radius:14px;
                padding:22px 28px;margin-bottom:18px;border:1px solid rgba(97,114,243,0.3);">
        <div style="font-size:1.15rem;font-weight:800;color:white;margin-bottom:4px;">
            📊 Comparaisons & Analyses — Tous les segments en face à face
        </div>
        <div style="font-size:0.8rem;color:rgba(255,255,255,0.65);">
            Sélectionnez librement les segments, zones, années et indicateurs à comparer
        </div>
    </div>""", unsafe_allow_html=True)

    comment(
        "Cette page vous permet de comparer librement <b>n'importe quels segments entre eux</b> "
        "sur les indicateurs qui vous intéressent : CA, paiements, nombre de clients, taux de recouvrement. "
        "Vous pouvez comparer <b>2 à 2</b> pour une analyse précise, ou <b>tous ensemble</b> pour une vue globale. "
        "Utilisez les filtres ci-dessous pour personnaliser votre analyse.",
        "📊 Comment utiliser cette page ?"
    )

    # ── PANNEAU DE CONTRÔLE DE COMPARAISON ─────────────────────────
    st.markdown(f'<div class="sec">⚙️  CONFIGUREZ VOTRE COMPARAISON</div>', unsafe_allow_html=True)

    ctrl1, ctrl2, ctrl3 = st.columns([2, 1, 1])
    with ctrl1:
        comp_segments = st.multiselect(
            "🏷️ Segments à comparer (sélectionnez 2 ou plus)",
            ALL_SEGMENTS,
            default=ALL_SEGMENTS[:4],
            help="Choisissez les segments à mettre en face à face. Minimum 2, pas de maximum."
        )
    with ctrl2:
        comp_annees = st.multiselect("🗓 Années", ALL_ANNEES, default=ALL_ANNEES)
    with ctrl3:
        comp_zones = st.multiselect("🗺 Zones", ALL_ZONES, default=ALL_ZONES)

    col_m1, col_m2 = st.columns(2)
    with col_m1:
        metrique = st.selectbox("📐 Indicateur principal",
            ["CA Total (FCFA)", "CA Payé (FCFA)", "CA Impayé (FCFA)",
             "Nombre de factures", "Taux de recouvrement (%)",
             "Panier moyen (FCFA)", "Nombre de clients actifs"])
    with col_m2:
        mode_comp = st.radio("🔀 Mode de comparaison",
            ["Tous les segments sélectionnés", "Comparaison 2 à 2"],
            horizontal=True)

    if len(comp_segments) < 2:
        st.warning("⚠️ Sélectionnez au moins 2 segments pour lancer une comparaison.")
        st.stop()

    # Données filtrées pour comparaison
    df_comp = df_fact[
        df_fact["SEGMENT"].isin(comp_segments) &
        df_fact["ANNEE"].isin(comp_annees) &
        df_fact["ZONE"].isin(comp_zones)
    ].copy()

    if df_comp.empty:
        st.warning("Aucune donnée pour ces filtres. Modifiez votre sélection.")
        st.stop()

    df_comp["Période"] = df_comp["DATE"].dt.to_period("M").astype(str)

    # ── Calculer toutes les métriques par segment
    def calc_metrics(df):
        grp = df.groupby("SEGMENT")
        ca = grp["MONTANT TTC"].sum()
        ca_p = df[df["ETAT DE PAIEMENT"]=="Payée"].groupby("SEGMENT")["MONTANT TTC"].sum()
        ca_i = df[df["ETAT DE PAIEMENT"]=="Impayée"].groupby("SEGMENT")["MONTANT TTC"].sum()
        nb_f = grp["MONTANT TTC"].count()
        panier = grp["MONTANT TTC"].mean()
        nb_cl = grp["CLIENT"].nunique()
        taux_r = (ca_p / ca * 100).fillna(0)
        df_out = pd.DataFrame({"CA Total": ca, "CA Payé": ca_p, "CA Impayé": ca_i,
                               "Nb Factures": nb_f, "Panier Moyen": panier,
                               "Nb Clients": nb_cl, "Taux Recouvrement": taux_r}).fillna(0)
        return df_out.reindex(comp_segments).fillna(0)

    metrics_df = calc_metrics(df_comp)

    metrique_col = {
        "CA Total (FCFA)": "CA Total", "CA Payé (FCFA)": "CA Payé",
        "CA Impayé (FCFA)": "CA Impayé", "Nombre de factures": "Nb Factures",
        "Taux de recouvrement (%)": "Taux Recouvrement",
        "Panier moyen (FCFA)": "Panier Moyen", "Nombre de clients actifs": "Nb Clients"
    }
    col_key = metrique_col[metrique]
    is_pct = "%" in metrique

    # ═══════════════════════════════════════
    # MODE 1 : TOUS LES SEGMENTS ENSEMBLE
    # ═══════════════════════════════════════
    if mode_comp == "Tous les segments sélectionnés":

        # ── KPIs de synthèse globale
        section(f"📋  SYNTHÈSE — {len(comp_segments)} segments comparés")
        kpi_cols = st.columns(min(len(comp_segments), 5))
        seg_colors_list = PIE + ["#4472C4","#ED7D31","#A5A5A5","#FFC000"]
        for i, seg in enumerate(comp_segments[:5]):
            val = metrics_df.loc[seg, col_key] if seg in metrics_df.index else 0
            valstr = f"{val:.1f}%" if is_pct else fmt(val,"").strip() if val > 1000 else f"{val:,.0f}"
            with kpi_cols[i]:
                kpi("🏷️", seg[:22], valstr,
                    f"{metrics_df.loc[seg,'Nb Factures']:.0f} factures · {metrics_df.loc[seg,'Nb Clients']:.0f} clients"
                    if seg in metrics_df.index else "",
                    ["bleu","rose","violet","orange","vert","or"][i % 6])

        st.markdown("<br>", unsafe_allow_html=True)

        # ── GRAPHIQUE 1 : Barres horizontales comparaison indicateur principal
        section(f"📊  {metrique.upper()} — Classement de tous les segments")
        comment(f"Chaque barre = un segment. Plus la barre est longue = meilleure performance sur <b>{metrique}</b>. "
                f"Ce graphique permet de voir d'un coup d'œil qui est en tête et qui est en retard.")

        m_sorted = metrics_df[[col_key]].sort_values(col_key, ascending=True)
        clrs = [seg_colors_list[comp_segments.index(s) % len(seg_colors_list)]
                for s in m_sorted.index]
        text_vals = [f"{v:.1f}%" if is_pct else fmt(v,"").strip() for v in m_sorted[col_key]]

        fig = go.Figure(go.Bar(
            x=m_sorted[col_key], y=m_sorted.index,
            orientation="h", marker_color=clrs,
            text=text_vals, textposition="outside", textfont=dict(size=11),
            hovertemplate="<b>%{y}</b><br>" + metrique + " : %{x:,.0f}<extra></extra>"))
        fig = excel_style(fig, max(340, len(comp_segments)*50), False)
        if is_pct:
            fig.update_xaxes(ticksuffix="%")
        else:
            fig.update_xaxes(tickformat=",.0f")
        st.plotly_chart(fig, use_container_width=True)

        st.markdown("<br>", unsafe_allow_html=True)

        # ── GRAPHIQUE 2 : Part de chaque segment dans le CA total (donut)
        section("🥧  PART DE CHAQUE SEGMENT DANS LE CA TOTAL")
        comment("Ce graphique montre comment le chiffre d'affaires est réparti entre les segments sélectionnés. "
                "Un segment qui occupe une grande part = notre principal moteur. "
                "Si un seul segment domine trop = risque de dépendance.")
        col_d1, col_d2 = st.columns([1.2, 1.8])
        with col_d1:
            ca_seg = metrics_df["CA Total"].sort_values(ascending=False)
            fig2 = go.Figure(go.Pie(
                labels=ca_seg.index, values=ca_seg.values, hole=0.55,
                marker=dict(colors=[seg_colors_list[comp_segments.index(s) % len(seg_colors_list)]
                                    for s in ca_seg.index],
                            line=dict(color="white",width=2)),
                textinfo="label+percent", textfont=dict(size=10),
                hovertemplate="<b>%{label}</b><br>%{value:,.0f} FCFA (%{percent})<extra></extra>"))
            fig2.update_layout(paper_bgcolor="#FFFFFF", height=340,
                margin=dict(t=10,b=30,l=10,r=10),
                legend=dict(orientation="v", x=1.02, y=0.5,
                            font=dict(color=C["texte"],size=9)))
            st.plotly_chart(fig2, use_container_width=True)

        with col_d2:
            # Tableau récap multi-indicateurs
            section("📋  TABLEAU COMPARATIF COMPLET")
            disp = metrics_df.copy()
            for c in ["CA Total","CA Payé","CA Impayé","Panier Moyen"]:
                disp[c] = disp[c].apply(lambda v: f"{v:,.0f} FCFA")
            disp["Nb Factures"] = disp["Nb Factures"].apply(lambda v: f"{int(v)}")
            disp["Nb Clients"]  = disp["Nb Clients"].apply(lambda v: f"{int(v)}")
            disp["Taux Recouvrement"] = disp["Taux Recouvrement"].apply(lambda v: f"{v:.1f}%")
            disp.index.name = "Segment"
            st.dataframe(disp.reset_index(), use_container_width=True, height=340)

        st.markdown("<br>", unsafe_allow_html=True)

        # ── GRAPHIQUE 3 : Évolution CA mensuelle — une courbe par segment
        section("📈  ÉVOLUTION MENSUELLE — Une courbe par segment")
        comment("Chaque courbe = un segment. Vous pouvez voir comment chaque segment évolue dans le temps. "
                "Cliquez sur un segment dans la légende pour l'afficher/masquer et comparer plus facilement.")
        ca_mois_seg = df_comp.groupby(["Période","SEGMENT"])["MONTANT TTC"].sum().reset_index()
        fig3 = go.Figure()
        for i, seg in enumerate(comp_segments):
            d = ca_mois_seg[ca_mois_seg["SEGMENT"]==seg]
            c = seg_colors_list[i % len(seg_colors_list)]
            fig3.add_trace(go.Scatter(x=d["Période"], y=d["MONTANT TTC"],
                name=seg, mode="lines+markers",
                line=dict(color=c, width=2.5),
                marker=dict(size=7, color=c, line=dict(color="white",width=1.5)),
                hovertemplate=f"<b>{seg}</b><br>%{{x}}<br>CA : %{{y:,.0f}} FCFA<extra></extra>"))
        fig3 = excel_style(fig3, 380)
        fig3.update_xaxes(tickangle=-35, tickfont=dict(size=8))
        fig3.update_yaxes(tickformat=",.0f", ticksuffix=" F")
        st.plotly_chart(fig3, use_container_width=True)

        st.markdown("<br>", unsafe_allow_html=True)

        # ── GRAPHIQUE 4 : Radar / Spider chart multi-indicateurs
        section("🕸️  PROFIL RADAR — Comparaison multi-indicateurs simultanée")
        comment("Ce graphique en toile d'araignée permet de comparer plusieurs indicateurs en même temps. "
                "Plus la surface couverte est grande = meilleure performance globale du segment. "
                "Idéal pour voir les forces et faiblesses de chaque segment d'un seul regard.")

        # Normaliser pour radar
        radar_cols = ["CA Total","CA Payé","Nb Factures","Panier Moyen","Nb Clients","Taux Recouvrement"]
        radar_labels = ["CA Total","CA Payé","Nb Factures","Panier Moyen","Nb Clients","Recouvrement %"]
        radar_data = metrics_df[radar_cols].copy()
        radar_norm = radar_data.copy()
        for col in radar_cols:
            max_v = radar_data[col].max()
            if max_v > 0:
                radar_norm[col] = radar_data[col] / max_v * 100

        fig4 = go.Figure()
        for i, seg in enumerate(comp_segments[:8]):  # Max 8 sur radar
            if seg not in radar_norm.index: continue
            vals = radar_norm.loc[seg, radar_cols].tolist()
            raw  = radar_data.loc[seg, radar_cols].tolist()
            vals += [vals[0]]  # fermer le polygone
            labels_ext = radar_labels + [radar_labels[0]]
            c = seg_colors_list[i % len(seg_colors_list)]
            fig4.add_trace(go.Scatterpolar(
                r=vals, theta=labels_ext, name=seg,
                fill="toself", fillcolor=f"rgba({int(c[1:3],16)},{int(c[3:5],16)},{int(c[5:7],16)},0.15)",
                line=dict(color=c, width=2),
                marker=dict(size=5, color=c),
                hovertemplate=(f"<b>{seg}</b><br>" +
                    "<br>".join([f"{radar_labels[j]}: {raw[j]:,.0f}" for j in range(len(radar_cols))]) +
                    "<extra></extra>"),
            ))
        fig4.update_layout(
            polar=dict(
                bgcolor="#FAFAFA",
                radialaxis=dict(visible=True, range=[0,100], ticksuffix="%",
                                gridcolor="#E2E8F0", tickfont=dict(size=9, color=C["muted"])),
                angularaxis=dict(gridcolor="#E2E8F0", tickfont=dict(size=10, color=C["texte"])),
            ),
            paper_bgcolor="#FFFFFF", height=450,
            legend=dict(bgcolor="rgba(0,0,0,0)", font=dict(color=C["texte"],size=10),
                        orientation="v", x=1.08, y=0.5),
            margin=dict(t=30,b=30,l=80,r=160),
        )
        st.plotly_chart(fig4, use_container_width=True)

        st.markdown("<br>", unsafe_allow_html=True)

        # ── GRAPHIQUE 5 : Payé vs Impayé par segment — barres côte à côte
        section("💳  PAYÉ VS IMPAYÉ PAR SEGMENT — Qui paie bien, qui paie mal ?")
        comment("Ce graphique sépare pour chaque segment ce qui a été payé (bleu) et ce qui reste impayé (rose). "
                "Les segments avec une grande barre rose ont un risque de trésorerie plus élevé. "
                "Le taux de recouvrement (% payé) est affiché pour chaque segment.")
        fig5 = go.Figure()
        segs_ord = metrics_df.sort_values("CA Total",ascending=False).index.tolist()
        fig5.add_trace(go.Bar(
            x=segs_ord,
            y=[metrics_df.loc[s,"CA Payé"] if s in metrics_df.index else 0 for s in segs_ord],
            name="Payé ✅", marker_color=C["bleu"],
            hovertemplate="<b>%{x}</b><br>CA Payé : %{y:,.0f} FCFA<extra></extra>"))
        fig5.add_trace(go.Bar(
            x=segs_ord,
            y=[metrics_df.loc[s,"CA Impayé"] if s in metrics_df.index else 0 for s in segs_ord],
            name="Impayé ⏳", marker_color=C["rose"],
            hovertemplate="<b>%{x}</b><br>CA Impayé : %{y:,.0f} FCFA<extra></extra>"))
        # Annotations taux recouvrement
        for s in segs_ord:
            if s in metrics_df.index:
                taux_v = metrics_df.loc[s,"Taux Recouvrement"]
                fig5.add_annotation(x=s, y=metrics_df.loc[s,"CA Total"]*1.05,
                    text=f"{taux_v:.0f}%", showarrow=False,
                    font=dict(size=10, color=C["texte"], family="Inter"),
                    bgcolor="rgba(255,255,255,0.8)",
                    bordercolor=C["bleu"] if taux_v>=80 else C["rouge"], borderwidth=1,
                    borderpad=3)
        fig5.update_layout(barmode="stack")
        fig5 = excel_style(fig5, 380)
        fig5.update_xaxes(tickangle=-20, tickfont=dict(size=9))
        fig5.update_yaxes(tickformat=",.0f", ticksuffix=" F")
        st.plotly_chart(fig5, use_container_width=True)

        st.markdown("<br>", unsafe_allow_html=True)

        # ── GRAPHIQUE 6 : Comparaison annuelle par segment (barres groupées)
        section("📅  ÉVOLUTION ANNUELLE PAR SEGMENT — Qui progresse ?")
        comment("Chaque groupe de barres = une année. À l'intérieur du groupe, chaque couleur = un segment. "
                "Si une barre grandit d'une année à l'autre = le segment est en croissance.")
        ca_ann_seg = df_comp.groupby(["ANNEE","SEGMENT"])["MONTANT TTC"].sum().reset_index()
        fig6 = go.Figure()
        for i, seg in enumerate(comp_segments):
            d = ca_ann_seg[ca_ann_seg["SEGMENT"]==seg]
            c = seg_colors_list[i % len(seg_colors_list)]
            fig6.add_trace(go.Bar(
                x=d["ANNEE"].astype(str), y=d["MONTANT TTC"],
                name=seg, marker_color=c,
                hovertemplate=f"<b>{seg}</b><br>%{{x}}<br>CA : %{{y:,.0f}} FCFA<extra></extra>"))
        fig6.update_layout(barmode="group")
        fig6 = excel_style(fig6, 360)
        fig6.update_xaxes(tickfont=dict(size=11))
        fig6.update_yaxes(tickformat=",.0f", ticksuffix=" F")
        st.plotly_chart(fig6, use_container_width=True)

    # ═══════════════════════════════════════
    # MODE 2 : COMPARAISON 2 À 2
    # ═══════════════════════════════════════
    else:
        st.markdown("<br>", unsafe_allow_html=True)
        comment("Choisissez <b>deux segments</b> à comparer en détail. Vous verrez leur performance "
                "côte à côte sur tous les indicateurs clés. Changez la paire à tout moment.",
                "🔀 Comparaison 2 à 2")

        col_s1, col_s2 = st.columns(2)
        with col_s1:
            seg_A = st.selectbox("Segment A", comp_segments, index=0)
        with col_s2:
            seg_B = st.selectbox("Segment B", comp_segments,
                                  index=min(1, len(comp_segments)-1))

        if seg_A == seg_B:
            st.warning("Choisissez deux segments différents.")
            st.stop()

        df_A = df_comp[df_comp["SEGMENT"]==seg_A]
        df_B = df_comp[df_comp["SEGMENT"]==seg_B]

        def seg_kpis(df, seg, color):
            ca   = df["MONTANT TTC"].sum()
            paye = df[df["ETAT DE PAIEMENT"]=="Payée"]["MONTANT TTC"].sum()
            imp  = df[df["ETAT DE PAIEMENT"]=="Impayée"]["MONTANT TTC"].sum()
            taux = paye/ca*100 if ca else 0
            nb_f = len(df); nb_cl = df["CLIENT"].nunique()
            panier = df["MONTANT TTC"].mean() if len(df) else 0
            return {"CA":ca,"Payé":paye,"Impayé":imp,"Taux":taux,
                    "Factures":nb_f,"Clients":nb_cl,"Panier":panier,"color":color,"seg":seg}

        kA = seg_kpis(df_A, seg_A, C["bleu"])
        kB = seg_kpis(df_B, seg_B, C["rose"])

        # ── Header comparatif
        section(f"⚡  {seg_A}  vs  {seg_B}")

        # KPI visuels côte à côte
        for label, keyA, keyB, fmtfn in [
            ("CA Total",   kA["CA"],      kB["CA"],      lambda v: fmt(v,"").strip()),
            ("CA Payé",    kA["Payé"],    kB["Payé"],    lambda v: fmt(v,"").strip()),
            ("CA Impayé",  kA["Impayé"],  kB["Impayé"],  lambda v: fmt(v,"").strip()),
            ("Taux Recouv.",kA["Taux"],   kB["Taux"],    lambda v: f"{v:.1f}%"),
            ("Nb Factures",kA["Factures"],kB["Factures"],lambda v: f"{int(v)}"),
            ("Nb Clients", kA["Clients"], kB["Clients"], lambda v: f"{int(v)}"),
        ]:
            winner_A = keyA >= keyB if label != "CA Impayé" else keyA <= keyB
            col_a_disp, col_mid, col_b_disp = st.columns([2, 1, 2])
            with col_a_disp:
                border = C["vert"] if winner_A else C["muted"]
                st.markdown(f"""
                <div style="background:#FFFFFF;border:2px solid {border};border-radius:10px;
                            padding:12px 16px;text-align:right;">
                    <div style="font-size:0.68rem;color:{C['muted']};text-transform:uppercase;font-weight:700;">{seg_A[:25]}</div>
                    <div style="font-size:1.35rem;font-weight:800;color:{C['bleu']};margin-top:3px;">{fmtfn(keyA)}</div>
                    {"<div style='font-size:0.75rem;color:"+C['vert']+"'>🏆 En tête</div>" if winner_A else ""}
                </div>""", unsafe_allow_html=True)
            with col_mid:
                st.markdown(f"""
                <div style="text-align:center;padding-top:18px;">
                    <div style="font-size:0.68rem;color:{C['muted']};font-weight:700;">{label}</div>
                    <div style="font-size:1rem;color:{C['texte']};font-weight:600;">VS</div>
                </div>""", unsafe_allow_html=True)
            with col_b_disp:
                border = C["vert"] if not winner_A else C["muted"]
                st.markdown(f"""
                <div style="background:#FFFFFF;border:2px solid {border};border-radius:10px;
                            padding:12px 16px;text-align:left;">
                    <div style="font-size:0.68rem;color:{C['muted']};text-transform:uppercase;font-weight:700;">{seg_B[:25]}</div>
                    <div style="font-size:1.35rem;font-weight:800;color:{C['rose']};margin-top:3px;">{fmtfn(keyB)}</div>
                    {"<div style='font-size:0.75rem;color:"+C['vert']+"'>🏆 En tête</div>" if not winner_A else ""}
                </div>""", unsafe_allow_html=True)
            st.markdown("<div style='margin:4px 0;'></div>", unsafe_allow_html=True)

        st.markdown("<br>", unsafe_allow_html=True)

        # ── GRAPHIQUE 1 : Courbes évolution CA mensuelle côte à côte
        section("📈  ÉVOLUTION MENSUELLE — Courbes séparées")
        comment(f"Comparez mois par mois les ventes de <b>{seg_A}</b> (bleu) et <b>{seg_B}</b> (rose). "
                "Si une courbe monte plus vite = ce segment est en forte croissance.")
        ca_m_AB = df_comp[df_comp["SEGMENT"].isin([seg_A,seg_B])].groupby(["Période","SEGMENT"])["MONTANT TTC"].sum().reset_index()
        fig_ev = go.Figure()
        for seg, c, dash in [(seg_A, C["bleu"],"solid"), (seg_B, C["rose"],"dash")]:
            d = ca_m_AB[ca_m_AB["SEGMENT"]==seg]
            fig_ev.add_trace(go.Scatter(x=d["Période"], y=d["MONTANT TTC"],
                name=seg, mode="lines+markers",
                line=dict(color=c, width=3, dash=dash),
                marker=dict(size=8, color=c, line=dict(color="white",width=2)),
                fill="tozeroy",
                fillcolor=f"rgba({int(c[1:3],16)},{int(c[3:5],16)},{int(c[5:7],16)},0.08)",
                hovertemplate=f"<b>{seg}</b><br>%{{x}}<br>CA : %{{y:,.0f}} FCFA<extra></extra>"))
        fig_ev = excel_style(fig_ev, 340)
        fig_ev.update_xaxes(tickangle=-35, tickfont=dict(size=8))
        fig_ev.update_yaxes(tickformat=",.0f", ticksuffix=" F")
        st.plotly_chart(fig_ev, use_container_width=True)

        st.markdown("<br>", unsafe_allow_html=True)

        col_g1, col_g2 = st.columns(2)
        with col_g1:
            # ── GRAPHIQUE 2 : Barres groupées annuelles
            section("📅  CA PAR ANNÉE")
            comment(f"Comparez l'évolution annuelle de <b>{seg_A}</b> vs <b>{seg_B}</b>. "
                    "Qui a progressé le plus ? Qui a reculé ?")
            ca_ann_AB = df_comp[df_comp["SEGMENT"].isin([seg_A,seg_B])].groupby(["ANNEE","SEGMENT"])["MONTANT TTC"].sum().reset_index()
            fig_an = go.Figure()
            for seg, c in [(seg_A, C["bleu"]), (seg_B, C["rose"])]:
                d = ca_ann_AB[ca_ann_AB["SEGMENT"]==seg]
                fig_an.add_trace(go.Bar(x=d["ANNEE"].astype(str), y=d["MONTANT TTC"],
                    name=seg, marker_color=c,
                    hovertemplate=f"<b>{seg}</b><br>%{{x}}<br>%{{y:,.0f}} FCFA<extra></extra>"))
            fig_an.update_layout(barmode="group")
            fig_an = excel_style(fig_an, 300)
            fig_an.update_yaxes(tickformat=",.0f", ticksuffix=" F")
            st.plotly_chart(fig_an, use_container_width=True)

        with col_g2:
            # ── GRAPHIQUE 3 : Payé vs Impayé pour A et B
            section("💳  PAYÉ VS IMPAYÉ")
            comment(f"Qui a le meilleur taux de paiement ? "
                    f"{seg_A} : {kA['Taux']:.0f}% · {seg_B} : {kB['Taux']:.0f}%")
            fig_pay = go.Figure()
            for seg, kk, c_p, c_i in [(seg_A, kA, C["bleu"], "#C5CBFB"), (seg_B, kB, C["rose"], "#FFB7D4")]:
                fig_pay.add_trace(go.Bar(x=[seg], y=[kk["Payé"]], name=f"Payé – {seg}",
                    marker_color=c_p, hovertemplate=f"Payé : %{{y:,.0f}} FCFA<extra></extra>"))
                fig_pay.add_trace(go.Bar(x=[seg], y=[kk["Impayé"]], name=f"Impayé – {seg}",
                    marker_color=c_i, hovertemplate=f"Impayé : %{{y:,.0f}} FCFA<extra></extra>"))
            fig_pay.update_layout(barmode="stack")
            fig_pay = excel_style(fig_pay, 300)
            fig_pay.update_yaxes(tickformat=",.0f", ticksuffix=" F")
            st.plotly_chart(fig_pay, use_container_width=True)

        st.markdown("<br>", unsafe_allow_html=True)

        # ── GRAPHIQUE 4 : Top clients de chaque segment
        col_g3, col_g4 = st.columns(2)
        with col_g3:
            section(f"🏆  TOP CLIENTS — {seg_A[:20]}")
            top_A = df_A.groupby("CLIENT")["MONTANT TTC"].sum().sort_values(ascending=True).tail(6).reset_index()
            fig_cA = go.Figure(go.Bar(x=top_A["MONTANT TTC"], y=top_A["CLIENT"],
                orientation="h", marker_color=C["bleu"],
                text=top_A["MONTANT TTC"].apply(lambda v: fmt(v,"").strip()),
                textposition="outside", textfont=dict(size=10),
                hovertemplate="<b>%{y}</b><br>%{x:,.0f} FCFA<extra></extra>"))
            fig_cA = excel_style(fig_cA, 300, False)
            fig_cA.update_xaxes(tickformat=",.0f")
            st.plotly_chart(fig_cA, use_container_width=True)

        with col_g4:
            section(f"🏆  TOP CLIENTS — {seg_B[:20]}")
            top_B = df_B.groupby("CLIENT")["MONTANT TTC"].sum().sort_values(ascending=True).tail(6).reset_index()
            fig_cB = go.Figure(go.Bar(x=top_B["MONTANT TTC"], y=top_B["CLIENT"],
                orientation="h", marker_color=C["rose"],
                text=top_B["MONTANT TTC"].apply(lambda v: fmt(v,"").strip()),
                textposition="outside", textfont=dict(size=10),
                hovertemplate="<b>%{y}</b><br>%{x:,.0f} FCFA<extra></extra>"))
            fig_cB = excel_style(fig_cB, 300, False)
            fig_cB.update_xaxes(tickformat=",.0f")
            st.plotly_chart(fig_cB, use_container_width=True)

        st.markdown("<br>", unsafe_allow_html=True)

        # ── Synthèse narrative automatique
        section("📝  SYNTHÈSE AUTOMATIQUE DE LA COMPARAISON")
        winner_ca  = seg_A if kA["CA"] >= kB["CA"] else seg_B
        winner_rec = seg_A if kA["Taux"] >= kB["Taux"] else seg_B
        loser_rec  = seg_B if kA["Taux"] >= kB["Taux"] else seg_A
        diff_ca    = abs(kA["CA"]-kB["CA"])
        diff_pct   = diff_ca/max(kA["CA"],kB["CA"])*100 if max(kA["CA"],kB["CA"])>0 else 0
        st.markdown(f"""
        <div class="comment-box">
            <div class="ct">📝 Analyse automatique : {seg_A} vs {seg_B}</div>
            <ul style="margin:6px 0;padding-left:18px;">
                <li><b>CA :</b> <b>{winner_ca}</b> est en tête avec {fmt(max(kA['CA'],kB['CA']),'').strip()} FCFA,
                soit <b>{diff_pct:.0f}%</b> de plus que l'autre segment ({fmt(diff_ca,'').strip()} FCFA d'écart).</li>
                <li><b>Recouvrement :</b> <b>{winner_rec}</b> paie mieux ({max(kA['Taux'],kB['Taux']):.0f}% payé).
                <b>{loser_rec}</b> présente un risque plus élevé ({min(kA['Taux'],kB['Taux']):.0f}% payé).</li>
                <li><b>Activité :</b> {seg_A} a généré <b>{kA['Factures']} factures</b> via <b>{kA['Clients']} clients</b>
                — {seg_B} a généré <b>{kB['Factures']} factures</b> via <b>{kB['Clients']} clients</b>.</li>
                <li><b>Panier moyen :</b> {seg_A} : {fmt(kA['Panier'],'').strip()} FCFA/commande
                · {seg_B} : {fmt(kB['Panier'],'').strip()} FCFA/commande.</li>
            </ul>
            <b>Recommandation :</b> Concentrez les efforts commerciaux sur <b>{winner_ca}</b> pour maximiser le CA,
            et renforcez le suivi des paiements chez <b>{loser_rec}</b> pour sécuriser la trésorerie.
        </div>""", unsafe_allow_html=True)


# PAGE 9 : ALERTES & CONSEILS
# ══════════════════════════════════════════════════════════════


# ══════════════════════════════════════════════════════════════
# PAGE CONTRÔLE DE GESTION — Module complet
# ══════════════════════════════════════════════════════════════
elif page == "🧮  Contrôle de Gestion":


    # ── Paramètres financiers MULTIPACK (hypothèses réalistes)
    # Ces paramètres peuvent être ajustés dans la sidebar par le contrôleur
    PARAMS = {
        # Charges fixes mensuelles (FCFA)
        "masse_salariale":    3_500_000,
        "loyer_usine":          800_000,
        "energie_electricite":  450_000,
        "maintenance_machines": 280_000,
        "assurances":           150_000,
        "frais_generaux":       320_000,
        "amortissements":       600_000,
        # Charges variables (% du CA HT)
        "transport_logistique": 0.035,
        "emballages_fournitures":0.018,
        "commissions_commerciales":0.025,
        # Taux TVA
        "taux_tva": 0.18,
        # Objectif CA mensuel (FCFA TTC)
        "objectif_ca_mensuel": None,  # calculé dynamiquement
        # Coût matières premières (% CA HT estimé)
        "taux_matieres":  0.42,
    }

    # ── Charges fixes totales mensuelles
    charges_fixes_mois = sum([
        PARAMS["masse_salariale"], PARAMS["loyer_usine"],
        PARAMS["energie_electricite"], PARAMS["maintenance_machines"],
        PARAMS["assurances"], PARAMS["frais_generaux"], PARAMS["amortissements"],
    ])

    # ── Données de base (toujours sur df_fact complet pour cohérence CG)
    df_cg = df_f.copy()  # Respecte les filtres sidebar (années, segments, zones)
    df_cg["Période"]    = df_cg["DATE"].dt.to_period("M").astype(str)
    df_cg["ANNEE"]      = df_cg["DATE"].dt.year
    df_cg["MOIS"]       = df_cg["DATE"].dt.month
    df_cg["CA_HT"]      = df_cg["MONTANT HT"] if "MONTANT HT" in df_cg.columns else df_cg["MONTANT TTC"] / 1.18
    df_cg["Matieres"]   = df_cg["CA_HT"] * PARAMS["taux_matieres"]
    df_cg["Marge_Brute"]= df_cg["CA_HT"] - df_cg["Matieres"]
    df_cg["Taux_MB"]    = df_cg["Marge_Brute"] / df_cg["CA_HT"] * 100

    # ── Agrégation mensuelle P&L
    pl_mois = df_cg.groupby("Période").agg(
        CA_TTC=("MONTANT TTC","sum"),
        CA_HT=("CA_HT","sum"),
        Matieres=("Matieres","sum"),
        Marge_Brute=("Marge_Brute","sum"),
        Nb_Factures=("MONTANT TTC","count"),
    ).reset_index().sort_values("Période")

    pl_mois["Charges_Var"] = pl_mois["CA_HT"] * (
        PARAMS["transport_logistique"] + PARAMS["emballages_fournitures"] +
        PARAMS["commissions_commerciales"])
    pl_mois["Marge_Semi_Nette"] = pl_mois["Marge_Brute"] - pl_mois["Charges_Var"]
    pl_mois["Charges_Fixes"]    = charges_fixes_mois
    pl_mois["EBITDA"]           = pl_mois["Marge_Semi_Nette"] - pl_mois["Charges_Fixes"]
    pl_mois["Amortissements"]   = PARAMS["amortissements"]
    pl_mois["EBIT"]             = pl_mois["EBITDA"] - pl_mois["Amortissements"]
    pl_mois["Taux_EBITDA"]      = pl_mois["EBITDA"] / pl_mois["CA_HT"] * 100
    pl_mois["Taux_MB_Brute"]    = pl_mois["Marge_Brute"] / pl_mois["CA_HT"] * 100

    # ── Objectif CA (moyenne + 10%)
    ca_moy_hist = pl_mois["CA_TTC"].mean()
    objectif_ca = ca_moy_hist * 1.10
    pl_mois["Objectif_CA"]  = objectif_ca
    pl_mois["Ecart_CA"]     = pl_mois["CA_TTC"] - objectif_ca
    pl_mois["Ecart_CA_Pct"] = pl_mois["Ecart_CA"] / objectif_ca * 100

    # ── Point mort mensuel
    taux_mv = 1 - PARAMS["taux_matieres"] - PARAMS["transport_logistique"] - \
              PARAMS["emballages_fournitures"] - PARAMS["commissions_commerciales"]
    point_mort_ht = charges_fixes_mois / taux_mv if taux_mv > 0 else 0
    point_mort_ttc = point_mort_ht * (1 + PARAMS["taux_tva"])

    # ── DSO (Days Sales Outstanding) — délai moyen de paiement
    df_paye = df_fact[df_fact["ETAT DE PAIEMENT"] == "Payée"]
    ca_jour_moy = df_fact["MONTANT TTC"].sum() / 365 if len(df_fact) > 0 else 1
    encours_clients = df_fact[df_fact["ETAT DE PAIEMENT"] == "Impayée"]["MONTANT TTC"].sum()
    dso = encours_clients / ca_jour_moy if ca_jour_moy > 0 else 0

    # ── BFR estimé
    stock_val   = df_inv["Valeur"].sum()
    creances    = encours_clients
    dettes_fourn= pl_mois["Matieres"].mean() * 1.5  # estimation 45j de délai fournisseur
    bfr         = stock_val + creances - dettes_fourn

    # ═══════════════════════════════════════════════
    # HEADER PAGE CG
    # ═══════════════════════════════════════════════
    st.markdown(f"""
    <div style="background:linear-gradient(135deg,#1E2A4A,#2D3A6B,{C['or']});
                border-radius:14px;padding:22px 28px;margin-bottom:18px;
                border:1px solid rgba(214,158,46,0.3);">
        <div style="font-size:1.15rem;font-weight:800;color:white;margin-bottom:4px;">
            🧮 Contrôle de Gestion — MULTIPACK SA
        </div>
        <div style="font-size:0.8rem;color:rgba(255,255,255,0.65);">
            Compte de résultat · Budget vs Réalisé · BFR · DSO · Point mort · Analyse des écarts
        </div>
    </div>""", unsafe_allow_html=True)

    # Panneau paramètres ajustables
    with st.expander("⚙️ Paramètres financiers — Cliquez pour ajuster les hypothèses"):
        st.markdown(f"<div style='font-size:0.8rem;color:{C['muted']};margin-bottom:12px;'>"
                    "Ces paramètres alimentent tous les calculs du module. "
                    "Ajustez-les selon les données réelles de MULTIPACK.</div>",
                    unsafe_allow_html=True)
        pc1, pc2, pc3 = st.columns(3)
        with pc1:
            st.markdown("**Charges fixes mensuelles (FCFA)**")
            ms  = st.number_input("Masse salariale",    value=PARAMS["masse_salariale"],    step=100_000, format="%d")
            lo  = st.number_input("Loyer usine",        value=PARAMS["loyer_usine"],        step=50_000,  format="%d")
            en  = st.number_input("Énergie/électricité",value=PARAMS["energie_electricite"],step=50_000,  format="%d")
            ma  = st.number_input("Maintenance",        value=PARAMS["maintenance_machines"],step=50_000,  format="%d")
        with pc2:
            st.markdown("**Autres charges fixes (FCFA)**")
            as_ = st.number_input("Assurances",         value=PARAMS["assurances"],         step=20_000,  format="%d")
            fg  = st.number_input("Frais généraux",     value=PARAMS["frais_generaux"],     step=50_000,  format="%d")
            am  = st.number_input("Amortissements",     value=PARAMS["amortissements"],     step=50_000,  format="%d")
        with pc3:
            st.markdown("**Taux variables (% CA HT)**")
            tm  = st.slider("Coût matières premières",      0.30, 0.70, PARAMS["taux_matieres"],        0.01)
            tt  = st.slider("Transport & logistique",        0.01, 0.10, PARAMS["transport_logistique"], 0.005)
            ef  = st.slider("Emballages & fournitures",      0.01, 0.05, PARAMS["emballages_fournitures"],0.002)
            cc  = st.slider("Commissions commerciales",      0.01, 0.08, PARAMS["commissions_commerciales"],0.005)

        # Recalculer avec paramètres ajustés
        charges_fixes_mois = ms + lo + en + ma + as_ + fg + am
        PARAMS.update({"taux_matieres":tm,"transport_logistique":tt,
                       "emballages_fournitures":ef,"commissions_commerciales":cc,
                       "amortissements":am})
        df_cg["Matieres"]    = df_cg["CA_HT"] * tm
        df_cg["Marge_Brute"] = df_cg["CA_HT"] - df_cg["Matieres"]
        pl_mois["Matieres"]       = pl_mois["CA_HT"] * tm
        pl_mois["Marge_Brute"]    = pl_mois["CA_HT"] - pl_mois["Matieres"]
        pl_mois["Charges_Var"]    = pl_mois["CA_HT"] * (tt + ef + cc)
        pl_mois["Marge_Semi_Nette"]= pl_mois["Marge_Brute"] - pl_mois["Charges_Var"]
        pl_mois["Charges_Fixes"]  = charges_fixes_mois
        pl_mois["EBITDA"]         = pl_mois["Marge_Semi_Nette"] - pl_mois["Charges_Fixes"]
        pl_mois["EBIT"]           = pl_mois["EBITDA"] - am
        pl_mois["Taux_EBITDA"]    = pl_mois["EBITDA"] / pl_mois["CA_HT"] * 100
        pl_mois["Taux_MB_Brute"]  = pl_mois["Marge_Brute"] / pl_mois["CA_HT"] * 100
        taux_mv = 1 - tm - tt - ef - cc
        point_mort_ht  = charges_fixes_mois / taux_mv if taux_mv > 0 else 0
        point_mort_ttc = point_mort_ht * 1.18
        dettes_fourn   = pl_mois["Matieres"].mean() * 1.5
        bfr            = stock_val + creances - dettes_fourn

    # ── ONGLETS DU MODULE CG ─────────────────────────────────
    cg1, cg2, cg3, cg4, cg5, cg6 = st.tabs([
        "📋 Compte de Résultat",
        "🎯 Budget vs Réalisé",
        "💹 Marges & Rentabilité",
        "🏦 BFR & Trésorerie",
        "⚖️ Point Mort",
        "📐 Analyse des Écarts",
    ])

    # ════════════════════════════════════════════
    # ONGLET 1 — COMPTE DE RÉSULTAT (P&L)
    # ════════════════════════════════════════════
    with cg1:
        st.markdown("<br>", unsafe_allow_html=True)
        comment(
            "Le compte de résultat montre comment MULTIPACK passe du <b>chiffre d'affaires</b> "
            "au <b>résultat net</b>, étape par étape. "
            "Chaque ligne défalque une catégorie de coûts. "
            "Un EBITDA positif = l'activité est rentable avant financement et investissements.",
            "📋 Comment lire le compte de résultat ?"
        )

        # KPIs P&L globaux (totaux sur toute la période)
        ca_ht_tot  = pl_mois["CA_HT"].sum()
        mat_tot    = pl_mois["Matieres"].sum()
        mb_tot     = pl_mois["Marge_Brute"].sum()
        cv_tot     = pl_mois["Charges_Var"].sum()
        msn_tot    = pl_mois["Marge_Semi_Nette"].sum()
        cf_tot     = pl_mois["Charges_Fixes"].sum()
        ebitda_tot = pl_mois["EBITDA"].sum()
        ebit_tot   = pl_mois["EBIT"].sum()

        # Tableau P&L visuel
        section("📋  COMPTE DE RÉSULTAT CONSOLIDÉ — Toute la période")

        def pl_row(label, value, indent=0, bold=False, color=None, separator=False):
            pad = indent * 20
            b_open  = "<b>" if bold else ""
            b_close = "</b>" if bold else ""
            if color is None:
                color = C["vert"] if value >= 0 else C["rouge"]
            bg = "rgba(97,114,243,0.06)" if bold else "transparent"
            border_top = f"border-top:2px solid {C['bordure']};" if separator else ""
            st.markdown(f"""
            <div style="display:flex;justify-content:space-between;align-items:center;
                        padding:9px 16px;border-bottom:1px solid {C['bordure']};
                        background:{bg};{border_top}margin:0;">
                <span style="font-size:0.84rem;color:{C['texte']};padding-left:{pad}px;">
                    {b_open}{label}{b_close}
                </span>
                <span style="font-size:0.9rem;font-weight:{'800' if bold else '500'};color:{color};">
                    {b_open}{fmt(value, '').strip()} FCFA{b_close}
                </span>
            </div>""", unsafe_allow_html=True)

        def pl_pct_row(label, value, pct, indent=0, bold=False):
            pad = indent * 20
            b_open  = "<b>" if bold else ""
            b_close = "</b>" if bold else ""
            color_v = C["vert"] if value >= 0 else C["rouge"]
            color_p = C["vert"] if pct >= 0 else C["rouge"]
            bg = "rgba(97,114,243,0.06)" if bold else "transparent"
            st.markdown(f"""
            <div style="display:flex;justify-content:space-between;align-items:center;
                        padding:9px 16px;border-bottom:1px solid {C['bordure']};background:{bg};">
                <span style="font-size:0.84rem;color:{C['texte']};padding-left:{pad}px;">
                    {b_open}{label}{b_close}
                </span>
                <div style="display:flex;gap:32px;align-items:center;">
                    <span style="font-size:0.78rem;color:{color_p};">{pct:+.1f}% du CA HT</span>
                    <span style="font-size:0.9rem;font-weight:{'800' if bold else '500'};color:{color_v};">
                        {b_open}{fmt(value,'').strip()} FCFA{b_close}
                    </span>
                </div>
            </div>""", unsafe_allow_html=True)

        st.markdown(f"""
        <div style="background:{C['blanc']};border:1px solid {C['bordure']};border-radius:12px;
                    overflow:hidden;box-shadow:0 2px 8px rgba(0,0,0,0.06);">""",
                    unsafe_allow_html=True)
        pl_row("📦 Chiffre d'Affaires TTC",  pl_mois["CA_TTC"].sum(), bold=True, color=C["bleu"])
        pl_row("  − TVA collectée",           pl_mois["CA_TTC"].sum() - ca_ht_tot, indent=1, color=C["muted"])
        pl_row("= Chiffre d'Affaires HT",     ca_ht_tot, bold=True, color=C["bleu"], separator=True)
        pl_pct_row("  − Coût des Matières Premières", -mat_tot,
                   -mat_tot/ca_ht_tot*100 if ca_ht_tot else 0, indent=1)
        pl_pct_row("= MARGE BRUTE",            mb_tot,
                   mb_tot/ca_ht_tot*100 if ca_ht_tot else 0, bold=True, indent=0)
        pl_pct_row("  − Transport & Logistique",
                   -pl_mois["CA_HT"].sum()*PARAMS["transport_logistique"],
                   -PARAMS["transport_logistique"]*100, indent=1)
        pl_pct_row("  − Emballages & Fournitures",
                   -pl_mois["CA_HT"].sum()*PARAMS["emballages_fournitures"],
                   -PARAMS["emballages_fournitures"]*100, indent=1)
        pl_pct_row("  − Commissions Commerciales",
                   -pl_mois["CA_HT"].sum()*PARAMS["commissions_commerciales"],
                   -PARAMS["commissions_commerciales"]*100, indent=1)
        pl_pct_row("= MARGE SUR COÛTS VARIABLES", msn_tot,
                   msn_tot/ca_ht_tot*100 if ca_ht_tot else 0, bold=True, indent=0)
        pl_pct_row("  − Masse Salariale",      -PARAMS['masse_salariale'] * len(pl_mois), -PARAMS['masse_salariale']/ca_ht_tot*len(pl_mois)*100 if ca_ht_tot else 0, indent=1)
        pl_pct_row("  − Loyer Usine",          -PARAMS['loyer_usine'] * len(pl_mois), -PARAMS['loyer_usine']/ca_ht_tot*len(pl_mois)*100 if ca_ht_tot else 0, indent=1)
        pl_pct_row("  − Énergie & Électricité",-PARAMS['energie_electricite'] * len(pl_mois), -PARAMS['energie_electricite']/ca_ht_tot*len(pl_mois)*100 if ca_ht_tot else 0, indent=1)
        pl_pct_row("  − Maintenance Machines", -PARAMS['maintenance_machines'] * len(pl_mois), -PARAMS['maintenance_machines']/ca_ht_tot*len(pl_mois)*100 if ca_ht_tot else 0, indent=1)
        pl_pct_row("  − Frais Généraux",       -(PARAMS['assurances']+PARAMS['frais_generaux']) * len(pl_mois), -(PARAMS['assurances']+PARAMS['frais_generaux'])/ca_ht_tot*len(pl_mois)*100 if ca_ht_tot else 0, indent=1)
        pl_pct_row("= EBITDA",                 ebitda_tot,
                   ebitda_tot/ca_ht_tot*100 if ca_ht_tot else 0, bold=True, indent=0)
        pl_pct_row("  − Amortissements",       -am * len(pl_mois), -am/ca_ht_tot*len(pl_mois)*100 if ca_ht_tot else 0, indent=1)
        pl_pct_row("= RÉSULTAT D'EXPLOITATION (EBIT)", ebit_tot,
                   ebit_tot/ca_ht_tot*100 if ca_ht_tot else 0, bold=True, indent=0)
        st.markdown("</div>", unsafe_allow_html=True)

        st.markdown("<br>", unsafe_allow_html=True)

        # KPIs synthèse
        k1,k2,k3,k4,k5 = st.columns(5)
        with k1: kpi("💰","CA HT Total", fmt(ca_ht_tot), "Hors taxes","bleu")
        with k2: kpi("📊","Marge Brute", fmt(mb_tot),
                     f"{mb_tot/ca_ht_tot*100:.1f}% du CA HT","vert" if mb_tot>0 else "rouge")
        with k3: kpi("⚙️","EBITDA", fmt(ebitda_tot),
                     f"{ebitda_tot/ca_ht_tot*100:.1f}% du CA HT","vert" if ebitda_tot>0 else "rouge")
        with k4: kpi("📈","EBIT", fmt(ebit_tot),
                     f"{ebit_tot/ca_ht_tot*100:.1f}% du CA HT","vert" if ebit_tot>0 else "rouge")
        with k5: kpi("🏭","Charges Fixes Tot.", fmt(cf_tot),
                     f"{cf_tot/ca_ht_tot*100:.1f}% du CA HT","orange")

        st.markdown("<br>", unsafe_allow_html=True)

        # GRAPHIQUE — Cascade P&L mensuel (waterfall)
        section("📊  ÉVOLUTION MENSUELLE DU P&L — Du CA à l'EBIT")
        comment("Chaque courbe montre l'évolution d'un indicateur clé du P&L. "
                "Si l'EBITDA (bleu foncé) est positif = l'activité génère de la trésorerie. "
                "L'écart entre la marge brute et l'EBITDA = le poids des charges fixes.")
        fig_pl = go.Figure()
        fig_pl.add_trace(go.Scatter(x=pl_mois["Période"], y=pl_mois["CA_HT"],
            name="CA HT", mode="lines+markers",
            line=dict(color=C["bleu"], width=2.5),
            marker=dict(size=5, color=C["bleu"]),
            hovertemplate="<b>%{x}</b><br>CA HT : %{y:,.0f} FCFA<extra></extra>"))
        fig_pl.add_trace(go.Scatter(x=pl_mois["Période"], y=pl_mois["Marge_Brute"],
            name="Marge Brute", mode="lines+markers",
            line=dict(color=C["vert"], width=2.5),
            marker=dict(size=5, color=C["vert"]),
            hovertemplate="<b>%{x}</b><br>Marge Brute : %{y:,.0f} FCFA<extra></extra>"))
        fig_pl.add_trace(go.Scatter(x=pl_mois["Période"], y=pl_mois["EBITDA"],
            name="EBITDA", mode="lines+markers",
            line=dict(color=C["orange"], width=2.5),
            marker=dict(size=5, color=C["orange"]),
            fill="tozeroy", fillcolor="rgba(243,156,18,0.08)",
            hovertemplate="<b>%{x}</b><br>EBITDA : %{y:,.0f} FCFA<extra></extra>"))
        fig_pl.add_trace(go.Scatter(x=pl_mois["Période"], y=pl_mois["EBIT"],
            name="EBIT (Résultat exploit.)", mode="lines+markers",
            line=dict(color=C["violet"], width=2, dash="dash"),
            marker=dict(size=5, color=C["violet"]),
            hovertemplate="<b>%{x}</b><br>EBIT : %{y:,.0f} FCFA<extra></extra>"))
        fig_pl.add_hline(y=0, line_color=C["rouge"], line_width=1, line_dash="dot")
        fig_pl = excel_style(fig_pl, 360)
        fig_pl.update_xaxes(tickangle=-35, tickfont=dict(size=8))
        fig_pl.update_yaxes(tickformat=",.0f", ticksuffix=" F")
        st.plotly_chart(fig_pl, use_container_width=True)

        # Export P&L
        buf_pl = io.BytesIO()
        pl_export = pl_mois[["Période","CA_TTC","CA_HT","Matieres","Marge_Brute",
                              "Charges_Var","Marge_Semi_Nette","Charges_Fixes","EBITDA","EBIT",
                              "Taux_MB_Brute","Taux_EBITDA"]].copy()
        pl_export.columns = ["Période","CA TTC","CA HT","Matières","Marge Brute",
                              "Charges Var.","Marge Semi-Nette","Charges Fixes","EBITDA","EBIT",
                              "Taux MB%","Taux EBITDA%"]
        with pd.ExcelWriter(buf_pl, engine="openpyxl") as w:
            pl_export.to_excel(w, sheet_name="P&L Mensuel", index=False)
        buf_pl.seek(0)
        st.download_button("⬇️ Exporter le P&L (.xlsx)", buf_pl,
            file_name="MULTIPACK_PL.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")


    # ════════════════════════════════════════════
    # ONGLET 2 — BUDGET VS RÉALISÉ
    # ════════════════════════════════════════════
    with cg2:
        st.markdown("<br>", unsafe_allow_html=True)
        comment(
            "Le budget vs réalisé compare ce qui était <b>prévu</b> avec ce qui a été <b>réellement fait</b>. "
            "Un écart positif (vert) = on a dépassé l'objectif. "
            "Un écart négatif (rouge) = on est en dessous. "
            "C'est le principal outil de pilotage du contrôleur de gestion.",
            "🎯 Comment lire le budget vs réalisé ?"
        )

        # Objectifs paramétrables
        cob1, cob2, cob3 = st.columns(3)
        with cob1:
            obj_ca     = st.number_input("Objectif CA TTC mensuel (FCFA)",
                value=int(ca_moy_hist*1.1), step=500_000, format="%d")
        with cob2:
            obj_mb_pct = st.slider("Objectif Taux Marge Brute (%)", 30, 70, 58, 1)
        with cob3:
            obj_ebitda_pct = st.slider("Objectif Taux EBITDA (%)", 5, 40, 15, 1)

        pl_mois["Ecart_CA"]     = pl_mois["CA_TTC"] - obj_ca
        pl_mois["Ecart_CA_Pct"] = pl_mois["Ecart_CA"] / obj_ca * 100
        pl_mois["Ecart_MB"]     = pl_mois["Taux_MB_Brute"] - obj_mb_pct
        pl_mois["Ecart_EBITDA"] = pl_mois["Taux_EBITDA"] - obj_ebitda_pct
        pl_mois["Statut_CA"]    = pl_mois["Ecart_CA"].apply(
            lambda v: "✅ Atteint" if v>=0 else ("⚠️ Proche" if v>=-obj_ca*0.05 else "❌ Manqué"))

        # KPIs budget globaux
        mois_atteints = (pl_mois["Ecart_CA"] >= 0).sum()
        mois_total    = len(pl_mois)
        ca_manquant   = pl_mois[pl_mois["Ecart_CA"]<0]["Ecart_CA"].sum()
        ca_surplus    = pl_mois[pl_mois["Ecart_CA"]>=0]["Ecart_CA"].sum()
        meilleur_mois = pl_mois.loc[pl_mois["Ecart_CA_Pct"].idxmax(), "Période"]
        pire_mois     = pl_mois.loc[pl_mois["Ecart_CA_Pct"].idxmin(), "Période"]

        bk1,bk2,bk3,bk4,bk5 = st.columns(5)
        with bk1: kpi("✅","Mois Objectif Atteint", f"{mois_atteints}/{mois_total}",
                      f"{mois_atteints/mois_total*100:.0f}% des mois","vert" if mois_atteints>mois_total/2 else "rouge")
        with bk2: kpi("📈","Surplus Cumulé",    fmt(ca_surplus), "CA au-dessus de l'objectif","vert")
        with bk3: kpi("📉","Manque à Gagner",   fmt(abs(ca_manquant)), "CA en dessous de l'objectif","rouge")
        with bk4: kpi("🌟","Meilleur Mois",     meilleur_mois,
                       f"+{pl_mois[pl_mois['Période']==meilleur_mois]['Ecart_CA_Pct'].values[0]:.1f}%","vert")
        with bk5: kpi("⚠️","Mois le + Difficile", pire_mois,
                       f"{pl_mois[pl_mois['Période']==pire_mois]['Ecart_CA_Pct'].values[0]:.1f}%","rouge")

        st.markdown("<br>", unsafe_allow_html=True)

        # GRAPHIQUE 1 — CA réalisé vs objectif avec zones
        section("📊  CA RÉALISÉ VS OBJECTIF MENSUEL")
        comment("Les barres bleues = CA réalisé. La ligne pointillée rouge = objectif. "
                "Les barres au-dessus de la ligne = objectif dépassé ✅. En dessous = manqué ❌.")
        fig_bvr = go.Figure()
        clrs_bar = [C["vert"] if v>=0 else C["rouge"] for v in pl_mois["Ecart_CA"]]
        fig_bvr.add_trace(go.Bar(x=pl_mois["Période"], y=pl_mois["CA_TTC"],
            name="CA Réalisé", marker_color=clrs_bar, opacity=0.85,
            hovertemplate="<b>%{x}</b><br>CA Réalisé : %{y:,.0f} FCFA<extra></extra>"))
        fig_bvr.add_hline(y=obj_ca, line_dash="dot", line_color=C["rouge"], line_width=2,
            annotation_text=f"Objectif : {fmt(obj_ca,'').strip()}",
            annotation_font=dict(color=C["rouge"],size=10), annotation_position="top right")
        fig_bvr.add_trace(go.Scatter(x=[None], y=[None], mode="markers",
            marker=dict(color=C["vert"],size=10,symbol="square"), name="✅ Objectif atteint"))
        fig_bvr.add_trace(go.Scatter(x=[None], y=[None], mode="markers",
            marker=dict(color=C["rouge"],size=10,symbol="square"), name="❌ Objectif manqué"))
        fig_bvr = excel_style(fig_bvr, 340)
        fig_bvr.update_xaxes(tickangle=-35, tickfont=dict(size=8))
        fig_bvr.update_yaxes(tickformat=",.0f", ticksuffix=" F")
        st.plotly_chart(fig_bvr, use_container_width=True)

        st.markdown("<br>", unsafe_allow_html=True)

        # GRAPHIQUE 2 — Écarts en % (waterfall mensuel)
        section("📐  ÉCARTS BUDGET VS RÉALISÉ (%) — Mois par mois")
        comment("Chaque barre montre l'écart entre le CA réalisé et l'objectif en pourcentage. "
                "🟢 Positif = dépassement. 🔴 Négatif = retard. "
                "Un retard supérieur à -10% est un signal d'alerte fort.")
        clrs_ecart = [C["vert"] if v>=0 else C["rouge"] for v in pl_mois["Ecart_CA_Pct"]]
        fig_ec = go.Figure(go.Bar(x=pl_mois["Période"], y=pl_mois["Ecart_CA_Pct"],
            marker_color=clrs_ecart,
            text=pl_mois["Ecart_CA_Pct"].apply(lambda v: f"{v:+.1f}%"),
            textposition="outside", textfont=dict(size=9),
            hovertemplate="<b>%{x}</b><br>Écart : %{y:+.1f}%<extra></extra>"))
        fig_ec.add_hline(y=0, line_color=C["muted"], line_width=1)
        fig_ec.add_hline(y=-10, line_dash="dot", line_color=C["rouge"], line_width=1,
            annotation_text="Seuil alerte -10%",
            annotation_font=dict(color=C["rouge"],size=9))
        fig_ec = excel_style(fig_ec, 300, False)
        fig_ec.update_xaxes(tickangle=-35, tickfont=dict(size=8))
        fig_ec.update_yaxes(ticksuffix="%")
        st.plotly_chart(fig_ec, use_container_width=True)

        st.markdown("<br>", unsafe_allow_html=True)

        # Tableau budget
        section("📋  TABLEAU DE SUIVI BUDGÉTAIRE COMPLET")
        bud_disp = pl_mois[["Période","CA_TTC","Ecart_CA","Ecart_CA_Pct",
                              "Taux_MB_Brute","Ecart_MB","Taux_EBITDA","Ecart_EBITDA","Statut_CA"]].copy()
        bud_disp["CA_TTC"]       = bud_disp["CA_TTC"].apply(lambda v: f"{v:,.0f} FCFA")
        bud_disp["Ecart_CA"]     = bud_disp["Ecart_CA"].apply(lambda v: f"{v:+,.0f} FCFA")
        bud_disp["Ecart_CA_Pct"] = bud_disp["Ecart_CA_Pct"].apply(lambda v: f"{v:+.1f}%")
        bud_disp["Taux_MB_Brute"]= bud_disp["Taux_MB_Brute"].apply(lambda v: f"{v:.1f}%")
        bud_disp["Ecart_MB"]     = bud_disp["Ecart_MB"].apply(lambda v: f"{v:+.1f}%")
        bud_disp["Taux_EBITDA"]  = bud_disp["Taux_EBITDA"].apply(lambda v: f"{v:.1f}%")
        bud_disp["Ecart_EBITDA"] = bud_disp["Ecart_EBITDA"].apply(lambda v: f"{v:+.1f}%")
        bud_disp.columns = ["Mois","CA Réalisé","Écart CA","Écart %",
                             "Taux MB%","Écart MB","Taux EBITDA%","Écart EBITDA","Statut"]
        st.dataframe(bud_disp, use_container_width=True, height=400)


    # ════════════════════════════════════════════
    # ONGLET 3 — MARGES & RENTABILITÉ
    # ════════════════════════════════════════════
    with cg3:
        st.markdown("<br>", unsafe_allow_html=True)
        comment(
            "La marge brute = CA HT - Coût des matières. "
            "La marge nette = ce qui reste après toutes les charges. "
            "Un produit ou segment avec une faible marge consomme les ressources sans rapporter suffisamment.",
            "💹 Comment lire les marges ?"
        )

        # Marges par segment
        seg_margin = df_cg.groupby("SEGMENT").agg(
            CA_HT=("CA_HT","sum"),
            Matieres=("Matieres","sum"),
            Marge_Brute=("Marge_Brute","sum"),
            Nb_Factures=("MONTANT TTC","count"),
        ).reset_index()
        seg_margin["Charges_Fixes_Part"] = charges_fixes_mois * len(pl_mois) * \
            seg_margin["CA_HT"] / seg_margin["CA_HT"].sum()
        seg_margin["Marge_Nette"]  = seg_margin["Marge_Brute"] - seg_margin["Charges_Fixes_Part"]
        seg_margin["Taux_MB"]      = seg_margin["Marge_Brute"] / seg_margin["CA_HT"] * 100
        seg_margin["Taux_MN"]      = seg_margin["Marge_Nette"]  / seg_margin["CA_HT"] * 100
        seg_margin["CA_par_Fact"]  = seg_margin["CA_HT"] / seg_margin["Nb_Factures"]
        seg_margin = seg_margin.sort_values("Taux_MB", ascending=False)

        mk1,mk2,mk3,mk4 = st.columns(4)
        mb_global  = seg_margin["Marge_Brute"].sum()/seg_margin["CA_HT"].sum()*100
        mn_global  = seg_margin["Marge_Nette"].sum()/seg_margin["CA_HT"].sum()*100
        best_seg   = seg_margin.iloc[0]["SEGMENT"]
        worst_seg  = seg_margin.iloc[-1]["SEGMENT"]
        with mk1: kpi("📊","Taux MB Global",  f"{mb_global:.1f}%","Marge brute / CA HT","vert" if mb_global>40 else "orange")
        with mk2: kpi("💹","Taux Marge Nette",f"{mn_global:.1f}%","Après charges fixes","vert" if mn_global>0 else "rouge")
        with mk3: kpi("🏆","Segment + Rentable", best_seg[:18],
                      f"MB : {seg_margin.iloc[0]['Taux_MB']:.1f}%","vert")
        with mk4: kpi("⚠️","Segment Sensible",   worst_seg[:18],
                      f"MB : {seg_margin.iloc[-1]['Taux_MB']:.1f}%","orange")

        st.markdown("<br>", unsafe_allow_html=True)

        col_m1, col_m2 = st.columns(2)
        with col_m1:
            # GRAPHIQUE — Taux de marge brute par segment
            section("💹  TAUX DE MARGE BRUTE PAR SEGMENT")
            comment("Plus la barre est haute = plus le segment est rentable sur ses ventes. "
                    "La ligne rouge = seuil minimum de 40% recommandé.")
            clr_mb = [C["vert"] if v>=50 else (C["orange"] if v>=35 else C["rouge"])
                      for v in seg_margin["Taux_MB"]]
            fig_mb = go.Figure(go.Bar(
                x=seg_margin["SEGMENT"], y=seg_margin["Taux_MB"],
                marker_color=clr_mb,
                text=seg_margin["Taux_MB"].apply(lambda v: f"{v:.1f}%"),
                textposition="outside", textfont=dict(size=10),
                hovertemplate="<b>%{x}</b><br>Taux MB : %{y:.1f}%<extra></extra>"))
            fig_mb.add_hline(y=40, line_dash="dot", line_color=C["rouge"], line_width=1.5,
                annotation_text="Seuil 40%", annotation_font=dict(color=C["rouge"],size=9))
            fig_mb = excel_style(fig_mb, 340, False)
            fig_mb.update_xaxes(tickangle=-25, tickfont=dict(size=8))
            fig_mb.update_yaxes(ticksuffix="%", range=[0, seg_margin["Taux_MB"].max()*1.2])
            st.plotly_chart(fig_mb, use_container_width=True)

        with col_m2:
            # GRAPHIQUE — Marge brute vs Marge nette par segment (comparaison)
            section("📊  MARGE BRUTE VS MARGE NETTE PAR SEGMENT")
            comment("La différence entre les deux barres = le poids des charges fixes imputées. "
                    "Un segment avec une marge nette négative coûte plus qu'il ne rapporte.")
            seg_sorted = seg_margin.sort_values("Marge_Brute", ascending=True).tail(10)
            fig_mn = go.Figure()
            fig_mn.add_trace(go.Bar(x=seg_sorted["Marge_Brute"], y=seg_sorted["SEGMENT"],
                orientation="h", name="Marge Brute", marker_color=C["bleu"],
                hovertemplate="<b>%{y}</b><br>Marge Brute : %{x:,.0f} FCFA<extra></extra>"))
            fig_mn.add_trace(go.Bar(x=seg_sorted["Marge_Nette"], y=seg_sorted["SEGMENT"],
                orientation="h", name="Marge Nette",
                marker_color=[C["vert"] if v>=0 else C["rouge"] for v in seg_sorted["Marge_Nette"]],
                hovertemplate="<b>%{y}</b><br>Marge Nette : %{x:,.0f} FCFA<extra></extra>"))
            fig_mn.update_layout(barmode="overlay", bargap=0.3)
            fig_mn = excel_style(fig_mn, 340)
            fig_mn.update_xaxes(tickformat=",.0f")
            st.plotly_chart(fig_mn, use_container_width=True)

        st.markdown("<br>", unsafe_allow_html=True)

        # GRAPHIQUE — Évolution mensuelle taux de marge
        section("📈  ÉVOLUTION MENSUELLE DU TAUX DE MARGE BRUTE")
        comment("Ce taux doit rester stable ou croître. Une baisse = hausse des coûts matières "
                "ou pression sur les prix de vente. La ligne pointillée = taux cible.")
        fig_ev_mb = go.Figure()
        fig_ev_mb.add_trace(go.Scatter(x=pl_mois["Période"], y=pl_mois["Taux_MB_Brute"],
            name="Taux MB (%)", mode="lines+markers",
            line=dict(color=C["vert"], width=3),
            marker=dict(size=8, color=[C["vert"] if v>=obj_mb_pct else C["rouge"]
                                        for v in pl_mois["Taux_MB_Brute"]],
                        line=dict(color="white",width=2)),
            hovertemplate="<b>%{x}</b><br>Taux MB : %{y:.1f}%<extra></extra>"))
        fig_ev_mb.add_hline(y=obj_mb_pct, line_dash="dot", line_color=C["orange"], line_width=1.5,
            annotation_text=f"Objectif {obj_mb_pct}%",
            annotation_font=dict(color=C["orange"],size=10))
        fig_ev_mb = excel_style(fig_ev_mb, 300, False)
        fig_ev_mb.update_xaxes(tickangle=-35, tickfont=dict(size=8))
        fig_ev_mb.update_yaxes(ticksuffix="%")
        st.plotly_chart(fig_ev_mb, use_container_width=True)

        # Tableau marges par segment
        section("📋  TABLEAU DES MARGES PAR SEGMENT")
        sm_disp = seg_margin.copy()
        for c in ["CA_HT","Matieres","Marge_Brute","Charges_Fixes_Part","Marge_Nette","CA_par_Fact"]:
            sm_disp[c] = sm_disp[c].apply(lambda v: f"{v:,.0f} FCFA")
        sm_disp["Taux_MB"] = sm_disp["Taux_MB"].apply(lambda v: f"{v:.1f}%")
        sm_disp["Taux_MN"] = sm_disp["Taux_MN"].apply(lambda v: f"{v:.1f}%")
        sm_disp = sm_disp.rename(columns={"SEGMENT":"Segment","CA_HT":"CA HT","Matieres":"Matières",
            "Marge_Brute":"Marge Brute","Charges_Fixes_Part":"Charg. Fixes Imputées",
            "Marge_Nette":"Marge Nette","Taux_MB":"Taux MB%","Taux_MN":"Taux MN%",
            "Nb_Factures":"Nb Fact.","CA_par_Fact":"CA/Facture"})
        st.dataframe(sm_disp, use_container_width=True)


    # ════════════════════════════════════════════
    # ONGLET 4 — BFR & TRÉSORERIE
    # ════════════════════════════════════════════
    with cg4:
        st.markdown("<br>", unsafe_allow_html=True)
        comment(
            "Le <b>BFR (Besoin en Fonds de Roulement)</b> = ce que l'entreprise doit financer "
            "pour fonctionner au quotidien. "
            "Un BFR élevé = beaucoup d'argent immobilisé dans les stocks et les créances. "
            "Le <b>DSO</b> = combien de jours en moyenne nos clients mettent à payer.",
            "🏦 Comprendre le BFR et le DSO"
        )

        # KPIs BFR/DSO
        bfk1,bfk2,bfk3,bfk4,bfk5 = st.columns(5)
        with bfk1: kpi("🏦","BFR Estimé", fmt(bfr),
                       "Stock + Créances - Dettes fourn.","orange" if bfr>0 else "vert")
        with bfk2: kpi("📦","Stock Immobilisé", fmt(stock_val),
                       f"{len(df_inv)} références","bleu")
        with bfk3: kpi("💳","Créances Clients", fmt(creances),
                       "Factures impayées","rouge" if creances>0 else "vert")
        with bfk4: kpi("🏭","Dettes Fournisseurs", fmt(dettes_fourn),
                       "Estimation 45j (matières)","vert")
        with bfk5:
            kpi("⏱️","DSO", f"{dso:.0f} jours",
                "Délai moyen de paiement","vert" if dso<=30 else ("orange" if dso<=60 else "rouge"))

        st.markdown("<br>", unsafe_allow_html=True)

        # GRAPHIQUE 1 — Décomposition BFR (waterfall)
        section("🏦  DÉCOMPOSITION DU BFR")
        comment("Le BFR est la somme de 3 éléments : "
                "Stock (argent immobilisé en produits) + Créances clients (argent dû) "
                "- Dettes fournisseurs (argent qu'on doit). "
                "Réduire le BFR améliore la trésorerie sans changer le CA.")
        fig_bfr = go.Figure(go.Waterfall(
            name="BFR",
            orientation="v",
            measure=["relative","relative","relative","total"],
            x=["📦 Stocks", "💳 Créances Clients", "🏭 − Dettes Fourn.", "= BFR Total"],
            y=[stock_val, creances, -dettes_fourn, 0],
            connector=dict(line=dict(color=C["muted"], width=1, dash="dot")),
            increasing=dict(marker=dict(color=C["rouge"])),
            decreasing=dict(marker=dict(color=C["vert"])),
            totals=dict(marker=dict(color=C["bleu"])),
            text=[fmt(stock_val,"").strip(), fmt(creances,"").strip(),
                  fmt(-dettes_fourn,"").strip(), fmt(bfr,"").strip()],
            textposition="outside", textfont=dict(size=11),
        ))
        fig_bfr = excel_style(fig_bfr, 360, False)
        fig_bfr.update_yaxes(tickformat=",.0f", ticksuffix=" F")
        st.plotly_chart(fig_bfr, use_container_width=True)

        st.markdown("<br>", unsafe_allow_html=True)

        col_bf1, col_bf2 = st.columns(2)

        with col_bf1:
            # GRAPHIQUE 2 — DSO simulé par segment
            section("⏱️  DSO ESTIMÉ PAR SEGMENT")
            comment("Le DSO = nombre de jours avant d'être payé. "
                    "Plus le DSO est élevé = plus on attend l'argent. "
                    "Un DSO > 60 jours est préoccupant pour la trésorerie.")
            dso_seg = df_fact.groupby("SEGMENT").apply(
                lambda g: (g[g["ETAT DE PAIEMENT"]=="Impayée"]["MONTANT TTC"].sum() /
                           (g["MONTANT TTC"].sum() / 365)) if g["MONTANT TTC"].sum() > 0 else 0
            ).reset_index()
            dso_seg.columns = ["SEGMENT","DSO"]
            dso_seg = dso_seg.sort_values("DSO", ascending=True)
            clr_dso = [C["vert"] if v<=30 else (C["orange"] if v<=60 else C["rouge"])
                       for v in dso_seg["DSO"]]
            fig_dso = go.Figure(go.Bar(
                x=dso_seg["DSO"], y=dso_seg["SEGMENT"], orientation="h",
                marker_color=clr_dso,
                text=dso_seg["DSO"].apply(lambda v: f"{v:.0f}j"),
                textposition="outside", textfont=dict(size=10),
                hovertemplate="<b>%{y}</b><br>DSO : %{x:.0f} jours<extra></extra>"))
            fig_dso.add_vline(x=30, line_dash="dot", line_color=C["vert"], line_width=1.5,
                annotation_text="30j ✅", annotation_font=dict(color=C["vert"],size=9))
            fig_dso.add_vline(x=60, line_dash="dot", line_color=C["rouge"], line_width=1.5,
                annotation_text="60j ⚠️", annotation_font=dict(color=C["rouge"],size=9))
            fig_dso = excel_style(fig_dso, 380, False)
            fig_dso.update_xaxes(ticksuffix=" j")
            st.plotly_chart(fig_dso, use_container_width=True)

        with col_bf2:
            # GRAPHIQUE 3 — Évolution créances mensuelles
            section("📈  ÉVOLUTION DES CRÉANCES CLIENTS")
            comment("Les créances = factures impayées à date. "
                    "Une courbe qui monte = les impayés s'accumulent → risque de trésorerie. "
                    "L'objectif est de maintenir les créances sous contrôle.")
            creances_mois = df_fact[df_fact["ETAT DE PAIEMENT"]=="Impayée"].copy()
            creances_mois["Période"] = creances_mois["DATE"].dt.to_period("M").astype(str)
            cr_m = creances_mois.groupby("Période")["MONTANT TTC"].sum().reset_index()
            fig_cr = go.Figure(go.Scatter(
                x=cr_m["Période"], y=cr_m["MONTANT TTC"],
                mode="lines+markers",
                line=dict(color=C["rouge"], width=2.5),
                marker=dict(size=7, color=C["rouge"], line=dict(color="white",width=2)),
                fill="tozeroy", fillcolor="rgba(231,76,60,0.08)",
                hovertemplate="<b>%{x}</b><br>Créances : %{y:,.0f} FCFA<extra></extra>"))
            fig_cr = excel_style(fig_cr, 380, False)
            fig_cr.update_xaxes(tickangle=-35, tickfont=dict(size=8))
            fig_cr.update_yaxes(tickformat=",.0f", ticksuffix=" F")
            st.plotly_chart(fig_cr, use_container_width=True)

        prevision_comment(
            f"Pour réduire le BFR de MULTIPACK : "
            f"(1) Réduire le DSO → relancer plus tôt, passer à 30j max. "
            f"(2) Optimiser les stocks → éviter les sur-stockages (valeur actuelle : {fmt(stock_val)}). "
            f"(3) Négocier des délais fournisseurs plus longs → passer de 45j à 60j. "
            f"Impact estimé : réduction BFR de {fmt(stock_val*0.15)}.",
            "💡 Actions pour réduire le BFR"
        )


    # ════════════════════════════════════════════
    # ONGLET 5 — POINT MORT
    # ════════════════════════════════════════════
    with cg5:
        st.markdown("<br>", unsafe_allow_html=True)
        comment(
            "Le <b>point mort</b> (ou seuil de rentabilité) est le CA minimum que MULTIPACK doit réaliser "
            "chaque mois pour couvrir toutes ses charges et ne pas perdre d'argent. "
            "En dessous = on perd de l'argent. Au-dessus = on est rentable.",
            "⚖️ Comprendre le point mort"
        )

        # KPIs point mort
        marge_secu = ca_moy_hist - point_mort_ttc
        taux_secu  = marge_secu / ca_moy_hist * 100 if ca_moy_hist else 0
        mois_pm_atteints = (pl_mois["CA_TTC"] >= point_mort_ttc).sum()

        pmk1,pmk2,pmk3,pmk4 = st.columns(4)
        with pmk1: kpi("⚖️","Point Mort Mensuel", fmt(point_mort_ttc),
                       "CA minimum pour être rentable","bleu")
        with pmk2: kpi("🛡️","Marge de Sécurité", fmt(marge_secu),
                       f"{taux_secu:.1f}% au-dessus du point mort",
                       "vert" if taux_secu>15 else ("orange" if taux_secu>5 else "rouge"))
        with pmk3: kpi("✅","Mois Rentables", f"{mois_pm_atteints}/{len(pl_mois)}",
                       f"CA > Point mort","vert" if mois_pm_atteints>len(pl_mois)*0.7 else "rouge")
        with pmk4:
            jours_pm = (point_mort_ttc / (ca_moy_hist / 30)) if ca_moy_hist else 0
            kpi("📅","Jours pour Atteindre le PM", f"{jours_pm:.0f} jours",
                "Par mois (sur ~30j)","vert" if jours_pm<20 else ("orange" if jours_pm<25 else "rouge"))

        st.markdown("<br>", unsafe_allow_html=True)

        # GRAPHIQUE 1 — CA vs Point mort (zone rentabilité)
        section("⚖️  CA MENSUEL VS POINT MORT — Zone de rentabilité")
        comment("La ligne rouge = seuil de rentabilité (point mort). "
                "Zone verte = on est rentable. Zone rouge = on perd de l'argent. "
                "Plus la barre est loin au-dessus de la ligne = plus la marge de sécurité est confortable.")
        fig_pm = go.Figure()
        # Zone rentabilité
        fig_pm.add_hrect(y0=point_mort_ttc, y1=pl_mois["CA_TTC"].max()*1.15,
            fillcolor="rgba(46,204,113,0.05)", layer="below", line_width=0,
            annotation_text="Zone rentable ✅",
            annotation_position="top right",
            annotation_font=dict(color=C["vert"], size=10))
        fig_pm.add_hrect(y0=0, y1=point_mort_ttc,
            fillcolor="rgba(231,76,60,0.05)", layer="below", line_width=0,
            annotation_text="Zone déficitaire ❌",
            annotation_position="bottom right",
            annotation_font=dict(color=C["rouge"], size=10))
        # Barres CA
        clr_pm = [C["vert"] if v>=point_mort_ttc else C["rouge"] for v in pl_mois["CA_TTC"]]
        fig_pm.add_trace(go.Bar(x=pl_mois["Période"], y=pl_mois["CA_TTC"],
            name="CA Réalisé", marker_color=clr_pm, opacity=0.85,
            hovertemplate="<b>%{x}</b><br>CA : %{y:,.0f} FCFA<extra></extra>"))
        # Ligne point mort
        fig_pm.add_hline(y=point_mort_ttc, line_dash="solid", line_color=C["rouge"], line_width=2.5,
            annotation_text=f"Point mort : {fmt(point_mort_ttc,'').strip()} FCFA",
            annotation_font=dict(color=C["rouge"], size=11, family="Inter"),
            annotation_position="top left")
        fig_pm = excel_style(fig_pm, 380)
        fig_pm.update_xaxes(tickangle=-35, tickfont=dict(size=8))
        fig_pm.update_yaxes(tickformat=",.0f", ticksuffix=" F")
        st.plotly_chart(fig_pm, use_container_width=True)

        st.markdown("<br>", unsafe_allow_html=True)

        col_pm1, col_pm2 = st.columns(2)
        with col_pm1:
            # Décomposition point mort
            section("🔍  DÉCOMPOSITION DU POINT MORT")
            pm_decomp = {
                "Masse salariale": ms,
                "Loyer usine":     lo,
                "Énergie":         en,
                "Maintenance":     ma,
                "Assurances":      as_,
                "Frais généraux":  fg,
                "Amortissements":  am,
            }
            df_pm_d = pd.DataFrame(list(pm_decomp.items()), columns=["Poste","Charges Fixes"])
            df_pm_d["Part du PM"] = df_pm_d["Charges Fixes"] / charges_fixes_mois * 100
            df_pm_d = df_pm_d.sort_values("Charges Fixes", ascending=True)
            fig_pm_d = go.Figure(go.Bar(
                x=df_pm_d["Charges Fixes"], y=df_pm_d["Poste"],
                orientation="h",
                marker_color=[C["bleu"],C["violet"],C["rose"],C["orange"],
                              C["vert"],C["bleu_clair"],C["bleu_fonce"]][:len(df_pm_d)],
                text=df_pm_d["Charges Fixes"].apply(lambda v: fmt(v,"").strip()),
                textposition="outside", textfont=dict(size=9),
                hovertemplate="<b>%{y}</b><br>%{x:,.0f} FCFA / mois<extra></extra>"))
            fig_pm_d = excel_style(fig_pm_d, 340, False)
            fig_pm_d.update_xaxes(tickformat=",.0f")
            st.plotly_chart(fig_pm_d, use_container_width=True)

        with col_pm2:
            # Simulation "What if"
            section("🔬  SIMULATION — Que se passe-t-il si on change ?")
            comment("Déplacez les curseurs pour voir l'impact sur le point mort.")
            sim_taux_mat = st.slider("Si le coût matières passe à (%)", 30, 65, int(tm*100), 1, key="sim_mat")
            sim_cf_pct   = st.slider("Si les charges fixes varient de (%)", -30, 30, 0, 5, key="sim_cf")

            sim_taux_mv   = 1 - sim_taux_mat/100 - tt - ef - cc
            sim_cf        = charges_fixes_mois * (1 + sim_cf_pct/100)
            sim_pm_ht     = sim_cf / sim_taux_mv if sim_taux_mv > 0 else 0
            sim_pm_ttc    = sim_pm_ht * 1.18
            delta_pm      = sim_pm_ttc - point_mort_ttc
            sim_ms        = marge_secu - delta_pm
            sim_taux_secu = sim_ms / ca_moy_hist * 100 if ca_moy_hist else 0

            cols_sim = st.columns(2)
            with cols_sim[0]:
                kpi("⚖️","Nouveau Point Mort", fmt(sim_pm_ttc),
                    f"{'▲' if delta_pm>0 else '▼'} {fmt(abs(delta_pm),'').strip()} vs actuel",
                    "rouge" if delta_pm>0 else "vert")
            with cols_sim[1]:
                kpi("🛡️","Nouvelle Marge Sécu.", fmt(sim_ms),
                    f"{sim_taux_secu:.1f}%",
                    "vert" if sim_taux_secu>10 else ("orange" if sim_taux_secu>0 else "rouge"))

            st.markdown(f"""
            <div class="{'prevision-box' if delta_pm>0 else 'comment-box'}">
                <div class="ct">📊 Impact de la simulation</div>
                {"Une hausse des coûts matières augmente le point mort de " + fmt(abs(delta_pm),'').strip() + " FCFA. "
                 if delta_pm>0 else
                 "Cette configuration améliore le point mort de " + fmt(abs(delta_pm),'').strip() + " FCFA. "}
                Marge de sécurité : <b>{sim_taux_secu:.1f}%</b>
                {"— ⚠️ Situation précaire" if sim_taux_secu<5 else
                 " — ✅ Situation correcte" if sim_taux_secu>10 else " — À surveiller"}.
            </div>""", unsafe_allow_html=True)


    # ════════════════════════════════════════════
    # ONGLET 6 — ANALYSE DES ÉCARTS
    # ════════════════════════════════════════════
    with cg6:
        st.markdown("<br>", unsafe_allow_html=True)
        comment(
            "L'analyse des écarts décompose la différence entre ce qui était prévu "
            "et ce qui a été réellement réalisé. "
            "Elle permet d'identifier si l'écart vient d'un <b>problème de volume</b> "
            "(on a moins vendu) ou d'un <b>problème de prix/marge</b> (on a moins bien vendu).",
            "📐 Comprendre l'analyse des écarts"
        )

        # Calcul des écarts par période
        pl_mois["Ecart_Volume"]    = (pl_mois["Nb_Factures"] - pl_mois["Nb_Factures"].mean()) * \
                                      pl_mois["CA_HT"].mean() / pl_mois["Nb_Factures"].mean()
        pl_mois["Ecart_Prix"]      = pl_mois["CA_HT"] - pl_mois["Nb_Factures"] * \
                                      (pl_mois["CA_HT"].mean() / pl_mois["Nb_Factures"].mean())
        pl_mois["Ecart_Marge"]     = pl_mois["EBITDA"] - pl_mois["EBITDA"].mean()
        pl_mois["Ecart_Charges"]   = -(pl_mois["Matieres"] - pl_mois["Matieres"].mean())
        pl_mois["Trim"]            = pl_mois["Période"].apply(
            lambda p: f"T{(int(p[5:7])-1)//3+1} {p[:4]}")

        # KPIs écarts globaux
        ek1,ek2,ek3,ek4 = st.columns(4)
        ev_pos = pl_mois[pl_mois["Ecart_Volume"]>0]["Ecart_Volume"].sum()
        ev_neg = pl_mois[pl_mois["Ecart_Volume"]<=0]["Ecart_Volume"].sum()
        em_pos = pl_mois[pl_mois["Ecart_Marge"]>0]["Ecart_Marge"].sum()
        em_neg = pl_mois[pl_mois["Ecart_Marge"]<=0]["Ecart_Marge"].sum()
        with ek1: kpi("📦","Écart Volume + Total", fmt(ev_pos), "Mois au-dessus de la moyenne","vert")
        with ek2: kpi("📉","Écart Volume − Total", fmt(ev_neg), "Mois en dessous de la moyenne","rouge")
        with ek3: kpi("💹","Écart Marge + Total",  fmt(em_pos), "Mois avec marge supérieure","vert")
        with ek4: kpi("⚠️","Écart Marge − Total",  fmt(em_neg), "Mois avec marge inférieure","rouge")

        st.markdown("<br>", unsafe_allow_html=True)

        # GRAPHIQUE 1 — Décomposition des écarts (Volume vs Prix)
        section("📐  DÉCOMPOSITION DES ÉCARTS — Volume vs Prix/Mix")
        comment("Les barres bleues = écart dû au <b>volume</b> (on a eu plus ou moins de commandes). "
                "Les barres roses = écart dû au <b>prix/mix</b> (les commandes valaient plus ou moins cher). "
                "Deux informations distinctes pour trouver la vraie cause d'un mauvais mois.")
        fig_ec2 = go.Figure()
        fig_ec2.add_trace(go.Bar(x=pl_mois["Période"], y=pl_mois["Ecart_Volume"],
            name="Écart Volume", marker_color=C["bleu"], opacity=0.85,
            hovertemplate="<b>%{x}</b><br>Écart volume : %{y:+,.0f} FCFA<extra></extra>"))
        fig_ec2.add_trace(go.Bar(x=pl_mois["Période"], y=pl_mois["Ecart_Prix"],
            name="Écart Prix/Mix", marker_color=C["rose"], opacity=0.85,
            hovertemplate="<b>%{x}</b><br>Écart prix/mix : %{y:+,.0f} FCFA<extra></extra>"))
        fig_ec2.update_layout(barmode="group")
        fig_ec2 = excel_style(fig_ec2, 340)
        fig_ec2.update_xaxes(tickangle=-35, tickfont=dict(size=8))
        fig_ec2.update_yaxes(tickformat="+,.0f", ticksuffix=" F")
        fig_ec2.add_hline(y=0, line_color=C["muted"], line_width=1)
        st.plotly_chart(fig_ec2, use_container_width=True)

        st.markdown("<br>", unsafe_allow_html=True)

        col_ec1, col_ec2 = st.columns(2)
        with col_ec1:
            # Écarts par trimestre
            section("📅  ÉCARTS CUMULÉS PAR TRIMESTRE")
            comment("Vue trimestrielle des écarts pour dégager les tendances sur un horizon plus long.")
            trim_agg = pl_mois.groupby("Trim").agg(
                Ecart_CA=("Ecart_CA","sum"),
                Ecart_Marge=("Ecart_Marge","sum"),
            ).reset_index()
            fig_trim = go.Figure()
            fig_trim.add_trace(go.Bar(x=trim_agg["Trim"], y=trim_agg["Ecart_CA"],
                name="Écart CA vs Objectif",
                marker_color=[C["vert"] if v>=0 else C["rouge"] for v in trim_agg["Ecart_CA"]],
                hovertemplate="<b>%{x}</b><br>Écart CA : %{y:+,.0f} FCFA<extra></extra>"))
            fig_trim.add_trace(go.Scatter(x=trim_agg["Trim"], y=trim_agg["Ecart_Marge"],
                name="Écart EBITDA", mode="lines+markers",
                line=dict(color=C["orange"], width=2, dash="dash"),
                marker=dict(size=8),
                hovertemplate="<b>%{x}</b><br>Écart EBITDA : %{y:+,.0f} FCFA<extra></extra>"))
            fig_trim.add_hline(y=0, line_color=C["muted"], line_width=1)
            fig_trim = excel_style(fig_trim, 320)
            fig_trim.update_xaxes(tickangle=-20, tickfont=dict(size=9))
            fig_trim.update_yaxes(tickformat="+,.0f", ticksuffix=" F")
            st.plotly_chart(fig_trim, use_container_width=True)

        with col_ec2:
            # Top 5 mois avec meilleurs/pires écarts
            section("🏆  TOP MOIS — Meilleurs et Pires Performances")
            comment("Les 5 meilleurs et les 5 pires mois en termes d'écart vs objectif.")
            top5_pos  = pl_mois.nlargest(5, "Ecart_CA")[["Période","CA_TTC","Ecart_CA","Ecart_CA_Pct"]]
            top5_neg  = pl_mois.nsmallest(5,"Ecart_CA")[["Période","CA_TTC","Ecart_CA","Ecart_CA_Pct"]]
            all_top   = pd.concat([top5_pos, top5_neg]).sort_values("Ecart_CA", ascending=True)
            clr_top   = [C["vert"] if v>=0 else C["rouge"] for v in all_top["Ecart_CA"]]
            fig_top = go.Figure(go.Bar(
                x=all_top["Ecart_CA"], y=all_top["Période"], orientation="h",
                marker_color=clr_top,
                text=all_top["Ecart_CA"].apply(lambda v: f"{v:+,.0f} FCFA"),
                textposition="outside", textfont=dict(size=9),
                hovertemplate="<b>%{y}</b><br>Écart : %{x:+,.0f} FCFA<extra></extra>"))
            fig_top.add_vline(x=0, line_color=C["muted"], line_width=1)
            fig_top = excel_style(fig_top, 320, False)
            fig_top.update_xaxes(tickformat="+,.0f")
            st.plotly_chart(fig_top, use_container_width=True)

        st.markdown("<br>", unsafe_allow_html=True)

        # Tableau d'analyse complet
        section("📋  TABLEAU D'ANALYSE DES ÉCARTS — Détail mensuel")
        ec_disp = pl_mois[["Période","CA_TTC","Ecart_CA","Ecart_CA_Pct",
                             "Ecart_Volume","Ecart_Prix","EBITDA","Ecart_Marge","Statut_CA"]].copy()
        for col in ["CA_TTC","Ecart_CA","Ecart_Volume","Ecart_Prix","EBITDA","Ecart_Marge"]:
            ec_disp[col] = ec_disp[col].apply(lambda v: f"{v:+,.0f} FCFA")
        ec_disp["Ecart_CA_Pct"] = ec_disp["Ecart_CA_Pct"].apply(lambda v: f"{v:+.1f}%")
        ec_disp.columns = ["Mois","CA Réalisé","Écart CA","Écart %",
                            "Écart Volume","Écart Prix/Mix","EBITDA","Écart EBITDA","Statut"]
        st.dataframe(ec_disp, use_container_width=True, height=400)

        # Export
        buf_cg = io.BytesIO()
        with pd.ExcelWriter(buf_cg, engine="openpyxl") as w:
            pl_mois.to_excel(w, sheet_name="P&L Mensuel", index=False)
            seg_margin.to_excel(w, sheet_name="Marges par Segment", index=False)
            dso_seg.to_excel(w, sheet_name="DSO par Segment", index=False)
        buf_cg.seek(0)
        st.download_button("⬇️ Exporter le rapport Contrôle de Gestion (.xlsx)",
            buf_cg, file_name=f"MULTIPACK_CG_{datetime.now().strftime('%Y%m%d')}.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")


elif page == "⚠️  Alertes":

    ca_total  = df_f["MONTANT TTC"].sum()
    ca_paye   = df_f[df_f["ETAT DE PAIEMENT"]=="Payée"]["MONTANT TTC"].sum()
    ca_impaye = df_f[df_f["ETAT DE PAIEMENT"]=="Impayée"]["MONTANT TTC"].sum()
    taux_rec  = ca_paye/ca_total*100 if ca_total>0 else 0
    nb_rupt   = (df_inv["Statut"]=="Non disponible").sum()
    nb_faib   = (df_inv["Statut"]=="Stock faible").sum()
    rend_moy  = df_prod["Taux rendement"].mean()

    comment("Cette page regroupe <b>toutes les informations essentielles</b> pour prendre les bonnes décisions. "
            "🔴 Urgent = action immédiate · 🟡 Attention = à surveiller · 🟢 Positif = bonne nouvelle.",
            "⚠️ Comment utiliser cette page ?")

    st.markdown("<br>", unsafe_allow_html=True)
    c1, c2, c3 = st.columns(3)

    def health_card_v2(col, icon, title, valeur, sous, status_txt, color):
        col.markdown(f"""
        <div style="background:#FFFFFF;border-radius:14px;padding:22px;border:1px solid {C['bordure']};
                    border-top:5px solid {color};box-shadow:0 2px 10px rgba(0,0,0,0.07);">
            <div style="font-size:1.8rem;margin-bottom:8px;">{icon}</div>
            <div style="font-size:0.7rem;font-weight:700;color:{C['muted']};text-transform:uppercase;letter-spacing:0.07em;">{title}</div>
            <div style="font-size:1.6rem;font-weight:800;color:{C['texte']};margin:8px 0 3px;">{valeur}</div>
            <div style="font-size:0.78rem;color:{C['muted']};margin-bottom:10px;">{sous}</div>
            <div style="font-size:0.88rem;font-weight:600;color:{color};">{status_txt}</div>
        </div>""", unsafe_allow_html=True)

    t_col = C["vert"] if taux_rec>=80 else (C["orange"] if taux_rec>=65 else C["rouge"])
    s_col = C["vert"] if (nb_rupt+nb_faib)==0 else (C["orange"] if (nb_rupt+nb_faib)<=3 else C["rouge"])
    p_col = C["vert"] if rend_moy>=85 else (C["orange"] if rend_moy>=75 else C["rouge"])

    health_card_v2(c1,"💰","TRÉSORERIE", f"{taux_rec:.1f}%", f"{fmt(ca_impaye)} non encaissés",
        ("🟢 Objectif atteint" if taux_rec>=80 else "🟡 En dessous de l'objectif" if taux_rec>=65 else "🔴 Situation critique"), t_col)
    health_card_v2(c2,"📦","STOCKS", f"{nb_rupt} rupture(s)", f"{nb_faib} stock(s) faible(s)",
        ("🟢 Situation normale" if (nb_rupt+nb_faib)==0 else "🟡 À surveiller" if (nb_rupt+nb_faib)<=3 else "🔴 Situation critique"), s_col)
    health_card_v2(c3,"🏭","PRODUCTION", f"{rend_moy:.1f}%", f"Rebuts : {df_prod['Rebuts'].sum():,.0f} u.",
        ("🟢 Objectif atteint" if rend_moy>=85 else "🟡 En dessous de l'objectif" if rend_moy>=75 else "🔴 Situation critique"), p_col)

    st.markdown("<br>", unsafe_allow_html=True)
    col_al, col_reco = st.columns(2)

    with col_al:
        section("🚨  ALERTES ACTIVES")
        rupt = df_inv[df_inv["Statut"]=="Non disponible"]
        faib = df_inv[df_inv["Statut"]=="Stock faible"]
        if len(rupt)==0 and len(faib)==0:
            st.markdown('<div class="alert-g">✅<div><b>Stocks en bonne santé</b><br>Aucune rupture ni stock critique détecté.</div></div>', unsafe_allow_html=True)
        for _, r in rupt.iterrows():
            st.markdown(f'<div class="alert-r">🔴<div><b>RUPTURE — {r["designation"]}</b><br><small>{r["categorie"]} · 0 unité · ➜ Commander immédiatement</small></div></div>', unsafe_allow_html=True)
        for _, r in faib.head(4).iterrows():
            st.markdown(f'<div class="alert-y">⚠️<div><b>STOCK FAIBLE — {r["designation"]}</b><br><small>{r["categorie"]} · {int(r["Stock final"])} u. restantes (seuil : {int(r["seuil"])} u.) · ➜ Planifier commande</small></div></div>', unsafe_allow_html=True)
        imp = df_f[df_f["ETAT DE PAIEMENT"]=="Impayée"].groupby("CLIENT")["MONTANT TTC"].sum()
        for cl, v in imp.nlargest(4).items():
            st.markdown(f'<div class="alert-r">💸<div><b>IMPAYÉ — {cl}</b><br><small>Créance : <b>{v:,.0f} FCFA</b> · ➜ Appel + lettre de relance</small></div></div>', unsafe_allow_html=True)
        if rend_moy < 80:
            st.markdown(f'<div class="alert-r">⚙️<div><b>RENDEMENT FAIBLE — {rend_moy:.1f}%</b><br><small>Sous 80% · ➜ Audit technique + maintenance machines</small></div></div>', unsafe_allow_html=True)

    with col_reco:
        section("💡  RECOMMANDATIONS STRATÉGIQUES")
        top_seg    = df_f.groupby("SEGMENT")["MONTANT TTC"].sum().idxmax()
        top_zone   = df_f.groupby("ZONE")["MONTANT TTC"].sum().idxmax()
        top_client = df_f.groupby("CLIENT")["MONTANT TTC"].sum().idxmax()
        ca21 = df_fact[df_fact["ANNEE"]==2022]["MONTANT TTC"].sum()
        ca22 = df_fact[df_fact["ANNEE"]==2023]["MONTANT TTC"].sum()
        crois = (ca22-ca21)/ca21*100 if ca21 else 0

        recos = []
        if taux_rec >= 80:
            recos.append(("g","✅ Trésorerie saine", f"Taux de recouvrement {taux_rec:.1f}% ≥ 80%. Maintenez la rigueur de suivi."))
        else:
            recos.append(("r","🔴 Améliorer le recouvrement", f"Taux à {taux_rec:.1f}% (objectif : 80%). Relances à J+30, J+60, J+90."))
        recos.append(("g",f"🏆 Fidéliser {top_client.split()[0]}", f"Meilleur client : <b>{top_client}</b>. Visiter régulièrement, proposer une remise fidélité."))
        recos.append(("g",f"🗺️ Renforcer la {top_zone}", f"Zone la plus rentable. Augmenter la fréquence des visites commerciales."))
        recos.append(("g",f"🏷️ Miser sur {top_seg.split()[0]}", f"Segment <b>{top_seg}</b> = moteur de croissance. Concentrer la prospection sur ce profil."))
        if crois > 0:
            recos.append(("g",f"📈 Croissance +{crois:.1f}%", f"CA en hausse entre 2022 et 2023. Maintenez en fixant des objectifs mensuels ambitieux."))
        else:
            recos.append(("y","📉 Relancer la croissance", f"CA en recul de {crois:.1f}%. Analyser les clients perdus et lancer un plan de reconquête."))
        if nb_faib > 0:
            recos.append(("y","📦 Optimiser les achats stock", f"{nb_faib} références faibles. Mettre en place une commande automatique dès le seuil atteint."))
        recos.append(("g","🔮 Consulter les Prévisions", "Rendez-vous sur la page <b>Prévisions & Anticipations</b> pour anticiper les 6 prochains mois."))

        for typ, titre, msg in recos:
            cls = {"r":"alert-r","y":"alert-y","g":"alert-g"}[typ]
            st.markdown(f'<div class="{cls}"><div><b>{titre}</b><br><small>{msg}</small></div></div>', unsafe_allow_html=True)

    # Export
    st.markdown("<br>", unsafe_allow_html=True)
    section("📥  EXPORTS — Rapport Excel & Présentation PowerPoint")
    comment(
        "Téléchargez le rapport Excel complet (données brutes, impayés, stock, production) "
        "ou la <b>présentation PowerPoint prête à projeter</b> pour la direction (6 diapositives, "
        "graphiques, KPIs, recommandations automatiques).",
        "📥 Formats disponibles"
    )

    col_dl1, col_dl2, _ = st.columns([1, 1, 2])

    with col_dl1:
        buf_xl = io.BytesIO()
        with pd.ExcelWriter(buf_xl, engine="openpyxl") as w:
            df_fact.to_excel(w, sheet_name="Factures", index=False)
            df_inv.to_excel(w, sheet_name="Inventaire", index=False)
            df_ent.to_excel(w, sheet_name="Entrées Stock", index=False)
            df_sor.to_excel(w, sheet_name="Sorties Stock", index=False)
            df_prod.to_excel(w, sheet_name="Production", index=False)
            df_f[df_f["ETAT DE PAIEMENT"]=="Impayée"].to_excel(w, sheet_name="Impayés Prioritaires", index=False)
        buf_xl.seek(0)
        st.download_button(
            label="⬇️  Rapport Excel Complet",
            data=buf_xl,
            file_name=f"MULTIPACK_Rapport_{datetime.now().strftime('%Y%m%d_%H%M')}.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            use_container_width=True)

    with col_dl2:
        try:
            buf_pptx = generer_rapport_pptx(df_fact, df_inv, df_ent, df_sor, df_prod, df_f, datetime.now())
            st.download_button(
                label="📊  Présentation PowerPoint Direction",
                data=buf_pptx,
                file_name=f"MULTIPACK_Direction_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True)
        except Exception as e:
            st.error(f"Erreur génération PowerPoint : {e}")
